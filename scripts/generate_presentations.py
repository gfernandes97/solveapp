"""
Generate FinanceApp PowerPoint presentations.
Usage: python scripts/generate_presentations.py
Output: docs/presentations/
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Color palette ────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
C_BRAND  = RGBColor(0x00, 0x84, 0xFF)   # electric blue
C_ACCENT = RGBColor(0x00, 0xD4, 0x8A)   # emerald green
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xFF)   # near-white blue
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_WARN   = RGBColor(0xF5, 0x9E, 0x0B)   # amber
C_RED    = RGBColor(0xEF, 0x44, 0x44)

VERSION_MARKET = "v1.1.0"
VERSION_COSTS  = "v1.1.0"
DATE           = "2026-05-11"

OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations")


# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=C_WHITE, bold=False, spacing=1.2):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(font_size * (spacing - 1))
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def slide_header(slide, title: str, subtitle: str = "", accent_color=C_BRAND):
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, accent_color)          # top stripe
    add_rect(slide, 0, 7.42, 13.33, 0.08, accent_color)       # bottom stripe
    add_text(slide, title, 0.5, 0.2, 12, 0.7,
             font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.9, 12, 0.4,
                 font_size=14, bold=False, color=C_GRAY)


def title_slide(prs, title, subtitle, version, date, accent=C_BRAND):
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK)

    add_rect(slide, 0, 0, 13.33, 0.12, accent)
    add_rect(slide, 0, 7.38, 13.33, 0.12, accent)
    add_rect(slide, 0, 2.8, 13.33, 0.04, accent)

    add_text(slide, "FinanceApp", 1, 0.5, 11, 0.9,
             font_size=18, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1, 1.5, 11, 1.4,
             font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 1, 3.1, 11, 0.8,
             font_size=20, bold=False, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, f"{version}  ·  {date}", 1, 6.8, 11, 0.5,
             font_size=13, bold=False, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 1 — MARKET RESEARCH
# ══════════════════════════════════════════════════════════════════════════════

def build_market_research(prs: Presentation):

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    title_slide(
        prs,
        "Estudo de Mercado",
        "Oportunidade, Concorrência e Estratégia de Go-to-Market",
        VERSION_MARKET, DATE, C_BRAND
    )

    # ── Slide 2: Visão Geral do Mercado Brasileiro ────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Mercado Brasileiro de Fintechs", "Contexto e tamanho do mercado")

    stats = [
        ("160M",  "adultos com\nsmartphone"),
        ("130M",  "usuários de\ninternet banking"),
        ("25M",   "investidores\nna B3"),
        ("70M",   "CPFs negativados"),
    ]
    colors = [C_BRAND, C_ACCENT, C_WARN, C_RED]
    for i, ((val, label), color) in enumerate(zip(stats, colors)):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 1.5, 2.8, 2.8, color)
        add_text(slide, val,   x, 1.7,  2.8, 1.2, font_size=42, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 3.0,  2.8, 0.9, font_size=14,
                 align=PP_ALIGN.CENTER)

    add_text(slide,
             "Brasil é o maior mercado de fintechs da América Latina, "
             "com penetração crescente e forte demanda reprimida por educação financeira.",
             0.5, 4.6, 12.3, 0.8, font_size=14, color=C_LIGHT)

    # ── Slide 3: TAM / SAM / SOM ──────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "TAM · SAM · SOM", "Dimensionamento do mercado endereçável")

    data = [
        ("TAM", "Total Addressable Market",
         "Adultos financeiramente ativos no Brasil", "120M pessoas", C_BRAND, 9.5),
        ("SAM", "Serviceable Addressable Market",
         "Millennials + Gen-Z (18-42 anos) com smartphone e renda formal", "45M pessoas", C_ACCENT, 6.5),
        ("SOM", "Serviceable Obtainable Market",
         "Early adopters dispostos a pagar por app financeiro", "3–5M pessoas", C_WARN, 4.0),
    ]
    for i, (abbr, full, desc, size, color, w) in enumerate(data):
        y = 1.5 + i * 1.8
        add_rect(slide, 0.4, y, w, 1.4, color)
        add_text(slide, abbr,  0.5, y + 0.05, 1.2, 0.6, font_size=28, bold=True)
        add_text(slide, full,  1.8, y + 0.05, w - 1.6, 0.5, font_size=13, bold=True)
        add_text(slide, desc,  1.8, y + 0.55, w - 1.6, 0.6, font_size=12)
        add_text(slide, size,  w - 1.5, y + 0.1, 1.8, 0.6,
                 font_size=20, bold=True, align=PP_ALIGN.RIGHT)

    # ── Slide 4: Mapa Competitivo ─────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Panorama Competitivo", "Apps de controle de gastos e investimentos")

    headers = ["App", "Preço/mês", "Gastos", "Investimentos", "Educação", "Design"]
    rows = [
        ["Organizze",    "R$19,90", "✔", "✘", "✘", "Bom"],
        ["Mobills",      "R$14,90", "✔", "✘", "✘", "Datado"],
        ["Kinvo",        "R$29,90", "✘", "✔", "Básico", "Ótimo"],
        ["Gorila",       "R$39,90", "✘", "✔", "✘", "Profissional"],
        ["Guiabolso",    "Descontinuado", "✔", "✘", "✘", "—"],
        ["FinanceApp ★", "R$19,90–34,90", "✔", "✔", "✔", "Mobile-first"],
    ]
    col_widths = [2.2, 1.8, 1.0, 1.5, 1.5, 2.0]
    col_x = [0.3]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    row_h = 0.55
    header_y = 1.4
    add_rect(slide, 0.3, header_y, 12.5, row_h, C_BRAND)
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_widths)):
        add_text(slide, h, x + 0.05, header_y + 0.05, w, row_h - 0.1,
                 font_size=13, bold=True, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        y = header_y + row_h * (i + 1)
        bg = C_ACCENT if row[0].startswith("Finance") else (C_DARK if i % 2 == 0 else RGBColor(0x16, 0x23, 0x3A))
        add_rect(slide, 0.3, y, 12.5, row_h, bg)
        for j, (cell, x, w) in enumerate(zip(row, col_x, col_widths)):
            cell_color = C_WARN if cell == "✔" else (C_RED if cell == "✘" else C_WHITE)
            if row[0].startswith("Finance"):
                cell_color = C_DARK
            add_text(slide, cell, x + 0.05, y + 0.07, w, row_h - 0.14,
                     font_size=12, color=cell_color, align=PP_ALIGN.CENTER)

    # ── Slide 5: Gap de Mercado ───────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Oportunidade de Mercado", "Gap identificado e benchmarks internacionais")

    add_rect(slide, 0.4, 1.4, 8.0, 2.6, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "Gap Principal", 0.6, 1.5, 7.6, 0.5,
             font_size=16, bold=True, color=C_ACCENT)
    gaps = [
        "• Nenhum app brasileiro integra gastos + investimentos + educação",
        "• Guiabolso descontinuado deixou ~5 milhões de usuários órfãos",
        "• Usuários usam 2–3 apps diferentes para o que 1 deveria resolver",
        "• Crescimento de investidores B3: 500K (2018) → 25M (2026)",
    ]
    add_multiline(slide, gaps, 0.6, 2.0, 7.6, 2.0, font_size=13, color=C_LIGHT)

    add_rect(slide, 8.8, 1.4, 4.1, 2.6, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "YNAB (EUA) — Referência", 8.9, 1.5, 3.9, 0.5,
             font_size=15, bold=True, color=C_WARN)
    ynab = [
        "USD 14,99/mês",
        "> 1 milhão de usuários pagantes",
        "Metodologia como produto",
        "Educação = retenção alta",
        "NPS > 80",
    ]
    add_multiline(slide, ynab, 8.9, 2.0, 3.9, 2.0, font_size=13, color=C_LIGHT)

    add_rect(slide, 0.4, 4.3, 12.5, 2.6, RGBColor(0x05, 0x3A, 0x1F))
    add_text(slide, "Disposição a Pagar no Brasil", 0.6, 4.4, 12, 0.45,
             font_size=16, bold=True, color=C_ACCENT)
    pay = [
        "• 15–20% dos usuários de apps financeiros pagam por premium (crescendo)",
        "• Faixa validada: R$15–45/mês  ·  Modelo Freemium com 3 tiers",
        "• Tier Essentials R$19,90  ·  Tier Premium R$34,90",
    ]
    add_multiline(slide, pay, 0.6, 4.9, 12, 1.8, font_size=13, color=C_LIGHT)

    # ── Slide 6: Personas ─────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Personas", "Perfis de usuário identificados na pesquisa")

    personas = [
        {
            "name": "Ana, 26 anos",
            "role": "1º emprego CLT",
            "income": "R$3.500–6.000/mês",
            "goal": "Guardar para imóvel em 5 anos",
            "pain": "Não sabe quanto sobra para investir",
            "pay": "R$20–30/mês",
            "color": C_BRAND,
        },
        {
            "name": "Carlos, 34 anos",
            "role": "Gerente de área",
            "income": "R$8.000–15.000/mês",
            "goal": "Saber se investe o suficiente",
            "pain": "Controla gastos no Excel (ineficiente)",
            "pay": "R$30–50/mês",
            "color": C_ACCENT,
        },
        {
            "name": "Lucas, 19 anos",
            "role": "Estagiário universitário",
            "income": "R$1.500/mês",
            "goal": "Aprender finanças e investir no início",
            "pain": "Precisa de educação básica primeiro",
            "pay": "R$10–15/mês",
            "color": C_WARN,
        },
    ]

    for i, p in enumerate(personas):
        x = 0.4 + i * 4.3
        add_rect(slide, x, 1.4, 4.0, 5.5, p["color"])
        add_text(slide, p["name"],   x + 0.15, 1.55, 3.7, 0.6,
                 font_size=20, bold=True)
        add_text(slide, p["role"],   x + 0.15, 2.15, 3.7, 0.4,
                 font_size=13)
        add_rect(slide, x, 2.6, 4.0, 0.04, C_DARK)
        lines = [
            f"Renda: {p['income']}",
            f"Objetivo: {p['goal']}",
            f"Dor: {p['pain']}",
            f"Disposição a pagar: {p['pay']}",
        ]
        add_multiline(slide, lines, x + 0.15, 2.75, 3.7, 4.0,
                      font_size=12, color=C_DARK, spacing=1.5)

    # ── Slide 7: Modelo de Monetização ───────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Modelo de Monetização", "Freemium + 3 tiers de assinatura")

    tiers = [
        {
            "name": "Free",
            "price": "R$0",
            "features": [
                "Lançamento manual de transações",
                "2 contas bancárias",
                "Relatórios básicos",
                "Metas simples",
            ],
            "color": C_GRAY,
        },
        {
            "name": "Essentials",
            "price": "R$19,90/mês",
            "features": [
                "Open Finance (contas automáticas)",
                "Investimentos básicos",
                "Metas avançadas",
                "Educação — módulo inicial",
            ],
            "color": C_BRAND,
        },
        {
            "name": "Premium",
            "price": "R$34,90/mês",
            "features": [
                "Categorização por IA",
                "Impostos automatizados",
                "Simuladores avançados",
                "Suporte prioritário",
            ],
            "color": C_ACCENT,
        },
    ]

    for i, t in enumerate(tiers):
        x = 0.5 + i * 4.2
        add_rect(slide, x, 1.4, 3.8, 5.5, t["color"])
        add_text(slide, t["name"],  x, 1.5,  3.8, 0.65,
                 font_size=24, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, t["price"], x, 2.2,  3.8, 0.65,
                 font_size=18, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, x, 2.9, 3.8, 0.04, C_DARK)
        add_multiline(slide, ["• " + f for f in t["features"]],
                      x + 0.15, 3.0, 3.5, 3.5, font_size=13, color=C_DARK)

    # ── Slide 8: Riscos de Mercado ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Riscos de Mercado", "Principais ameaças e mitigações")

    risks = [
        ("ALTO",   C_RED,    "Nubank lançar funcionalidades similares",
         "Distribuição de 100M+ clientes vs. nosso nicho premium/educação"),
        ("ALTO",   C_RED,    "CAC elevado",
         "Custo de aquisição high — mitigar com conteúdo orgânico e parcerias"),
        ("MÉDIO",  C_WARN,   "Mudanças regulatórias Open Finance",
         "BC atualizando normas continuamente — acompanhar via API aggregator"),
        ("MÉDIO",  C_WARN,   "Churn por LTV baixo",
         "Primeiros 90 dias críticos — onboarding e education loop são chave"),
    ]

    for i, (level, color, risk, mitigation) in enumerate(risks):
        y = 1.4 + i * 1.4
        add_rect(slide, 0.4, y, 1.2, 1.2, color)
        add_text(slide, level, 0.4, y + 0.25, 1.2, 0.7,
                 font_size=14, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.7, y, 11.0, 1.2, RGBColor(0x16, 0x23, 0x3A))
        add_text(slide, risk,       1.85, y + 0.07, 10.7, 0.45,
                 font_size=15, bold=True, color=C_WHITE)
        add_text(slide, mitigation, 1.85, y + 0.55, 10.7, 0.6,
                 font_size=12, color=C_GRAY)

    # ── Slide 9: Estratégia de Entrada ───────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Estratégia de Entrada no Mercado",
                 "Posicionamento, canais e validação antes do lançamento")

    # Left block: positioning
    add_rect(slide, 0.4, 1.4, 6.0, 5.5, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "Posicionamento", 0.6, 1.5, 5.6, 0.45,
             font_size=16, bold=True, color=C_BRAND)
    pos_lines = [
        "\"O único app brasileiro que une controle",
        " de gastos, carteira de investimentos",
        " e educação financeira em um só lugar.\"",
        "",
        "Público inicial prioritário:",
        "• Ex-usuários do Guiabolso (~5M órfãos)",
        "• Jovens investidores B3 (18–35 anos)",
        "• Profissionais que usam Excel p/ finanças",
    ]
    add_multiline(slide, pos_lines, 0.6, 2.05, 5.6, 4.5,
                  font_size=12, color=C_LIGHT, spacing=1.35)

    # Right block: channels + validation
    add_rect(slide, 6.7, 1.4, 6.3, 2.5, RGBColor(0x05, 0x2A, 0x1A))
    add_text(slide, "Canais de Aquisição", 6.9, 1.5, 5.9, 0.45,
             font_size=16, bold=True, color=C_ACCENT)
    channels = [
        "• Conteúdo orgânico: Instagram · TikTok Finance · YouTube",
        "• SEO: artigos longos sobre finanças pessoais BR",
        "• Parcerias: influenciadores de finanças (micro 50K–500K)",
        "• Comunidades: Reddit r/investimentos · Discord fintech BR",
    ]
    add_multiline(slide, channels, 6.9, 2.05, 5.9, 1.7,
                  font_size=11, color=C_LIGHT, spacing=1.3)

    add_rect(slide, 6.7, 4.15, 6.3, 2.75, RGBColor(0x1A, 0x14, 0x05))
    add_text(slide, "Pré-lançamento (antes do MVP)", 6.9, 4.25, 5.9, 0.45,
             font_size=16, bold=True, color=C_WARN)
    prelaunch = [
        "1. Landing page com proposta de valor clara",
        "2. Formulário de lista de espera (meta: 500 e-mails)",
        "3. Entrevistar 20+ pessoas da lista (Jobs to be Done)",
        "4. Validar willingness to pay com página de preços",
        "5. Só iniciar dev com lista validada",
    ]
    add_multiline(slide, prelaunch, 6.9, 4.75, 5.9, 2.0,
                  font_size=11, color=C_LIGHT, spacing=1.3)

    # ── Slide 10: Primeiros Usuários (0 → 1.000) ─────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Primeiros Usuários — 0 → 1.000",
                 "Beta fechado, onboarding e loop de ativação")

    # Activation funnel steps
    steps = [
        ("1", "Lista de Espera",
         "Beta fechado gera FOMO e filtra\nusuários mais motivados.",
         C_BRAND),
        ("2", "Convite & Cadastro",
         "Onboarding em 5 passos.\nConectar 1 conta bancária via Open Finance.",
         C_ACCENT),
        ("3", "Diagnóstico Instantâneo",
         "Mostrar gráfico de gastos dos últimos\n30 dias logo no 1º acesso.",
         C_WARN),
        ("4", "Criar 1ª Meta",
         "Guiar usuário a criar uma meta\nfinanceira no mesmo dia.",
         C_BRAND),
        ("5", "Compartilhar Conquista",
         "Ao atingir 25% da meta:\nshare card nas redes sociais.",
         C_ACCENT),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = 0.3 + i * 2.55
        add_rect(slide, x, 1.4, 2.3, 3.8, color)
        add_text(slide, num,   x, 1.55, 2.3, 0.6,
                 font_size=28, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, title, x, 2.2,  2.3, 0.65,
                 font_size=13, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_rect(slide, x, 2.9, 2.3, 0.04, C_DARK)
        add_text(slide, desc,  x + 0.1, 3.05, 2.1, 1.9,
                 font_size=11, color=C_DARK, align=PP_ALIGN.CENTER)

    # Bottom: KPIs and qualitative work
    add_rect(slide, 0.3, 5.45, 6.0, 1.7, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "KPIs de Ativação (meta)", 0.5, 5.55, 5.6, 0.45,
             font_size=15, bold=True, color=C_ACCENT)
    kpis = [
        "• 40% completam onboarding em 7 dias",
        "• 60% retornam na semana 2",
        "• NPS > 50 antes de escalar marketing",
    ]
    add_multiline(slide, kpis, 0.5, 6.05, 5.6, 1.0,
                  font_size=12, color=C_LIGHT, spacing=1.3)

    add_rect(slide, 6.6, 5.45, 6.4, 1.7, RGBColor(0x05, 0x2A, 0x1A))
    add_text(slide, "Trabalho Qualitativo (primeiros 100 users)", 6.8, 5.55, 6.0, 0.45,
             font_size=15, bold=True, color=C_ACCENT)
    qual = [
        "• Entrevista semanal de 30 min (Jobs to be Done)",
        "• Canal de feedback direto (WhatsApp/Discord)",
        "• Iterar features a cada 2 semanas com base no feedback",
    ]
    add_multiline(slide, qual, 6.8, 6.05, 6.0, 1.0,
                  font_size=12, color=C_LIGHT, spacing=1.3)

    # ── Slide 11: Escala (1.000 → 10.000+) ───────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Escala — 1.000 → 10.000+ Usuários",
                 "Growth loops, B2B e paid acquisition")

    # Growth loops
    add_rect(slide, 0.4, 1.4, 5.8, 5.5, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "Growth Loops", 0.6, 1.5, 5.4, 0.45,
             font_size=17, bold=True, color=C_BRAND)
    loops = [
        "Upgrade loop:",
        "  Free → Essentials quando atingir limite de contas",
        "  Essentials → Premium pela demanda de IA/impostos",
        "",
        "Viral loop:",
        "  Compartilhar card \"Meta atingida!\" nas redes",
        "  Indicação: 1 mês grátis por amigo convertido",
        "",
        "Content loop:",
        "  Usuários compartilham insights gerados pelo app",
        "  Gera tráfego orgânico de volta ao produto",
    ]
    add_multiline(slide, loops, 0.6, 2.05, 5.4, 4.6,
                  font_size=12, color=C_LIGHT, spacing=1.3)

    # Right: channels at scale
    add_rect(slide, 6.5, 1.4, 6.5, 2.5, RGBColor(0x05, 0x2A, 0x1A))
    add_text(slide, "Canal B2B — Benefício Corporativo", 6.7, 1.5, 6.1, 0.45,
             font_size=15, bold=True, color=C_ACCENT)
    b2b = [
        "• RH de empresas oferece o app como benefício",
        "• Ticket: R$15/funcionário/mês (desconto de volume)",
        "• Distribuição imediata sem CAC individual",
        "• Parceiros iniciais: startups, consultorias, fintechs",
    ]
    add_multiline(slide, b2b, 6.7, 2.05, 6.1, 1.65,
                  font_size=12, color=C_LIGHT, spacing=1.3)

    add_rect(slide, 6.5, 4.15, 6.5, 2.75, RGBColor(0x1A, 0x07, 0x07))
    add_text(slide, "Paid Acquisition (após CAC validado)", 6.7, 4.25, 6.1, 0.45,
             font_size=15, bold=True, color=C_WARN)
    paid = [
        "• Meta Ads: segmentação por interesse em finanças",
        "• Google Search: keywords de dor (\"controlar gastos\")",
        "• Só ativar paid quando CAC orgânico < R$50 validado",
        "• Budget inicial: R$3.000–5.000/mês para testes A/B",
    ]
    add_multiline(slide, paid, 6.7, 4.75, 6.1, 2.0,
                  font_size=12, color=C_LIGHT, spacing=1.3)

    # Milestones bar
    milestones_gtm = [
        ("Mês 1–3",   "Beta fechado\n100 usuários", C_GRAY),
        ("Mês 4–6",   "Lançamento público\n500 free", C_BRAND),
        ("Mês 7–12",  "2.000 pagantes\nbreak-even", C_ACCENT),
        ("Mês 13–24", "10.000 pagantes\ntime contratado", C_WARN),
    ]
    # render as footer strip
    for i, (period, label, color) in enumerate(milestones_gtm):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 6.9, 3.0, 0.45, color)
        add_text(slide, f"{period}: {label.replace(chr(10), ' · ')}",
                 x + 0.07, 6.93, 2.85, 0.38,
                 font_size=10, bold=False, color=C_DARK)

    slide_header(slide, "Conclusão & Próximos Passos", "")

    conclusions = [
        ("Mercado validado", "Brazil fintech é o maior da LATAM com gap claro de produto integrado."),
        ("Concorrência fragmentada", "Nenhum player combina gastos + investimentos + educação."),
        ("Monetização comprovada", "YNAB prova que metodologia financeira é vendável. ~15–20% convertem."),
        ("Break-even realista", "2.000 usuários pagantes no modelo solo. Atingível em 12–18 meses."),
    ]
    next_steps = [
        "1. Definir escopo exato do MVP",
        "2. Validar monetização com landing page",
        "3. Prototipar telas principais",
        "4. Iniciar desenvolvimento backend",
    ]

    for i, (title, desc) in enumerate(conclusions):
        y = 1.4 + i * 1.3
        add_rect(slide, 0.4, y, 0.5, 0.9, C_ACCENT)
        add_text(slide, str(i + 1), 0.4, y + 0.1, 0.5, 0.7,
                 font_size=22, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.0, y, 7.5, 0.9, RGBColor(0x16, 0x23, 0x3A))
        add_text(slide, title, 1.15, y + 0.05, 7.2, 0.4,
                 font_size=15, bold=True, color=C_ACCENT)
        add_text(slide, desc,  1.15, y + 0.48, 7.2, 0.38,
                 font_size=12, color=C_LIGHT)

    add_rect(slide, 9.0, 1.4, 4.0, 5.1, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "Próximos Passos", 9.1, 1.5, 3.8, 0.5,
             font_size=16, bold=True, color=C_BRAND)
    add_multiline(slide, next_steps, 9.1, 2.1, 3.7, 4.0,
                  font_size=13, color=C_LIGHT, spacing=1.6)

    add_text(slide, f"FinanceApp  ·  Estudo de Mercado {VERSION_MARKET}  ·  {DATE}",
             0, 7.1, 13.33, 0.35, font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 2 — EXPENSES / COSTS
# ══════════════════════════════════════════════════════════════════════════════

def build_costs(prs: Presentation):

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    title_slide(
        prs,
        "Análise de Gastos",
        "Infraestrutura, Equipe, Projeções e Break-even",
        VERSION_COSTS, DATE, C_ACCENT
    )

    # ── Slide 2: Fases do Projeto ─────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Fases e Custo por Fase", "Evolução dos gastos ao longo do roadmap",
                 accent_color=C_ACCENT)

    phases = [
        ("MVP Solo",        "0–6 meses",   "R$400–1.600/mês",
         "Fundador dev, infra mínima, sem equipe.", C_BRAND),
        ("Lançamento",      "6–18 meses",  "R$45.000–75.000/mês",
         "2–3 devs + designer + marketing inicial.", C_ACCENT),
        ("Escala",          "18+ meses",   "R$103.000–211.000/mês",
         "5–8 pessoas + marketing acelerado.", C_WARN),
    ]

    for i, (phase, period, cost, desc, color) in enumerate(phases):
        y = 1.5 + i * 1.8
        add_rect(slide, 0.4, y, 12.5, 1.6, color)
        add_text(slide, phase,  0.6, y + 0.1,  3.0, 0.6,
                 font_size=22, bold=True, color=C_DARK)
        add_text(slide, period, 0.6, y + 0.75, 3.0, 0.5,
                 font_size=14, color=C_DARK)
        add_text(slide, cost,   3.8, y + 0.1,  4.5, 0.75,
                 font_size=26, bold=True, color=C_DARK)
        add_text(slide, desc,   3.8, y + 0.85, 9.0, 0.65,
                 font_size=13, color=C_DARK)

    # ── Slide 3: Infra MVP — com prestadores ─────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Prestadores de Serviço — Infra MVP",
                 "Opção recomendada, alternativas e custo estimado por componente",
                 accent_color=C_ACCENT)

    # (componente, prestador_principal, why, alternativas, custo)
    infra = [
        (
            "Backend API",
            "Railway",
            "Deploy via Git, zero config, escala automática",
            "Render · Fly.io",
            "R$60–150/mês",
        ),
        (
            "Banco de Dados (PostgreSQL)",
            "Railway Postgres",
            "Managed, backups diários, mesma rede do backend",
            "Neon · Supabase",
            "R$30–90/mês",
        ),
        (
            "Cache & Filas (Redis)",
            "Upstash",
            "Pay-per-request, free tier generoso, serverless",
            "Railway Redis · Redis Cloud",
            "R$0–60/mês",
        ),
        (
            "Frontend Web (Next.js)",
            "Vercel",
            "Deploy automático, CDN global, preview por PR",
            "Netlify · Cloudflare Pages",
            "R$0–90/mês",
        ),
        (
            "Storage (arquivos/avatares)",
            "Cloudflare R2",
            "Sem taxa de egress, S3-compatible, 10 GB free",
            "AWS S3 · Backblaze B2",
            "R$0–30/mês",
        ),
        (
            "Open Finance (aggregator)",
            "Pluggy",
            "Melhor cobertura BR, SDK pronto, suporte rápido",
            "Belvo · Quanto",
            "R$800–1.500/mês*",
        ),
        (
            "Gateway de Pagamento",
            "Asaas",
            "Foco BR, PIX + cartão, menor burocracia",
            "Stripe · Iugu",
            "~3% por transação",
        ),
        (
            "E-mail Transacional",
            "Resend",
            "API moderna, 100 e-mails/dia grátis, React Email",
            "SendGrid · Postmark",
            "R$0–80/mês",
        ),
    ]

    # Colunas: Componente | Prestador | Por quê | Alternativas | Custo
    col_defs = [
        ("Componente",    2.4),
        ("Prestador",     1.6),
        ("Por quê",       3.8),
        ("Alternativas",  2.2),
        ("Custo/mês",     1.9),
    ]
    col_x3 = [0.3]
    for _, w in col_defs[:-1]:
        col_x3.append(col_x3[-1] + w)

    row_h3 = 0.56
    header_y3 = 1.35
    add_rect(slide, 0.3, header_y3, 12.6, row_h3, C_ACCENT)
    for (h, w), x in zip(col_defs, col_x3):
        add_text(slide, h, x + 0.06, header_y3 + 0.08, w, row_h3 - 0.16,
                 font_size=12, bold=True, color=C_DARK)

    for i, (comp, provider, why, alts, cost) in enumerate(infra):
        y3 = header_y3 + row_h3 * (i + 1)
        bg = RGBColor(0x16, 0x23, 0x3A) if i % 2 == 0 else C_DARK
        add_rect(slide, 0.3, y3, 12.6, row_h3, bg)
        cells = [comp, provider, why, alts, cost]
        cell_colors = [C_WHITE, C_ACCENT, C_LIGHT, C_GRAY, C_WARN]
        for (_, w), x, cell, cc in zip(col_defs, col_x3, cells, cell_colors):
            add_text(slide, cell, x + 0.06, y3 + 0.07, w - 0.1, row_h3 - 0.14,
                     font_size=10, color=cc)

    total_y3 = header_y3 + row_h3 * (len(infra) + 1)
    add_rect(slide, 0.3, total_y3, 12.6, 0.48, C_BRAND)
    add_text(slide, "TOTAL INFRA (excl. Open Finance e Gateway)",
             0.45, total_y3 + 0.07, 9.5, 0.34, font_size=12, bold=True, color=C_DARK)
    add_text(slide, "R$90–420/mês",
             10.5, total_y3 + 0.07, 2.3, 0.34,
             font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.RIGHT)

    add_text(slide,
             "* Pluggy cobrado por conexão ativa/mês — custo só incide com usuários Open Finance ativos.",
             0.3, total_y3 + 0.55, 12.6, 0.35, font_size=10, color=C_GRAY)

    # ── Slide 4: Salários e Equipe ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Faixas Salariais — Mercado Brasil 2026", "Valores CLT (salário × 1,8 = custo real)",
                 accent_color=C_ACCENT)

    roles = [
        ("Dev Júnior",          "R$4.000–7.000",   "R$7.200–12.600"),
        ("Dev Pleno",           "R$8.000–13.000",  "R$14.400–23.400"),
        ("Dev Sênior",          "R$14.000–22.000", "R$25.200–39.600"),
        ("UX/UI Designer Jr",   "R$3.500–6.000",   "R$6.300–10.800"),
        ("UX/UI Designer Sr",   "R$10.000–16.000", "R$18.000–28.800"),
        ("Product Manager",     "R$8.000–25.000",  "R$14.400–45.000"),
    ]

    headers = ["Cargo", "Salário (CLT)", "Custo Real (×1,8)"]
    col_widths2 = [5.5, 3.2, 3.5]
    col_x2 = [0.4]
    for w in col_widths2[:-1]:
        col_x2.append(col_x2[-1] + w)

    add_rect(slide, 0.4, 1.5, 12.3, 0.5, C_ACCENT)
    for h, x, w in zip(headers, col_x2, col_widths2):
        add_text(slide, h, x + 0.1, 1.55, w, 0.38,
                 font_size=14, bold=True, color=C_DARK)

    for i, (role, salary, real) in enumerate(roles):
        y = 2.0 + i * 0.72
        bg = RGBColor(0x16, 0x23, 0x3A) if i % 2 == 0 else C_DARK
        add_rect(slide, 0.4, y, 12.3, 0.7, bg)
        add_text(slide, role,   col_x2[0] + 0.1, y + 0.12, col_widths2[0], 0.45, font_size=13)
        add_text(slide, salary, col_x2[1] + 0.1, y + 0.12, col_widths2[1], 0.45,
                 font_size=13, color=C_WARN)
        add_text(slide, real,   col_x2[2] + 0.1, y + 0.12, col_widths2[2], 0.45,
                 font_size=13, bold=True, color=C_ACCENT)

    add_text(slide,
             "Multiplicador CLT 1,8× inclui: FGTS, férias, 13º, INSS patronal, vale-transporte, plano de saúde.",
             0.4, 6.85, 12.3, 0.5, font_size=11, color=C_GRAY)

    # ── Slide 5: Break-even ────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Análise de Break-even", "Usuários pagantes vs. custo operacional",
                 accent_color=C_ACCENT)

    data_be = [
        (500,    "R$12.500",  "R$12.500",   "Não sustentável",    C_RED,    1.5),
        (2000,   "R$50.000",  "R$35.000",   "Viável — fase solo", C_ACCENT, 4.0),
        (5000,   "R$125.000", "R$75.000",   "Viável — time pequeno", C_BRAND, 6.5),
        (10000,  "R$250.000", "R$150.000",  "Viável — time médio", C_WARN,   9.5),
    ]

    add_rect(slide, 0.4, 1.45, 12.5, 0.5, C_ACCENT)
    for h, x in zip(["Usuários Pag.", "Receita (R$25/mês)", "Custo Estimado", "Status"],
                     [0.5, 2.9, 6.1, 9.4]):
        add_text(slide, h, x, 1.5, 3.0, 0.4, font_size=13, bold=True, color=C_DARK)

    for i, (users, rev, cost, status, color, _bar_w) in enumerate(data_be):
        y = 1.95 + i * 1.15
        bg = RGBColor(0x16, 0x23, 0x3A) if i % 2 == 0 else C_DARK
        add_rect(slide, 0.4, y, 12.5, 1.1, bg)
        add_text(slide, f"{users:,}",  0.5, y + 0.15, 2.2, 0.6,
                 font_size=22, bold=True, color=color)
        add_text(slide, rev,   2.9, y + 0.15, 3.0, 0.6, font_size=18, color=C_ACCENT)
        add_text(slide, cost,  6.1, y + 0.15, 3.0, 0.6, font_size=18, color=C_WARN)
        add_text(slide, status, 9.4, y + 0.2, 3.3, 0.5, font_size=13,
                 bold=True, color=color)

    # ── Slide 6: Economia de Unidade ──────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Economia de Unidade (Unit Economics)", "CAC · LTV · Payback",
                 accent_color=C_ACCENT)

    metrics = [
        ("CAC", "Custo de Aquisição\nde Cliente",
         "R$30–80", "Meta conservadora.\nOrg. reduz via conteúdo.", C_BRAND),
        ("LTV", "Lifetime Value\n(24–48 meses)",
         "R$600–1.800", "R$25/mês × 24–72 meses\nde retenção média.", C_ACCENT),
        ("LTV/CAC", "Relação\nLTV:CAC",
         "8–22×", "Acima de 3× é saudável.\nBenchmark SaaS: 4–8×.", C_WARN),
        ("Payback", "Tempo de\nRecuperação",
         "1,5–4 meses", "Muito eficiente.\nFluxo de caixa rápido.", C_BRAND),
    ]

    for i, (abbr, label, value, note, color) in enumerate(metrics):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 1.5, 3.0, 5.3, color)
        add_text(slide, abbr,  x, 1.65, 3.0, 0.7,
                 font_size=28, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, label, x, 2.35, 3.0, 0.8,
                 font_size=12, align=PP_ALIGN.CENTER, color=C_DARK)
        add_rect(slide, x, 3.2, 3.0, 0.04, C_DARK)
        add_text(slide, value, x, 3.35, 3.0, 0.9,
                 font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, note,  x, 4.3,  3.0, 2.0,
                 font_size=12, align=PP_ALIGN.CENTER, color=C_DARK)

    # ── Slide 7: Resumo Financeiro ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Resumo Financeiro & Recomendações", "",
                 accent_color=C_ACCENT)

    add_rect(slide, 0.4, 1.4, 5.8, 5.5, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, "Marcos Financeiros", 0.6, 1.5, 5.4, 0.5,
             font_size=17, bold=True, color=C_ACCENT)
    milestones = [
        "MVP: R$400–1.600/mês (infra apenas)",
        "Break-even: 2.000 usuários pagantes",
        "Receita mínima viável: R$50.000/mês",
        "Primeiro contrato: ~R$19,90 × 2.000 users",
        "Escala (time 5-8): 5.000+ users",
        "LTV/CAC target: > 8×",
    ]
    add_multiline(slide, ["• " + m for m in milestones],
                  0.6, 2.05, 5.5, 4.5, font_size=13, color=C_LIGHT, spacing=1.55)

    add_rect(slide, 6.6, 1.4, 6.5, 5.5, RGBColor(0x05, 0x3A, 0x1F))
    add_text(slide, "Recomendações", 6.8, 1.5, 6.1, 0.5,
             font_size=17, bold=True, color=C_ACCENT)
    recs = [
        "Iniciar solo para reduzir burn rate",
        "Usar Railway no MVP (sem DevOps overhead)",
        "Validar monetização antes de contratar",
        "Manter CAC < R$50 com conteúdo orgânico",
        "Priorizar retenção (onboarding 90 dias)",
        "Contratar 1º dev somente com 1.000 users",
    ]
    add_multiline(slide, ["✔ " + r for r in recs],
                  6.8, 2.05, 6.1, 4.5, font_size=13, color=C_LIGHT, spacing=1.55)

    add_text(slide, f"FinanceApp  ·  Análise de Gastos {VERSION_COSTS}  ·  {DATE}",
             0, 7.1, 13.33, 0.35, font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    os.makedirs(OUT_DIR, exist_ok=True)

    # Market research
    prs1 = new_prs()
    build_market_research(prs1)
    out1 = os.path.join(OUT_DIR, f"market-research-{VERSION_MARKET}.pptx")
    prs1.save(out1)
    print(f"Saved: {out1}")

    # Costs
    prs2 = new_prs()
    build_costs(prs2)
    out2 = os.path.join(OUT_DIR, f"costs-{VERSION_COSTS}.pptx")
    prs2.save(out2)
    print(f"Saved: {out2}")


if __name__ == "__main__":
    main()
