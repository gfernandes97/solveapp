"""
Solve — Brand Manual Generator
Cria logos PNG + brand manual PPTX.
Usage: python scripts/generate_brand.py
"""

import os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── BRAND COLORS ──────────────────────────────────────────────────────────────
# hex (for matplotlib)
HEX_OBS  = "#0F172A"    # Obsidian   — o problema, o desconhecido
HEX_GRN  = "#00C97A"    # Solve      — a resolução, o crescimento
HEX_SIG  = "#2563EB"    # Signal     — digital, confiança
HEX_SOL  = "#F59E0B"    # Solar      — energia, jovialidade
HEX_CLO  = "#F8FAFC"    # Cloud      — clareza, espaço
HEX_SMO  = "#64748B"    # Smoke      — texto de apoio
HEX_WHI  = "#FFFFFF"

# pptx RGBColor
C_OBS = RGBColor(0x0F, 0x17, 0x2A)
C_GRN = RGBColor(0x00, 0xC9, 0x7A)
C_SIG = RGBColor(0x25, 0x63, 0xEB)
C_SOL = RGBColor(0xF5, 0x9E, 0x0B)
C_CLO = RGBColor(0xF8, 0xFA, 0xFC)
C_SMO = RGBColor(0x64, 0x74, 0x8B)
C_WHI = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xEF, 0x44, 0x44)
C_GRN_DK = RGBColor(0x00, 0x8F, 0x57)   # darker green for text on light bg

VERSION = "v1.0.0"
DATE    = "2026-05-12"

ROOT     = os.path.join(os.path.dirname(__file__), "..")
LOGO_DIR = os.path.join(ROOT, "brand", "logos")
PPTX_DIR = os.path.join(ROOT, "docs", "presentations")

plt.rcParams["font.family"]     = "sans-serif"
plt.rcParams["font.sans-serif"] = ["Segoe UI", "Arial", "Helvetica", "DejaVu Sans"]


# ══════════════════════════════════════════════════════════════════════════════
#  LOGO CREATION (matplotlib → PNG)
# ══════════════════════════════════════════════════════════════════════════════

def _draw_icon(ax, bg, fg, cx=50, cy=50, r=30, sw=6.5):
    """
    Core icon: circle ring + ascending arrow emerging from inside.
    Coordinate space 0-100.
    """
    ax.set_facecolor(bg)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.set_aspect("equal")
    ax.axis("off")

    # ── Ring ──────────────────────────────────────────────────────────────────
    ring = mpatches.Circle((cx, cy), r, fill=False,
                           edgecolor=fg, linewidth=sw,
                           capstyle="round", zorder=3)
    ax.add_patch(ring)

    # ── Arrow ─────────────────────────────────────────────────────────────────
    # 45° (NE): tail inside circle, head outside circle
    ang = np.radians(45)
    tx  = cx - r * 0.38 * np.cos(ang)
    ty  = cy - r * 0.38 * np.sin(ang)
    hx  = cx + r * 1.42 * np.cos(ang)
    hy  = cy + r * 1.42 * np.sin(ang)

    ax.annotate("", xy=(hx, hy), xytext=(tx, ty),
                arrowprops=dict(
                    arrowstyle="->",
                    color=fg,
                    lw=sw * 0.72,
                    mutation_scale=32,
                    shrinkA=0,
                    shrinkB=0,
                    connectionstyle="arc3,rad=0",
                ), zorder=6)

    # ── Gap at ring exit point (arrow pierces the ring) ───────────────────────
    # Cover the ring at 45° with bg-colored arc → visual break
    gap = 9   # degrees each side
    ax.add_patch(mpatches.Arc(
        (cx, cy), r * 2, r * 2,
        angle=0, theta1=45 - gap, theta2=45 + gap,
        color=bg, linewidth=sw * 3.2, zorder=5,
    ))
    # Redraw arrow on top
    ax.annotate("", xy=(hx, hy), xytext=(tx, ty),
                arrowprops=dict(
                    arrowstyle="->",
                    color=fg,
                    lw=sw * 0.72,
                    mutation_scale=32,
                    shrinkA=0,
                    shrinkB=0,
                ), zorder=7)


def _save_icon(path, bg, fg, size_in=4):
    fig = plt.figure(figsize=(size_in, size_in), facecolor=bg)
    ax  = fig.add_axes([0.06, 0.06, 0.88, 0.88])
    _draw_icon(ax, bg, fg)
    fig.savefig(path, dpi=300, facecolor=bg, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)


def _save_lockup_h(path, bg, icon_fg, word_color, size=(10, 3)):
    """Horizontal lockup: icon left + wordmark right."""
    fig = plt.figure(figsize=size, facecolor=bg)

    # Icon axes (left square portion)
    ax_icon = fig.add_axes([0.03, 0.08, 0.28, 0.84])
    _draw_icon(ax_icon, bg, icon_fg, cx=50, cy=50, r=30, sw=7)

    # Wordmark axes (right portion)
    ax_word = fig.add_axes([0.34, 0.1, 0.63, 0.8])
    ax_word.set_xlim(0, 10)
    ax_word.set_ylim(0, 3)
    ax_word.axis("off")
    ax_word.set_facecolor(bg)
    ax_word.text(0, 1.5, "solve", fontsize=72, fontweight="bold",
                 color=word_color, va="center", ha="left",
                 fontstyle="normal")

    fig.savefig(path, dpi=300, facecolor=bg, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)


def _save_lockup_v(path, bg, icon_fg, word_color, size=(4, 5.5)):
    """Vertical lockup: icon above + wordmark below."""
    fig = plt.figure(figsize=size, facecolor=bg)

    ax_icon = fig.add_axes([0.10, 0.38, 0.80, 0.58])
    _draw_icon(ax_icon, bg, icon_fg, cx=50, cy=50, r=30, sw=7)

    ax_word = fig.add_axes([0.05, 0.02, 0.90, 0.33])
    ax_word.set_xlim(0, 10)
    ax_word.set_ylim(0, 3)
    ax_word.axis("off")
    ax_word.set_facecolor(bg)
    ax_word.text(5, 1.5, "solve", fontsize=60, fontweight="bold",
                 color=word_color, va="center", ha="center")

    fig.savefig(path, dpi=300, facecolor=bg, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)


def _save_wordmark(path, bg, color, size=(8, 2)):
    """Wordmark only."""
    fig = plt.figure(figsize=size, facecolor=bg)
    ax  = fig.add_axes([0.02, 0.05, 0.96, 0.90])
    ax.set_xlim(0, 10); ax.set_ylim(0, 3)
    ax.axis("off"); ax.set_facecolor(bg)
    ax.text(0.3, 1.5, "solve", fontsize=80, fontweight="bold",
            color=color, va="center", ha="left")
    ax.text(0.5, 0.3, "Agora você sabe.", fontsize=22,
            color=color, va="bottom", ha="left", alpha=0.55)
    fig.savefig(path, dpi=300, facecolor=bg, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)


def create_all_logos():
    os.makedirs(LOGO_DIR, exist_ok=True)
    paths = {}

    # V1 — Icon green on white
    p = os.path.join(LOGO_DIR, "icon-green-white.png")
    _save_icon(p, HEX_WHI, HEX_GRN)
    paths["icon_gw"] = p

    # V2 — Icon white on obsidian
    p = os.path.join(LOGO_DIR, "icon-white-dark.png")
    _save_icon(p, HEX_OBS, HEX_WHI)
    paths["icon_wd"] = p

    # V3 — Icon white on green
    p = os.path.join(LOGO_DIR, "icon-white-green.png")
    _save_icon(p, HEX_GRN, HEX_WHI)
    paths["icon_wg"] = p

    # V4 — Primary lockup: H, white bg
    p = os.path.join(LOGO_DIR, "primary-white.png")
    _save_lockup_h(p, HEX_WHI, HEX_GRN, HEX_OBS)
    paths["primary_w"] = p

    # V5 — Primary lockup: H, dark bg
    p = os.path.join(LOGO_DIR, "primary-dark.png")
    _save_lockup_h(p, HEX_OBS, HEX_WHI, HEX_WHI)
    paths["primary_d"] = p

    # V6 — Stacked lockup, white bg
    p = os.path.join(LOGO_DIR, "stacked-white.png")
    _save_lockup_v(p, HEX_WHI, HEX_GRN, HEX_OBS)
    paths["stacked_w"] = p

    # V7 — Wordmark with tagline, white bg
    p = os.path.join(LOGO_DIR, "wordmark-white.png")
    _save_wordmark(p, HEX_WHI, HEX_OBS)
    paths["wordmark_w"] = p

    # V8 — Wordmark with tagline, dark bg
    p = os.path.join(LOGO_DIR, "wordmark-dark.png")
    _save_wordmark(p, HEX_OBS, HEX_WHI)
    paths["wordmark_d"] = p

    print(f"  Logos criados em {LOGO_DIR}")
    return paths


# ══════════════════════════════════════════════════════════════════════════════
#  PPTX HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=14, bold=False,
      color=C_OBS, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def ML(slide, items, l, t, w, h, sz=12, color=C_OBS, bold=False, sp=1.35):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(sz * (sp - 1))
        run = p.add_run(); run.text = item
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Segoe UI"

def img(slide, path, l, t, w):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def swatch(slide, l, t, w, h, hex_color, name, hex_txt, use):
    """Color swatch block."""
    r_val = int(hex_color[1:3], 16)
    g_val = int(hex_color[3:5], 16)
    b_val = int(hex_color[5:7], 16)
    col   = RGBColor(r_val, g_val, b_val)
    R(slide, l, t, w, h * 0.58, col)
    R(slide, l, t + h * 0.58, w, h * 0.42, C_WHI)
    lum = 0.299 * r_val + 0.587 * g_val + 0.114 * b_val
    txt_c = C_WHI if lum < 160 else C_OBS
    T(slide, hex_color.upper(), l + 0.08, t + h * 0.58 - 0.28, w - 0.12, 0.26,
      sz=9, bold=True, color=txt_c, align=PP_ALIGN.LEFT)
    T(slide, name, l + 0.08, t + h * 0.60, w - 0.12, 0.28,
      sz=11, bold=True, color=C_OBS)
    T(slide, use, l + 0.08, t + h * 0.78, w - 0.12, 0.32,
      sz=9, color=C_SMO)

def footer(slide, txt=""):
    R(slide, 0, 7.35, 13.33, 0.15, C_GRN)
    if txt:
        T(slide, txt, 0.25, 7.35, 10.0, 0.15, sz=8, color=C_OBS)
    T(slide, f"Solve · Manual de Marca {VERSION} · {DATE}",
      9.0, 7.35, 4.2, 0.15, sz=8, color=C_OBS, align=PP_ALIGN.RIGHT)

def chapter_tag(slide, text, l=0.3, t=0.18):
    R(slide, l, t, len(text) * 0.105 + 0.25, 0.26, C_GRN)
    T(slide, text, l + 0.1, t + 0.04, len(text) * 0.105 + 0.1, 0.20,
      sz=9, bold=True, color=C_OBS)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs, logos):
    slide = blank(prs)
    bg(slide, C_OBS)

    # Full-bleed green strip at bottom
    R(slide, 0, 6.8, 13.33, 0.70, C_GRN)
    T(slide, "Manual de Marca", 0.5, 6.87, 5.0, 0.40,
      sz=13, bold=True, color=C_OBS)
    T(slide, f"{VERSION}  ·  {DATE}", 9.5, 6.90, 3.6, 0.30,
      sz=11, color=C_OBS, align=PP_ALIGN.RIGHT)

    # Primary logo (white on dark)
    img(slide, logos["primary_d"], 2.0, 1.6, 9.3)

    # Tagline
    T(slide, "Agora você sabe.", 1.0, 5.5, 11.33, 0.8,
      sz=28, bold=False, color=C_GRN, align=PP_ALIGN.CENTER)

    # Subtle divider
    R(slide, 3.5, 5.32, 6.33, 0.04, RGBColor(0x1E, 0x2D, 0x42))


def slide_naming(prs, logos):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Nomenclatura & Missão")
    footer(slide)

    # Left: story
    T(slide, 'Por que\n“Solve”?', 0.5, 0.55, 6.5, 1.4,
      sz=38, bold=True, color=C_OBS)

    ML(slide, [
        "O brasileiro ganha, gasta, e no fim do mês não sabe o que "
        "aconteceu com o dinheiro. Sem visibilidade real dos seus gastos, "
        "ele não consegue sair do vermelho, montar uma reserva, "
        "nem investir para o futuro.",
        "",
        "Solve é a solução direta a esse problema.",
        "",
        "Uma palavra. Universal. Sem floreios.",
        "Fácil de pronunciar em português e inglês.",
        "Diz exatamente o que fazemos.",
    ], 0.5, 2.05, 6.2, 4.5, sz=13, color=C_SMO, sp=1.45)

    # Right: wordmark large + tagline explanation
    R(slide, 7.2, 0.5, 5.9, 6.7, C_OBS)
    img(slide, logos["primary_d"], 7.4, 1.4, 5.5)

    T(slide, "\"Agora você sabe.\"", 7.4, 3.5, 5.5, 0.55,
      sz=18, bold=True, color=C_GRN, align=PP_ALIGN.CENTER)

    ML(slide, [
        "A tagline responde diretamente ao pain point:",
        "\"O brasileiro não sabe para onde vai seu dinheiro.\"",
        "",
        "Com o Solve, agora ele sabe.",
        "Três palavras. Uma promessa inteira.",
    ], 7.5, 4.15, 5.3, 2.6, sz=11, color=RGBColor(0x94, 0xA3, 0xB8), sp=1.35)


def slide_logo_concept(prs, logos):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Logo — Conceito")
    footer(slide)

    T(slide, "O que o símbolo comunica", 0.5, 0.55, 12.5, 0.58,
      sz=28, bold=True, color=C_OBS)
    T(slide, "Cada elemento tem significado direto sobre o problema que resolvemos",
      0.5, 1.15, 12.5, 0.35, sz=13, color=C_SMO)

    # Large icon, centered
    img(slide, logos["icon_gw"], 4.5, 1.6, 4.33)

    # Annotation LEFT: the ring
    R(slide, 0.3, 3.2, 3.8, 1.8, RGBColor(0xF1, 0xF5, 0xF9))
    R(slide, 0.3, 3.2, 0.06, 1.8, C_SMO)
    T(slide, "O círculo", 0.5, 3.3, 3.5, 0.35,
      sz=13, bold=True, color=C_OBS)
    ML(slide, [
        "Representa o ciclo do não saber —",
        "o dinheiro some todo mês e a",
        "pessoa não entende por quê.",
        "O loop de invisibilidade financeira.",
    ], 0.5, 3.7, 3.5, 1.2, sz=11, color=C_SMO, sp=1.30)

    # Annotation RIGHT: the arrow
    R(slide, 9.2, 3.2, 3.8, 1.8, RGBColor(0xF0, 0xFD, 0xF4))
    R(slide, 9.2, 3.2, 0.06, 1.8, C_GRN)
    T(slide, "A seta", 9.4, 3.3, 3.5, 0.35,
      sz=13, bold=True, color=C_OBS)
    ML(slide, [
        "Nasce dentro do ciclo e rompe",
        "para fora, apontando para cima.",
        "Representa o solve: direção,",
        "clareza e movimento consciente.",
    ], 9.4, 3.7, 3.5, 1.2, sz=11, color=C_SMO, sp=1.30)

    # Bottom: the combined message
    R(slide, 0.3, 5.25, 12.73, 0.90, C_OBS)
    T(slide, "Juntos: \"do ciclo de não saber, emerge a direção.\"",
      0.5, 5.38, 12.3, 0.55,
      sz=18, bold=True, color=C_GRN, align=PP_ALIGN.CENTER)

    # Horizontal connector lines (annotation arrows to icon)
    R(slide, 4.1, 4.05, 0.4, 0.04, C_SMO)
    R(slide, 8.83, 4.05, 0.4, 0.04, C_GRN)


def slide_logo_variants(prs, logos):
    slide = blank(prs)
    bg(slide, C_CLO)
    chapter_tag(slide, "Logo — Variantes")
    footer(slide)

    T(slide, "As 5 variantes aprovadas", 0.5, 0.55, 10.0, 0.55,
      sz=28, bold=True, color=C_OBS)

    # ── Row 1: Primary (white bg) + Primary (dark bg) ─────────────────────────
    # V1: white bg
    R(slide, 0.3, 1.25, 6.1, 2.35, C_WHI)
    img(slide, logos["primary_w"], 0.5, 1.35, 5.7)
    T(slide, "V1 — Principal  ·  fundo claro",
      0.35, 3.30, 6.0, 0.28, sz=10, color=C_SMO)

    # V2: dark bg
    R(slide, 6.75, 1.25, 6.25, 2.35, C_OBS)
    img(slide, logos["primary_d"], 6.95, 1.35, 5.85)
    T(slide, "V2 — Inverso  ·  fundo escuro",
      6.8, 3.30, 6.0, 0.28, sz=10, color=C_SMO)

    # ── Row 2: Icon only | Stacked | Wordmark ─────────────────────────────────
    # V3: Icon only
    R(slide, 0.3, 3.75, 3.0, 3.35, C_WHI)
    img(slide, logos["icon_gw"], 0.65, 3.90, 2.30)
    T(slide, "V3 — Ícone puro",
      0.35, 6.87, 3.0, 0.28, sz=10, color=C_SMO, align=PP_ALIGN.CENTER)

    # V4: Stacked
    R(slide, 3.65, 3.75, 3.05, 3.35, C_WHI)
    img(slide, logos["stacked_w"], 3.75, 3.82, 2.85)
    T(slide, "V4 — Empilhado",
      3.7, 6.87, 3.0, 0.28, sz=10, color=C_SMO, align=PP_ALIGN.CENTER)

    # V5: Wordmark only
    R(slide, 7.05, 3.75, 5.95, 3.35, C_OBS)
    img(slide, logos["wordmark_d"], 7.15, 4.35, 5.75)
    T(slide, "V5 — Wordmark + tagline",
      7.1, 6.87, 5.9, 0.28, sz=10, color=C_SMO, align=PP_ALIGN.CENTER)


def slide_logo_rules(prs, logos):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Logo — Regras de Uso")
    footer(slide)

    T(slide, "O que respeitar", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)

    # DO column
    R(slide, 0.3, 1.25, 6.0, 0.40, C_GRN)
    T(slide, "✓  Pode", 0.45, 1.30, 5.7, 0.30, sz=14, bold=True, color=C_OBS)
    dos = [
        ("Usar em fundo branco, obsidian ou verde Solve.", C_WHI, 0.0),
        ("Manter a proporção original. Sempre.", C_WHI, 0.0),
        ("Usar versão monocromática preta em impressão P&B.", C_WHI, 0.0),
        ("Aplicar com área de proteção ≥ altura da letra 'S'.", C_WHI, 0.0),
        ("Usar V3 (ícone puro) em espaços pequenos (16px+).", C_WHI, 0.0),
    ]
    for i, (txt, bg_c, _) in enumerate(dos):
        y = 1.72 + i * 0.82
        R(slide, 0.3, y, 6.0, 0.72, RGBColor(0xF0, 0xFD, 0xF4) if i % 2 == 0 else C_WHI)
        R(slide, 0.3, y, 0.06, 0.72, C_GRN)
        T(slide, txt, 0.5, y + 0.17, 5.7, 0.38, sz=12, color=C_OBS)

    # DON'T column
    R(slide, 6.9, 1.25, 6.1, 0.40, C_RED)
    T(slide, "✕  Não pode", 7.05, 1.30, 5.7, 0.30, sz=14, bold=True, color=C_WHI)
    donts = [
        "Distorcer ou esticar o logo.",
        "Alterar as cores fora da paleta oficial.",
        "Adicionar sombra, brilho ou efeitos ao ícone.",
        "Usar sobre fundos fotográficos ou padrões.",
        "Recriar o ícone em ferramentas de apresentação.",
    ]
    for i, txt in enumerate(donts):
        y = 1.72 + i * 0.82
        R(slide, 6.9, y, 6.1, 0.72, RGBColor(0xFE, 0xF2, 0xF2) if i % 2 == 0 else C_WHI)
        R(slide, 6.9, y, 0.06, 0.72, C_RED)
        T(slide, txt, 7.1, y + 0.17, 5.7, 0.38, sz=12, color=C_OBS)

    # Clear space rule
    R(slide, 0.3, 6.65, 12.7, 0.55, RGBColor(0xF8, 0xFA, 0xFC))
    T(slide, 'Área de proteção: o logo deve ter ao redor um espaço livre equivalente à altura da letra "S" do wordmark em todas as direções.',
      0.5, 6.70, 12.4, 0.42, sz=10, color=C_SMO)


def slide_colors(prs):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Paleta de Cores")
    footer(slide)

    T(slide, "As cores da marca", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)
    T(slide, "Cada cor tem um papel. Obsidian e Solve são a dupla principal.",
      0.5, 1.12, 12.0, 0.30, sz=13, color=C_SMO)

    palette = [
        (HEX_OBS, "Obsidian",  "Cor principal escura — o problema, a profundidade"),
        (HEX_GRN, "Solve",     "Cor principal — resolução, crescimento, vida"),
        (HEX_SIG, "Signal",    "Secundária — digital, links, ações"),
        (HEX_SOL, "Solar",     "Acento — energia, jovialidade, alertas"),
        (HEX_CLO, "Cloud",     "Background — clareza, espaço, respiro"),
        (HEX_SMO, "Smoke",     "Texto de apoio — neutro, hierarquia"),
    ]

    sw_w = (13.33 - 0.6) / len(palette) - 0.12
    for i, (hex_c, name, use) in enumerate(palette):
        x = 0.3 + i * (sw_w + 0.12)
        swatch(slide, x, 1.55, sw_w, 5.4, hex_c, name, hex_c.upper(), use)

    # Ratio guidance
    R(slide, 0.3, 7.0, 12.73, 0.3, RGBColor(0xF1, 0xF5, 0xF9))
    for ratio, color, lbl in [
        (0.40, C_OBS, "Obsidian  40%"),
        (0.30, C_GRN, "Solve  30%"),
        (0.30, C_CLO, "Cloud  30%"),
    ]:
        pass  # simplified


def slide_color_combinations(prs, logos):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Combinações de Cores")
    footer(slide)

    T(slide, "Combinações aprovadas", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)

    combos = [
        # (bg_hex, text_hex, label, usage)
        (HEX_OBS, HEX_GRN,  "Obsidian + Solve",
         "Principal — capas, headers,\ndestaque máximo"),
        (HEX_WHI, HEX_OBS,  "Branco + Obsidian",
         "Textos longos, corpo de\ntexto, fundos de conteúdo"),
        (HEX_GRN, HEX_OBS,  "Solve + Obsidian",
         "CTAs primários, badges,\ndestaques de sucesso"),
        (HEX_OBS, HEX_WHI,  "Obsidian + Branco",
         "Navegação, sidebars,\nrodapés e modais"),
        (HEX_SOL, HEX_OBS,  "Solar + Obsidian",
         "Alertas de atenção,\nwarnings, onboarding"),
        (HEX_CLO, HEX_OBS,  "Cloud + Obsidian",
         "Cards, backgrounds\nde conteúdo secundário"),
    ]

    cw = (13.33 - 0.6) / 3 - 0.12
    ch = 2.85

    for i, (bg_h, fg_h, lbl, use) in enumerate(combos):
        col, row = i % 3, i // 3
        x = 0.3 + col * (cw + 0.12)
        y = 1.4 + row * (ch + 0.18)

        br = int(bg_h[1:3], 16); bg2 = int(bg_h[3:5], 16); bb = int(bg_h[5:7], 16)
        fr = int(fg_h[1:3], 16); fg2 = int(fg_h[3:5], 16); fb = int(fg_h[5:7], 16)

        R(slide, x, y, cw, ch, RGBColor(br, bg2, bb))

        T(slide, "solve", x + 0.2, y + 0.25, cw - 0.35, 0.65,
          sz=30, bold=True, color=RGBColor(fr, fg2, fb))
        T(slide, "Agora você sabe.", x + 0.2, y + 0.95, cw - 0.35, 0.30,
          sz=11, color=RGBColor(
              int((fr + br) // 2), int((fg2 + bg2) // 2), int((fb + bb) // 2)
          ))
        T(slide, lbl, x + 0.2, y + 1.45, cw - 0.35, 0.28,
          sz=10, bold=True, color=RGBColor(fr, fg2, fb))
        T(slide, use, x + 0.2, y + 1.78, cw - 0.35, 0.75,
          sz=9, color=RGBColor(
              min(fr + 40, 255), min(fg2 + 40, 255), min(fb + 40, 255)
          ))


def slide_typography(prs):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Tipografia")
    footer(slide)

    T(slide, "Inter — a tipografia da marca", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)
    T(slide, "Geométrica, humanista, criada para interfaces digitais. Google Fonts — gratuita.",
      0.5, 1.12, 12.0, 0.30, sz=13, color=C_SMO)

    # Left: hierarchy
    R(slide, 0.3, 1.55, 6.8, 5.6, RGBColor(0xF8, 0xFA, 0xFC))
    R(slide, 0.3, 1.55, 0.06, 5.6, C_GRN)

    styles = [
        ("Display — 48/56pt Bold",    "Agora você sabe.", 28, True),
        ("H1 — 32/40pt Bold",         "Onde foi o seu dinheiro?", 18, True),
        ("H2 — 24/32pt SemiBold",     "Gastos de maio 2026", 14, True),
        ("Body — 16/24pt Regular",    "Você gastou R$4.230 este mês em 127 transações.", 11, False),
        ("Caption — 12/16pt Regular", "Fonte: Open Finance · Atualizado agora", 9, False),
        ("Label — 11pt Bold Caps",    "VER FATURA →", 9, True),
    ]
    y = 1.68
    for style_label, sample, sz, bold in styles:
        T(slide, style_label, 0.5, y, 6.5, 0.22, sz=8, color=C_SMO)
        T(slide, sample, 0.5, y + 0.22, 6.5, 0.52, sz=sz, bold=bold, color=C_OBS)
        y += 0.85

    # Right: usage principles
    R(slide, 7.5, 1.55, 5.6, 5.6, C_OBS)
    T(slide, "Princípios tipográficos", 7.7, 1.68, 5.2, 0.35,
      sz=14, bold=True, color=C_GRN)

    principles = [
        ("Hierarquia clara",
         "Nunca mais de 3 tamanhos em uma mesma tela. O tamanho comunica importância."),
        ("Contraste alto",
         "Texto principal sempre em Obsidian sobre Cloud ou branco. Nunca cinza sobre cinza."),
        ("Peso como emoção",
         "Bold para impacto e dados críticos. Regular para explicação. Nunca Thin."),
        ("Espaçamento generoso",
         "Line-height mínimo de 1.5× o tamanho da fonte. Respiro é clareza."),
    ]
    py = 2.15
    for title, desc in principles:
        R(slide, 7.5, py, 5.6, 1.15, RGBColor(0x1A, 0x24, 0x38))
        T(slide, title, 7.68, py + 0.08, 5.2, 0.28, sz=11, bold=True, color=C_GRN)
        T(slide, desc, 7.68, py + 0.40, 5.2, 0.60, sz=10, color=RGBColor(0x94, 0xA3, 0xB8))
        py += 1.22


def slide_voice(prs):
    slide = blank(prs)
    bg(slide, C_WHI)
    chapter_tag(slide, "Voz & Tom")
    footer(slide)

    T(slide, "Como o Solve fala", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)

    # 3 words
    words = [
        (C_GRN, "Claro",
         "Zero jargão financeiro.\nSe precisar de um termo técnico, explica na mesma frase."),
        (C_SIG, "Direto",
         "A frase começa pelo dado que importa.\nNunca enrolação, nunca falsa intimidade."),
        (C_SOL, "Humano",
         "Reconhece que finanças têm emoção.\nNão julga. Encoraja. Celebra pequenas vitórias."),
    ]
    for i, (col, word, desc) in enumerate(words):
        x = 0.3 + i * 4.35
        R(slide, x, 1.35, 4.1, 2.55, col)
        T(slide, word, x, 1.45, 4.1, 0.85,
          sz=38, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)
        T(slide, desc, x + 0.15, 2.38, 3.8, 1.35, sz=12, color=C_OBS)

    # Before / After
    R(slide, 0.3, 4.10, 12.73, 0.40, C_OBS)
    T(slide, "Antes × Depois — o mesmo dado, o tom certo",
      0.5, 4.17, 12.3, 0.28, sz=13, bold=True, color=C_GRN)

    befores = [
        ("Seu saldo devedor no cartão de crédito é de R$ 1.847,50 "
         "com vencimento em 14/05/2026.",
         "Sua fatura Nubank vence em 3 dias. São R$1.847,50. "
         "Quer quitar agora e economizar R$14.290 em juros?"),
        ("Detectamos um padrão de gastos acima da média no "
         "segmento delivery nas últimas 4 semanas.",
         "Você pediu delivery 31 vezes este mês. Isso deu R$487. "
         "Na semana passada foram 9 pedidos — quer entender o porquê?"),
    ]
    for i, (before, after) in enumerate(befores):
        y = 4.65 + i * 1.35
        R(slide, 0.3, y, 6.0, 1.18, RGBColor(0xFE, 0xF2, 0xF2))
        R(slide, 0.3, y, 0.06, 1.18, C_RED)
        T(slide, "✕ Antes", 0.45, y + 0.05, 5.6, 0.22, sz=9, bold=True, color=C_RED)
        T(slide, before, 0.45, y + 0.30, 5.6, 0.80, sz=10, color=C_OBS)
        R(slide, 6.65, y, 6.35, 1.18, RGBColor(0xF0, 0xFD, 0xF4))
        R(slide, 6.65, y, 0.06, 1.18, C_GRN)
        T(slide, "✓ Depois (Solve)", 6.8, y + 0.05, 6.0, 0.22, sz=9, bold=True, color=C_GRN_DK)
        T(slide, after, 6.8, y + 0.30, 6.0, 0.80, sz=10, color=C_OBS)


def slide_tagline(prs, logos):
    slide = blank(prs)
    bg(slide, C_OBS)
    footer(slide)

    # Full-slide tagline display
    R(slide, 0, 2.5, 13.33, 0.07, C_GRN)
    T(slide, "Agora você sabe.", 0.5, 0.8, 12.33, 2.1,
      sz=64, bold=True, color=C_WHI, align=PP_ALIGN.CENTER)
    T(slide, "A tagline é a anti-tese do pain point.",
      1.0, 2.8, 11.33, 0.45, sz=18, color=C_GRN, align=PP_ALIGN.CENTER)

    # Split explanation
    left_items = [
        "O pain point:",
        "\"O brasileiro NÃO SABE",
        " para onde vai seu dinheiro.\"",
        "",
        "82,8 milhões com CPF negativado.",
        "48% sem controle dos gastos.",
        "31% sem reserva de emergência.",
    ]
    ML(slide, left_items, 0.5, 3.45, 5.8, 3.5,
       sz=13, color=RGBColor(0x94, 0xA3, 0xB8), sp=1.3)

    R(slide, 6.5, 3.45, 0.04, 3.5, RGBColor(0x1E, 0x2D, 0x42))

    right_items = [
        "A promessa do Solve:",
        "\"AGORA VOCÊ SABE",
        " exatamente para onde foi.\"",
        "",
        "Dashboard que responde em <3 seg.",
        "Alerta antes de você errar.",
        "Educação no momento da ação.",
    ]
    ML(slide, right_items, 6.9, 3.45, 5.9, 3.5,
       sz=13, color=C_GRN, sp=1.3)


def slide_brand_in_use(prs, logos):
    slide = blank(prs)
    bg(slide, C_CLO)
    chapter_tag(slide, "Marca em Uso")
    footer(slide)

    T(slide, "O Solve no mundo real", 0.5, 0.55, 12.0, 0.55,
      sz=28, bold=True, color=C_OBS)

    # App icon simulation
    R(slide, 0.3, 1.25, 2.8, 5.85, C_WHI)
    T(slide, "App Icon", 0.3, 1.28, 2.8, 0.28, sz=10, color=C_SMO, align=PP_ALIGN.CENTER)
    # 3 sizes
    for sz_in, y_pos, label in [(1.8, 1.65, "1024px"), (1.0, 3.65, "180px"), (0.5, 4.85, "60px")]:
        x_c = 0.3 + (2.8 - sz_in) / 2
        img(slide, logos["icon_wg"], x_c, y_pos, sz_in)
        T(slide, label, 0.3, y_pos + sz_in + 0.05, 2.8, 0.22,
          sz=9, color=C_SMO, align=PP_ALIGN.CENTER)

    # Notification card
    R(slide, 3.5, 1.25, 4.6, 2.70, C_WHI)
    T(slide, "Notificação Push", 3.55, 1.28, 4.5, 0.28, sz=10, color=C_SMO)
    R(slide, 3.6, 1.60, 4.3, 2.20, RGBColor(0xF8, 0xFA, 0xFC))
    img(slide, logos["icon_gw"], 3.65, 1.65, 0.50)
    T(slide, "solve", 4.22, 1.70, 1.5, 0.28, sz=10, bold=True, color=C_OBS)
    T(slide, "agora", 4.22, 1.98, 1.0, 0.18, sz=8, color=C_SMO)
    T(slide, "⚠  Sua fatura vence amanhã.",
      3.65, 2.30, 4.1, 0.28, sz=11, bold=True, color=C_OBS)
    T(slide, "R$ 1.847,50 · Nubank · Toque para ver",
      3.65, 2.60, 4.1, 0.24, sz=10, color=C_SMO)
    T(slide, "Quitar agora →",
      3.65, 2.95, 1.8, 0.24, sz=10, bold=True, color=C_GRN_DK)
    T(slide, "Ver detalhes",
      5.5, 2.95, 2.2, 0.24, sz=10, color=C_SMO)

    # Social media card (Instagram story mockup)
    R(slide, 3.5, 4.15, 4.6, 2.95, C_OBS)
    R(slide, 3.5, 4.15, 4.6, 0.06, C_GRN)
    img(slide, logos["icon_wd"], 3.6, 4.25, 0.55)
    T(slide, "solve", 4.22, 4.33, 2.0, 0.30, sz=12, bold=True, color=C_WHI)
    T(slide, "Você gastou R$ 4.230 em maio.", 3.6, 4.75, 4.3, 0.30,
      sz=13, bold=True, color=C_WHI)
    T(slide, "Quer saber onde foi cada centavo?",
      3.6, 5.10, 4.3, 0.30, sz=11, color=RGBColor(0xCB, 0xD5, 0xE1))
    R(slide, 3.65, 5.52, 2.5, 0.30, C_GRN)
    T(slide, "Agora você sabe →", 3.65, 5.55, 2.5, 0.25,
      sz=10, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)
    T(slide, "@usesolve.com.br", 3.6, 6.90, 4.3, 0.22,
      sz=9, color=C_SMO, align=PP_ALIGN.CENTER)

    # Email header
    R(slide, 8.5, 1.25, 4.6, 5.85, C_WHI)
    T(slide, "E-mail Transacional", 8.55, 1.28, 4.5, 0.28, sz=10, color=C_SMO)
    R(slide, 8.5, 1.58, 4.6, 1.25, C_OBS)
    img(slide, logos["primary_d"], 8.65, 1.68, 3.0)
    T(slide, "Relatório de maio", 8.6, 2.92, 4.3, 0.30,
      sz=13, bold=True, color=C_OBS)
    T(slide, "Ana, em maio você gastou R$4.230. Aqui está o resumo:",
      8.6, 3.28, 4.3, 0.40, sz=10, color=C_SMO)

    rows_e = [("Alimentação", "R$1.189", C_SOL),
              ("Moradia",     "R$1.058", C_SIG),
              ("Transporte",  "R$635",   C_GRN)]
    for j, (cat, val, col) in enumerate(rows_e):
        ey = 3.80 + j * 0.48
        R(slide, 8.6, ey, 4.2, 0.40, RGBColor(0xF8, 0xFA, 0xFC))
        R(slide, 8.6, ey, 0.05, 0.40, col)
        T(slide, cat, 8.72, ey + 0.08, 2.5, 0.24, sz=10, color=C_OBS)
        T(slide, val, 10.8, ey + 0.08, 1.9, 0.24,
          sz=10, bold=True, color=C_OBS, align=PP_ALIGN.RIGHT)

    R(slide, 8.6, 5.28, 4.2, 0.38, C_GRN)
    T(slide, "Ver relatório completo →", 8.6, 5.34, 4.2, 0.26,
      sz=11, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)
    T(slide, "Agora você sabe.", 8.6, 5.80, 4.2, 0.26,
      sz=10, color=C_SMO, align=PP_ALIGN.CENTER)


def slide_principles(prs, logos):
    slide = blank(prs)
    bg(slide, C_OBS)
    footer(slide)

    img(slide, logos["icon_wd"], 5.9, 0.3, 1.5)
    T(slide, "Os 4 princípios da marca", 1.0, 1.85, 11.33, 0.55,
      sz=24, bold=True, color=C_WHI, align=PP_ALIGN.CENTER)

    principles = [
        (C_GRN,  "Clareza antes de tudo",
         "Se o usuário não entendeu em 3 segundos, é culpa nossa, não dele."),
        (C_SIG,  "Dados com contexto",
         "Um número sem contexto é ruído. Um dado com ação é poder."),
        (C_SOL,  "Educação no momento",
         "A lição financeira deve aparecer quando a decisão está sendo tomada."),
        (C_RED if False else RGBColor(0xEC, 0x48, 0x99),
         "Sem julgamento",
         "Finanças têm emoção. A marca acolhe, encoraja e celebra o progresso — não o perfeito."),
    ]
    for i, (col, title, desc) in enumerate(principles):
        col_i, row_i = i % 2, i // 2
        x = 0.5 + col_i * 6.3
        y = 2.55 + row_i * 2.25
        R(slide, x, y, 6.0, 2.0, RGBColor(0x16, 0x23, 0x3A))
        R(slide, x, y, 0.08, 2.0, col)
        T(slide, title, x + 0.25, y + 0.18, 5.6, 0.42, sz=16, bold=True, color=C_WHI)
        T(slide, desc, x + 0.25, y + 0.68, 5.6, 0.80, sz=12, color=RGBColor(0x94, 0xA3, 0xB8))

    T(slide, "Solve — Agora você sabe.", 0, 7.0, 13.33, 0.35,
      sz=14, bold=True, color=C_GRN, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build_manual(logos):
    prs = new_prs()
    slide_cover(prs, logos)
    slide_naming(prs, logos)
    slide_logo_concept(prs, logos)
    slide_logo_variants(prs, logos)
    slide_logo_rules(prs, logos)
    slide_colors(prs)
    slide_color_combinations(prs, logos)
    slide_typography(prs)
    slide_voice(prs)
    slide_tagline(prs, logos)
    slide_brand_in_use(prs, logos)
    slide_principles(prs, logos)
    return prs


def main():
    os.makedirs(LOGO_DIR, exist_ok=True)
    os.makedirs(PPTX_DIR, exist_ok=True)

    print("Gerando logos...")
    logos = create_all_logos()

    print("Gerando brand manual...")
    prs = build_manual(logos)

    out = os.path.join(PPTX_DIR, f"brand-manual-{VERSION}.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
    print(f"Logos: {LOGO_DIR}")


if __name__ == "__main__":
    main()
