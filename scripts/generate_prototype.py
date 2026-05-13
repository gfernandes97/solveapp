"""
FinanceApp — Protótipo Web MVP
Wireframes para apresentação — 14 slides
Usage: python scripts/generate_prototype.py
Output: docs/presentations/prototype-v1.0.0.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Presentation palette ───────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
C_BRAND  = RGBColor(0x00, 0x84, 0xFF)
C_ACCENT = RGBColor(0x00, 0xD4, 0x8A)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_WARN   = RGBColor(0xF5, 0x9E, 0x0B)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_ORANGE = RGBColor(0xF9, 0x73, 0x16)
C_PINK   = RGBColor(0xEC, 0x48, 0x99)
C_TEAL   = RGBColor(0x14, 0xB8, 0xA6)

# ── Wireframe light-theme palette ─────────────────────────────────────────────
WF_BG    = RGBColor(0xF1, 0xF5, 0xF9)   # page background
WF_CARD  = RGBColor(0xFF, 0xFF, 0xFF)   # card / panel
WF_BORD  = RGBColor(0xE2, 0xE8, 0xF0)   # borders / dividers
WF_NAV   = RGBColor(0x0F, 0x17, 0x2A)   # top nav bar
WF_SIDE  = RGBColor(0x1E, 0x29, 0x3B)   # sidebar
WF_T1    = RGBColor(0x0F, 0x17, 0x2A)   # primary text
WF_T2    = RGBColor(0x64, 0x74, 0x8B)   # secondary text
WF_GRN   = RGBColor(0x05, 0xB0, 0x72)   # positive / income
WF_RED   = RGBColor(0xDC, 0x26, 0x26)   # negative / expense
WF_AMB   = RGBColor(0xB4, 0x5A, 0x09)   # warning text
WF_AMBB  = RGBColor(0xFE, 0xF3, 0xC7)   # warning background
WF_REDB  = RGBColor(0xFE, 0xE2, 0xE2)   # red background
WF_GRNB  = RGBColor(0xD1, 0xFA, 0xE5)   # green background
WF_BLUE  = RGBColor(0x00, 0x84, 0xFF)   # actions / links
WF_PURP  = RGBColor(0x8B, 0x5C, 0xF6)   # investment

VERSION = "v1.0.0"
DATE    = "2026-05-11"
OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations")

# ── Layout constants (all in Inches) ─────────────────────────────────────────
SL_W = 13.33
SL_H = 7.50

WF_L  = 0.20    # wireframe left edge
WF_T  = 1.05    # wireframe top edge (below slide header)
WF_W  = 12.93   # wireframe width
WF_H  = 5.45    # wireframe height  (bottom at 6.50)
CHR_H = 0.33    # browser chrome height
CT_T  = WF_T + CHR_H    # web content top
CT_H  = WF_H - CHR_H    # web content height

NAV_H = 0.38    # app top-nav height
SB_W  = 1.85    # sidebar width
MC_X  = WF_L + SB_W     # main content left (with sidebar)
MC_W  = WF_W - SB_W     # main content width
MC_T  = CT_T + NAV_H    # main content top
MC_H  = CT_H - NAV_H    # main content height

ANN_T = WF_T + WF_H + 0.05   # annotation strip top
ANN_H = 0.60


# ══════════════════════════════════════════════════════════════════════════════
#  PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(SL_W)
    prs.slide_height = Inches(SL_H)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def R(slide, l, t, w, h, color):
    """Rectangle"""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=10, bold=False,
      color=C_WHITE, align=PP_ALIGN.LEFT):
    """Text box"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def ML(slide, items, l, t, w, h, sz=10, color=C_WHITE, bold=False, sp=1.3):
    """Multi-line text box"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(sz * (sp - 1))
        run = p.add_run(); run.text = item
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
#  WIREFRAME COMPONENTS
# ══════════════════════════════════════════════════════════════════════════════

def chrome(slide, url="financeapp.com.br"):
    """Browser chrome bar"""
    R(slide, WF_L, WF_T, WF_W, WF_H, WF_BG)
    R(slide, WF_L, WF_T, WF_W, CHR_H, RGBColor(0xD1, 0xD5, 0xDB))
    for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFF,0xBD,0x2E), RGBColor(0x28,0xC8,0x40)]):
        R(slide, WF_L + 0.10 + i*0.20, WF_T + 0.10, 0.13, 0.13, c)
    R(slide, WF_L + 0.82, WF_T + 0.06, WF_W - 1.35, 0.21, WF_CARD)
    T(slide, url, WF_L + 0.87, WF_T + 0.07, WF_W - 1.45, 0.19,
      sz=8, color=WF_T2, align=PP_ALIGN.CENTER)

def app_nav(slide, items=None, active=None):
    """Dark app top-nav"""
    if items is None:
        items = ["Início", "Gastos", "Investimentos", "Metas", "Educação"]
    R(slide, WF_L, CT_T, WF_W, NAV_H, WF_NAV)
    T(slide, "FinanceApp", WF_L+0.15, CT_T+0.06, 1.5, NAV_H-0.12,
      sz=11, bold=True, color=C_BRAND)
    ix = WF_L + 1.85
    for item in items:
        col = C_ACCENT if item == active else RGBColor(0xCB,0xD5,0xE1)
        T(slide, item, ix, CT_T+0.09, 1.2, NAV_H-0.18, sz=9, color=col, align=PP_ALIGN.CENTER)
        ix += 1.35
    # Avatar
    R(slide, WF_L+WF_W-0.48, CT_T+0.06, 0.26, 0.26, C_ACCENT)
    T(slide, "AS", WF_L+WF_W-0.48, CT_T+0.07, 0.26, 0.22,
      sz=8, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    # Notif bell
    T(slide, "🔔", WF_L+WF_W-0.82, CT_T+0.07, 0.28, 0.24,
      sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def app_sidebar(slide, items=None, active=None):
    """Dark sidebar"""
    if items is None:
        items = ["Dashboard", "Gastos", "Investimentos", "Metas", "Relatórios", "Configurações"]
    R(slide, WF_L, MC_T, SB_W, MC_H, WF_SIDE)
    R(slide, WF_L, MC_T, SB_W, 0.60, RGBColor(0x14,0x1D,0x2E))
    R(slide, WF_L+0.15, MC_T+0.12, 0.34, 0.34, C_ACCENT)
    T(slide, "AS", WF_L+0.15, MC_T+0.13, 0.34, 0.30,
      sz=8, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    T(slide, "Ana Souza", WF_L+0.56, MC_T+0.11, SB_W-0.65, 0.22,
      sz=9, bold=True, color=C_WHITE)
    T(slide, "Saúde: 68/100", WF_L+0.56, MC_T+0.34, SB_W-0.65, 0.20,
      sz=8, color=C_ACCENT)
    y = MC_T + 0.68
    for item in items:
        is_act = (item == active)
        if is_act:
            R(slide, WF_L, y, SB_W, 0.34, C_BRAND)
        T(slide, ("▸ " if is_act else "  ") + item,
          WF_L+0.15, y+0.07, SB_W-0.2, 0.24,
          sz=9, bold=is_act, color=C_WHITE if is_act else WF_T2)
        y += 0.37

def land_nav(slide):
    """Landing page nav"""
    R(slide, WF_L, CT_T, WF_W, NAV_H, WF_NAV)
    T(slide, "FinanceApp", WF_L+0.2, CT_T+0.06, 1.6, NAV_H-0.12,
      sz=11, bold=True, color=C_BRAND)
    for i, lnk in enumerate(["Como funciona", "Para quem", "Preços"]):
        T(slide, lnk, WF_L+2.6+i*1.7, CT_T+0.09, 1.5, NAV_H-0.18,
          sz=9, color=C_LIGHT)
    R(slide, WF_L+WF_W-1.55, CT_T+0.07, 1.35, 0.24, C_BRAND)
    T(slide, "Começar grátis", WF_L+WF_W-1.55, CT_T+0.09, 1.35, 0.20,
      sz=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def wbtn(slide, l, t, w, text, bg=C_BRAND, tc=C_WHITE, sz=9):
    """Web button"""
    R(slide, l, t, w, 0.27, bg)
    T(slide, text, l, t+0.05, w, 0.20, sz=sz, bold=True,
      color=tc, align=PP_ALIGN.CENTER)

def wcard(slide, l, t, w, h, bg=None):
    R(slide, l, t, w, h, WF_CARD if bg is None else bg)

def hbars(slide, l, t, w, h, segs, lw=1.4, vw=0.85):
    """
    Horizontal bar chart.
    segs: [(label, value_str, pct_0_to_100, bar_color), ...]
    """
    n = len(segs)
    bh = h / n
    baw = w - lw - vw - 0.1
    for i, (lbl, val, pct, col) in enumerate(segs):
        y = t + i * bh
        T(slide, lbl, l, y+0.03, lw, bh-0.06, sz=8, color=WF_T1)
        R(slide, l+lw, y+0.08, baw, bh-0.18, WF_BORD)
        R(slide, l+lw, y+0.08, max(0.04, baw*pct/100), bh-0.18, col)
        T(slide, val, l+lw+baw+0.05, y+0.03, vw, bh-0.06,
          sz=8, bold=True, color=WF_T1, align=PP_ALIGN.RIGHT)

def stat_box(slide, l, t, w, h, val, lbl, sub=None, val_col=C_BRAND):
    """Metric card"""
    wcard(slide, l, t, w, h)
    T(slide, val, l+0.1, t+0.08, w-0.2, 0.42, sz=17, bold=True, color=val_col)
    T(slide, lbl, l+0.1, t+0.52, w-0.2, 0.20, sz=8, color=WF_T2)
    if sub:
        T(slide, sub, l+0.1, t+0.72, w-0.2, 0.20, sz=8, color=val_col)

def progress_bar(slide, l, t, w, h, pct, bg=WF_BORD, fill=C_BRAND):
    R(slide, l, t, w, h, bg)
    R(slide, l, t, max(0.04, w*pct/100), h, fill)

def ann(slide, items):
    """
    Annotation strip: items = [(icon, title, body, color), ...]
    """
    R(slide, 0.0, ANN_T, SL_W, ANN_H, RGBColor(0x16, 0x23, 0x3A))
    iw = SL_W / len(items)
    for i, (icon, title, body, col) in enumerate(items):
        x = i * iw
        R(slide, x, ANN_T, 0.06, ANN_H, col)
        T(slide, f"{icon}  {title}", x+0.15, ANN_T+0.06, iw-0.2, 0.24,
          sz=10, bold=True, color=col)
        T(slide, body, x+0.15, ANN_T+0.32, iw-0.2, 0.24, sz=9, color=C_GRAY)

def ftag(slide, text, l, t, bg=C_BRAND, tc=C_WHITE, sz=8):
    """Small floating tag/badge"""
    est_w = len(text)*0.076 + 0.25
    R(slide, l, t, est_w, 0.22, bg)
    T(slide, text, l+0.06, t+0.03, est_w-0.12, 0.18,
      sz=sz, bold=True, color=tc, align=PP_ALIGN.CENTER)
    return l + est_w + 0.1


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE FRAMING
# ══════════════════════════════════════════════════════════════════════════════

def pres_header(slide, title, subtitle="", accent=C_ORANGE):
    bg(slide, C_DARK)
    R(slide, 0, 0, SL_W, 0.07, accent)
    R(slide, 0, SL_H-0.07, SL_W, 0.07, accent)
    T(slide, title, 0.2, 0.12, 10.5, 0.60,
      sz=24, bold=True, color=C_WHITE)
    if subtitle:
        T(slide, subtitle, 0.2, 0.72, 12.0, 0.30,
          sz=11, color=C_GRAY)

def wf_header(slide, num, title, personas=None, pain=None):
    """Header for wireframe slides — compact"""
    bg(slide, C_DARK)
    R(slide, 0, 0, SL_W, 0.07, C_ORANGE)
    R(slide, 0, SL_H-0.07, SL_W, 0.07, C_ORANGE)
    # Screen number badge
    R(slide, 0.20, 0.13, 0.56, 0.75, C_ORANGE)
    T(slide, num, 0.20, 0.16, 0.56, 0.68,
      sz=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    T(slide, title, 0.86, 0.13, 8.5, 0.50,
      sz=20, bold=True, color=C_WHITE)
    # Persona and pain badges (right side)
    bx = 9.6
    if personas:
        for p in personas:
            R(slide, bx, 0.20, 0.80, 0.25, C_PURPLE)
            T(slide, p, bx+0.05, 0.23, 0.72, 0.20,
              sz=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            bx += 0.88
    if pain:
        R(slide, bx, 0.20, len(pain)*0.08+0.3, 0.25, C_RED)
        T(slide, pain, bx+0.06, 0.23, len(pain)*0.08+0.22, 0.20,
          sz=8, bold=True, color=C_WHITE)

def footer_tag(slide):
    T(slide, f"FinanceApp · Protótipo Web MVP {VERSION} · {DATE}",
      0, 7.12, SL_W, 0.32, sz=9, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = blank(prs)
    bg(slide, C_DARK)
    R(slide, 0, 0, SL_W, 0.12, C_ORANGE)
    R(slide, 0, SL_H-0.12, SL_W, 0.12, C_ORANGE)
    R(slide, 0, 3.05, SL_W, 0.04, C_ORANGE)
    T(slide, "FinanceApp", 1, 0.45, 11.33, 0.8,
      sz=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    T(slide, "Protótipo Web MVP", 1, 1.45, 11.33, 1.3,
      sz=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(slide, "Wireframes  ·  Fluxo de telas  ·  Ataque ao pain point central",
      1, 3.22, 11.33, 0.65, sz=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    ML(slide,
       ["Baseado em: Estudo de Personas v1.0.0  ·  Pain Point: "
        "\"O brasileiro não sabe para onde vai seu dinheiro\"",
        "8 perfis · 10 dores rankeadas · 3 pilares do problema"],
       1, 4.1, 11.33, 0.9, sz=12, color=C_GRAY, sp=1.4)
    T(slide, f"{VERSION}  ·  {DATE}",
      1, 6.85, 11.33, 0.40, sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ÂNCORA: PROBLEM STATEMENT → SOLUÇÃO
# ══════════════════════════════════════════════════════════════════════════════

def slide_anchor(prs):
    slide = blank(prs)
    pres_header(slide, "O Problema que Vamos Resolver",
                "Cada tela deste protótipo existe para atacar um desses três pilares")

    # Problem statement box
    R(slide, 0.3, 1.15, 12.73, 1.55, RGBColor(0x1A, 0x07, 0x07))
    R(slide, 0.3, 1.15, 0.07, 1.55, C_RED)
    T(slide, "O brasileiro ganha, gasta, e no fim do mês não sabe o que aconteceu com o dinheiro.",
      0.55, 1.25, 12.2, 0.72, sz=20, bold=True, color=C_WHITE)
    T(slide, "Sem visibilidade real dos seus gastos, ele não consegue sair do vermelho, "
             "montar uma reserva, nem investir para o futuro — mesmo quando quer.",
      0.55, 1.97, 12.2, 0.60, sz=13, color=C_LIGHT)

    # 3 pillars → solution
    pillars = [
        (C_RED, "Pilar 1", "Falta de\nVisibilidade",
         "Não enxerga onde o\ndinheiro vai em tempo real.",
         "Dashboard central que\nresponde AGORA: \"onde foi?\"",
         "Telas 3, 4"),
        (C_WARN, "Pilar 2", "Falta de\nOrientação",
         "Sabe que tem problema\nmas não sabe o próximo passo.",
         "Alertas inteligentes,\nplano de saída, metas.",
         "Telas 5, 6"),
        (C_ACCENT, "Pilar 3", "Falta do\nMomento Certo",
         "Educação financeira chega\ntarde, fora do contexto.",
         "Insights inline, educação\nno momento da ação.",
         "Telas 4, 7"),
    ]

    for i, (col, tag, title, problem, solution, screens) in enumerate(pillars):
        x = 0.3 + i * 4.25
        R(slide, x, 2.9, 4.0, 4.3, col)
        T(slide, tag, x, 2.97, 4.0, 0.30,
          sz=10, color=C_DARK, align=PP_ALIGN.CENTER)
        T(slide, title, x, 3.27, 4.0, 0.65,
          sz=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        R(slide, x, 3.92, 4.0, 0.03, C_DARK)
        T(slide, "Problema:", x+0.15, 4.00, 3.7, 0.25,
          sz=9, bold=True, color=C_DARK)
        T(slide, problem, x+0.15, 4.26, 3.7, 0.65,
          sz=11, color=C_DARK)
        R(slide, x, 4.95, 4.0, 0.03, RGBColor(0x00,0x00,0x00))
        T(slide, "Solução no protótipo:", x+0.15, 5.03, 3.7, 0.25,
          sz=9, bold=True, color=C_DARK)
        T(slide, solution, x+0.15, 5.29, 3.7, 0.65,
          sz=11, color=C_DARK)
        R(slide, x, 6.0, 4.0, 0.35, RGBColor(0x00,0x00,0x00))
        T(slide, screens, x, 6.05, 4.0, 0.27,
          sz=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — ARQUITETURA DE TELAS
# ══════════════════════════════════════════════════════════════════════════════

def slide_arch(prs):
    slide = blank(prs)
    pres_header(slide, "Arquitetura de Telas — MVP Web",
                "7 telas · fluxo do usuário · personas atendidas por bloco")

    screens = [
        ("T1", "Landing\nPage",     C_ORANGE, 0.25, 1.3, "Todos"),
        ("T2", "Onboarding\nDiagnóstico", C_WARN,   3.5, 1.3, "P1–P5"),
        ("T3", "Dashboard\n«Onde foi?»",  C_RED,    6.75, 1.3, "P1–P5, P8"),
        ("T4", "Análise de\nGastos",      C_BRAND,  0.25, 4.2, "P1–P5, P8"),
        ("T5", "Alerta\nInteligente",     C_RED,    3.5,  4.2, "P2, P1, P5"),
        ("T6", "Metas e\nReserva",        C_ACCENT, 6.75, 4.2, "P1, P3, P4"),
        ("T7", "Educação\nContextual",    C_PURPLE, 10.0, 1.3, "P1–P4"),
    ]

    for code, name, col, x, y, pers in screens:
        R(slide, x, y, 3.0, 2.6, col)
        T(slide, code, x, y+0.08, 3.0, 0.45,
          sz=20, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        R(slide, x, y+0.53, 3.0, 0.03, C_DARK)
        T(slide, name, x, y+0.60, 3.0, 0.75,
          sz=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        T(slide, pers, x, y+1.55, 3.0, 0.35,
          sz=10, color=C_DARK, align=PP_ALIGN.CENTER)
        R(slide, x, y+1.95, 3.0, 0.55, RGBColor(0,0,0))
        T(slide, "◀ clique para ver tela ▶", x, y+2.05, 3.0, 0.35,
          sz=9, color=C_DARK, align=PP_ALIGN.CENTER)

    # Arrows: T1→T2→T3
    for ax, ay in [(3.25, 2.6), (6.5, 2.6), (9.75, 2.6)]:
        T(slide, "→", ax, ay, 0.4, 0.35, sz=16, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # "Core loop" label
    R(slide, 0.25, 3.95, 10.5, 0.18, RGBColor(0x16, 0x23, 0x3A))
    T(slide, "▶  Loop de uso contínuo (abre app → vê dashboard → age)",
      0.4, 3.96, 10.2, 0.17, sz=9, color=C_ORANGE)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — T1: LANDING PAGE
# ══════════════════════════════════════════════════════════════════════════════

def slide_landing(prs):
    slide = blank(prs)
    wf_header(slide, "T1", "Landing Page — Entrada Pública",
              ["P1","P2","P3","P4","P5"], "Visibilidade")
    chrome(slide, "financeapp.com.br")
    land_nav(slide)

    # Hero section (dark background)
    hero_t = CT_T + NAV_H
    hero_h = 2.5
    R(slide, WF_L, hero_t, WF_W, hero_h, WF_NAV)

    # Problem statement headline
    T(slide, "Você sabe para onde foi o seu dinheiro este mês?",
      WF_L+0.4, hero_t+0.2, 8.0, 0.85,
      sz=20, bold=True, color=C_WHITE)
    T(slide, "48% dos brasileiros não controlam seus gastos. "
             "82,8 milhões têm CPF negativado. "
             "O problema começa sempre no mesmo lugar.",
      WF_L+0.4, hero_t+1.08, 7.8, 0.60,
      sz=10, color=RGBColor(0xCB,0xD5,0xE1))

    # CTAs
    wbtn(slide, WF_L+0.4, hero_t+1.82, 1.8,
         "Ver meu diagnóstico grátis", C_ACCENT, WF_T1, sz=9)
    wbtn(slide, WF_L+2.3, hero_t+1.82, 1.3,
         "Ver como funciona", RGBColor(0x2D,0x3A,0x4E), C_WHITE, sz=9)

    # Right stats panel
    R(slide, WF_L+WF_W-4.2, hero_t+0.15, 4.0, 2.2, RGBColor(0x1A,0x24,0x38))
    for i, (n, label, col) in enumerate([
        ("82,8M", "com CPF negativado", C_RED),
        ("80,9%", "das famílias endividadas", C_WARN),
        ("48%",   "sem controle dos gastos", C_ORANGE),
    ]):
        y = hero_t + 0.25 + i * 0.65
        T(slide, n, WF_L+WF_W-4.05, y, 1.3, 0.42, sz=18, bold=True, color=col)
        T(slide, label, WF_L+WF_W-2.65, y+0.1, 2.5, 0.28, sz=9, color=WF_T2)

    # "Como funciona" section
    how_t = hero_t + hero_h + 0.1
    R(slide, WF_L, how_t, WF_W, 0.28, WF_CARD)
    T(slide, "Como funciona em 3 passos",
      WF_L+0.3, how_t+0.05, 5.0, 0.22, sz=11, bold=True, color=WF_T1)

    step_cols = [C_BRAND, C_ACCENT, C_ORANGE]
    step_titles = ["Conecte seu banco", "Veja o diagnóstico", "Tome o controle"]
    step_descs = [
        "Open Finance: conecta em 2 minutos.\nSem digitar uma senha.",
        "Em segundos veja onde\nseu dinheiro foi este mês.",
        "Metas, alertas e plano\nde ação personalizado.",
    ]
    step_w = WF_W / 3 - 0.2
    for i, (col, title, desc) in enumerate(zip(step_cols, step_titles, step_descs)):
        x = WF_L + i * (step_w + 0.25) + 0.1
        y = how_t + 0.33
        R(slide, x, y, step_w, 1.8, WF_CARD)
        R(slide, x, y, step_w, 0.05, col)
        R(slide, x+0.15, y+0.12, 0.3, 0.3, col)
        T(slide, str(i+1), x+0.15, y+0.12, 0.3, 0.28,
          sz=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        T(slide, title, x+0.55, y+0.16, step_w-0.65, 0.26,
          sz=10, bold=True, color=WF_T1)
        T(slide, desc, x+0.15, y+0.52, step_w-0.25, 0.65,
          sz=9, color=WF_T2)
        wbtn(slide, x+0.15, y+1.48, step_w-0.3, "Saiba mais →",
             WF_BORD, WF_T2, sz=8)

    ann(slide, [
        ("🎯", "Pain Point atacado", "Visibilidade — mostra que o problema existe ANTES do login", C_RED),
        ("👤", "Personas",           "Todos os perfis · entrada universal do produto", C_PURPLE),
        ("✅", "Princípio",          "Problem Statement visível no hero em <3 segundos de leitura", C_ACCENT),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — T2: ONBOARDING / DIAGNÓSTICO
# ══════════════════════════════════════════════════════════════════════════════

def slide_onboarding(prs):
    slide = blank(prs)
    wf_header(slide, "T2", "Onboarding — Diagnóstico Imediato",
              ["P1","P2","P3","P4","P5"], "Visibilidade")
    chrome(slide, "financeapp.com.br/entrar")
    land_nav(slide)

    cont_t = CT_T + NAV_H

    # Step progress bar
    R(slide, WF_L, cont_t, WF_W, 0.30, WF_CARD)
    steps = ["Criar conta", "Conectar banco", "Ver diagnóstico", "Definir meta"]
    sw = WF_W / len(steps)
    for i, s in enumerate(steps):
        col = C_ACCENT if i < 2 else (C_BRAND if i == 2 else WF_BORD)
        R(slide, WF_L + i*sw, cont_t, sw, 0.30, col)
        T(slide, f"{i+1}. {s}", WF_L+i*sw+0.1, cont_t+0.06, sw-0.15, 0.20,
          sz=8, bold=(i == 2), color=C_DARK if i <= 2 else WF_T2,
          align=PP_ALIGN.CENTER)

    # Main content: split layout
    left_x = WF_L + 0.2
    left_w = 5.5
    right_x = WF_L + 6.2
    right_w = WF_W - 6.4

    # LEFT: Connect bank step
    wcard(slide, left_x, cont_t+0.38, left_w, 4.55)
    T(slide, "Conecte seu banco em 2 minutos",
      left_x+0.2, cont_t+0.52, left_w-0.35, 0.45,
      sz=14, bold=True, color=WF_T1)
    T(slide, "Via Open Finance — o mesmo padrão do Banco Central. "
             "Seus dados ficam criptografados e você pode revogar o acesso a qualquer hora.",
      left_x+0.2, cont_t+1.00, left_w-0.35, 0.60,
      sz=9, color=WF_T2)

    # Bank logos (placeholder boxes)
    banks = ["Nubank", "Itaú", "Bradesco", "Caixa", "Inter", "Santander"]
    for j, bank in enumerate(banks):
        bx = left_x + 0.2 + (j % 3) * 1.72
        by = cont_t + 1.75 + (j // 3) * 0.7
        R(slide, bx, by, 1.55, 0.55, WF_BORD)
        T(slide, bank, bx, by+0.14, 1.55, 0.26,
          sz=9, color=WF_T1, align=PP_ALIGN.CENTER)

    wbtn(slide, left_x+0.2, cont_t+3.45, left_w-0.4,
         "Conectar via Open Finance →", C_BRAND, C_WHITE)
    T(slide, "🔒  Conexão segura · Banco Central · Revogar a qualquer momento",
      left_x+0.2, cont_t+3.82, left_w-0.35, 0.25,
      sz=8, color=WF_T2, align=PP_ALIGN.CENTER)

    # RIGHT: Diagnóstico result (the "aha moment")
    wcard(slide, right_x, cont_t+0.38, right_w, 4.55, RGBColor(0x05,0x1A,0x10))
    R(slide, right_x, cont_t+0.38, right_w, 0.05, C_ACCENT)
    T(slide, "Seu diagnóstico de abril",
      right_x+0.2, cont_t+0.5, right_w-0.3, 0.30,
      sz=11, bold=True, color=C_ACCENT)
    T(slide, "Em abril você gastou...",
      right_x+0.2, cont_t+0.88, right_w-0.3, 0.28,
      sz=10, color=RGBColor(0xCB,0xD5,0xE1))
    T(slide, "R$ 3.847,00",
      right_x+0.2, cont_t+1.18, right_w-0.3, 0.65,
      sz=28, bold=True, color=C_WHITE)
    T(slide, "em 127 transações  ·  Você sabia disso?",
      right_x+0.2, cont_t+1.88, right_w-0.3, 0.26,
      sz=9, color=WF_T2)

    # Mini bars for the diagnosis
    diag_segs = [
        ("Alimentação", "R$1.077", 28, C_ORANGE),
        ("Moradia",     "R$962",   25, C_BRAND),
        ("Transporte",  "R$577",   15, C_TEAL),
        ("Lazer",       "R$462",   12, C_PURPLE),
        ("Cartão",      "R$769",   20, C_RED),
    ]
    hbars(slide, right_x+0.2, cont_t+2.25, right_w-0.3, 1.55,
          diag_segs, lw=1.0, vw=0.75)

    T(slide, "💡  Você gasta 23% acima da média no delivery",
      right_x+0.2, cont_t+3.92, right_w-0.3, 0.25,
      sz=9, color=C_WARN)
    wbtn(slide, right_x+0.2, cont_t+4.22, right_w-0.35,
         "Entender e definir minha 1ª meta →", C_ACCENT, C_DARK)

    ann(slide, [
        ("⚡", "Aha moment",  "Usuário vê onde foi seu dinheiro em <30 seg após conectar o banco", C_ACCENT),
        ("🎯", "Dor atacada", "Dor #1: não saber para onde o dinheiro vai — resolvida no onboarding", C_RED),
        ("👤", "Personas",    "P1 Jovem · P2 Endividado · P3 Iniciante · P4 Mulher · P5 Profissional", C_PURPLE),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — T3: DASHBOARD "ONDE FOI MEU DINHEIRO?" (core)
# ══════════════════════════════════════════════════════════════════════════════

def slide_dashboard(prs):
    slide = blank(prs)
    wf_header(slide, "T3", "Dashboard — «Onde foi meu dinheiro?»  ★ CORE",
              ["P1","P2","P3","P4","P5","P8"], "Dor #1")
    chrome(slide, "app.financeapp.com.br/dashboard")
    app_nav(slide, active="Início")
    app_sidebar(slide, active="Dashboard")

    # ── MAIN CONTENT AREA ────────────────────────────────────────────────────
    pad = 0.12
    mx = MC_X + pad
    mw = MC_W - pad*2

    # Alert banner (red)
    R(slide, MC_X, MC_T, MC_W, 0.28, WF_REDB)
    T(slide, "⚠  Fatura Nubank vence em 3 dias · R$ 1.847,50",
      mx, MC_T+0.04, mw*0.75, 0.22, sz=9, color=WF_RED)
    wbtn(slide, MC_X+MC_W-1.55, MC_T+0.02, 1.4, "Ver fatura →",
         WF_RED, C_WHITE, sz=8)

    # THE CENTRAL QUESTION CARD
    qy = MC_T + 0.33
    R(slide, MC_X, qy, MC_W, 0.90, WF_NAV)
    R(slide, MC_X, qy, 0.06, 0.90, C_RED)
    T(slide, "Em maio, você gastou...",
      mx, qy+0.06, 5.0, 0.26, sz=10, color=RGBColor(0xCB,0xD5,0xE1))
    T(slide, "R$ 4.230,00",
      mx, qy+0.32, 4.5, 0.50, sz=28, bold=True, color=C_WHITE)
    T(slide, "de R$ 6.000 disponíveis  ·  11 dias restantes  ·  R$ 161/dia para manter o ritmo",
      mx, qy+0.80, 7.5, 0.20, sz=8, color=WF_T2)

    # Progress bar
    R(slide, mx, qy+0.98, mw*0.55, 0.12, WF_BORD)
    R(slide, mx, qy+0.98, mw*0.55*0.705, 0.12, C_RED)
    T(slide, "70% do orçamento usado",
      mx+mw*0.57, qy+0.98, mw*0.38, 0.16, sz=8, color=WF_AMB)

    # 4 stat cards
    sy = qy + 1.18
    sh = 0.88
    sw4 = mw / 4 - 0.08
    stats = [
        ("R$ 1.770", "Disponível", "≈ R$161/dia", WF_GRN),
        ("R$ 6.000", "Entradas maio", "", WF_BLUE),
        ("R$ 4.230", "Gastos maio", "▲ 8% vs abr", WF_RED),
        ("R$ 850",   "Investido", "Meta: R$1.200", WF_PURP),
    ]
    for i, (val, lbl, sub, col) in enumerate(stats):
        sx = MC_X + i * (sw4+0.1)
        stat_box(slide, sx, sy, sw4, sh, val, lbl, sub or None, col)

    # 2-column section
    col_t = sy + sh + 0.1
    col_h = MC_H - (col_t - MC_T) - 0.05
    lw2 = mw * 0.52
    rw2 = mw - lw2 - 0.1

    # LEFT: Where did it go?
    wcard(slide, MC_X, col_t, lw2+pad, col_h)
    T(slide, "Onde foi?", mx, col_t+0.08, lw2, 0.26,
      sz=11, bold=True, color=WF_T1)
    cat_segs = [
        ("Alimentação", "R$1.189", 28, C_ORANGE),
        ("Moradia",     "R$1.058", 25, C_BRAND),
        ("Transporte",  "R$635",   15, C_TEAL),
        ("Lazer",       "R$508",   12, C_PURPLE),
        ("Assinaturas", "R$338",    8, WF_PURP),
        ("Outros",      "R$502",   12, C_GRAY),
    ]
    hbars(slide, mx, col_t+0.38, lw2-0.1, col_h-0.48, cat_segs, lw=1.1, vw=0.75)

    # RIGHT: Insights
    rx2 = MC_X + lw2 + pad + 0.12
    wcard(slide, rx2-0.05, col_t, rw2+0.05, col_h)
    T(slide, "Insights", rx2, col_t+0.08, rw2-0.1, 0.26,
      sz=11, bold=True, color=WF_T1)

    insights = [
        (WF_AMBB, WF_AMB,
         "🍕  Delivery +23% vs abril",
         "R$487 em 31 pedidos este mês"),
        (WF_REDB, WF_RED,
         "💳  3 assinaturas não utilizadas",
         "Netflix Família · Gym · R$198/mês"),
        (WF_GRNB, WF_GRN,
         "✅  Você investiu R$850",
         "Meta de R$1.200 · faltam R$350"),
    ]
    iy = col_t + 0.42
    for bg_col, tc, title, body in insights:
        R(slide, rx2, iy, rw2-0.12, 0.72, bg_col)
        T(slide, title, rx2+0.1, iy+0.06, rw2-0.25, 0.26,
          sz=10, bold=True, color=tc)
        T(slide, body, rx2+0.1, iy+0.34, rw2-0.25, 0.28,
          sz=9, color=WF_T2)
        iy += 0.80

    ann(slide, [
        ("🎯", "Pain Point Central",    "Responde \"onde foi meu dinheiro?\" em <3 segundos de abertura", C_RED),
        ("💡", "Insight contextual",    "Identifica anomalias (delivery +23%) antes que o usuário perceba", C_WARN),
        ("📊", "Visão consolidada",     "Open Finance: todas as contas em um lugar (Pilar 1 + Pilar 2)", C_BRAND),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — T4: ANÁLISE DE GASTOS
# ══════════════════════════════════════════════════════════════════════════════

def slide_gastos(prs):
    slide = blank(prs)
    wf_header(slide, "T4", "Análise de Gastos — Drill-down por Categoria",
              ["P1","P2","P3","P5"], "Dor #1 + #7")
    chrome(slide, "app.financeapp.com.br/gastos")
    app_nav(slide, active="Gastos")
    app_sidebar(slide, active="Gastos")

    pad = 0.12
    mx = MC_X + pad
    mw = MC_W - pad * 2

    # Header + period selector
    R(slide, MC_X, MC_T, MC_W, 0.38, WF_CARD)
    T(slide, "Análise de Gastos — Maio 2026",
      mx, MC_T+0.07, 5.0, 0.26, sz=12, bold=True, color=WF_T1)
    for i, period in enumerate(["Jan","Fev","Mar","Abr","Mai ✓","Jun"]):
        px = mx + 5.5 + i * 0.82
        col = C_BRAND if "✓" in period else WF_BORD
        tc  = C_WHITE if "✓" in period else WF_T2
        R(slide, px, MC_T+0.08, 0.72, 0.24, col)
        T(slide, period.replace(" ✓",""), px, MC_T+0.10, 0.72, 0.20,
          sz=8, color=tc, align=PP_ALIGN.CENTER)

    # 3-column layout
    col_t = MC_T + 0.44
    col_h = MC_H - 0.44
    lw3 = mw * 0.35
    cw3 = mw * 0.38
    rw3 = mw - lw3 - cw3 - 0.2

    # LEFT: Category totals
    wcard(slide, MC_X, col_t, lw3+pad, col_h)
    T(slide, "Por categoria", mx, col_t+0.1, lw3, 0.26,
      sz=10, bold=True, color=WF_T1)
    segs_g = [
        ("Alimentação", "R$1.189", 28, C_ORANGE),
        ("Moradia",     "R$1.058", 25, C_BRAND),
        ("Transporte",  "R$635",   15, C_TEAL),
        ("Lazer",       "R$508",   12, C_PURPLE),
        ("Saúde",       "R$253",    6, C_PINK),
        ("Assinaturas", "R$338",    8, WF_PURP),
        ("Outros",      "R$249",    6, C_GRAY),
    ]
    hbars(slide, mx, col_t+0.42, lw3-0.05, col_h-0.52,
          segs_g, lw=1.0, vw=0.72)

    # CENTER: Month comparison
    cx3 = MC_X + lw3 + pad + 0.1
    wcard(slide, cx3, col_t, cw3, col_h)
    T(slide, "Comparativo mensal", cx3+0.12, col_t+0.1, cw3-0.2, 0.26,
      sz=10, bold=True, color=WF_T1)
    T(slide, "vs. abril ↓ melhorou  ↑ piorou",
      cx3+0.12, col_t+0.38, cw3-0.2, 0.22, sz=8, color=WF_T2)
    comp_rows = [
        ("Alimentação", "R$1.189", "R$1.045", "+13,8%", WF_RED),
        ("Moradia",     "R$1.058", "R$1.058",   "=",    WF_T2),
        ("Transporte",  "R$635",   "R$710",    "-10,6%", WF_GRN),
        ("Lazer",       "R$508",   "R$420",    "+21,0%", WF_RED),
        ("Saúde",       "R$253",   "R$0",      "novo",   WF_AMB),
        ("Assinaturas", "R$338",   "R$338",      "=",    WF_T2),
    ]
    ry = col_t + 0.68
    for cat, maio, abril, diff, dc in comp_rows:
        T(slide, cat,   cx3+0.12, ry, 1.2, 0.30, sz=8, color=WF_T1)
        T(slide, maio,  cx3+1.4,  ry, 0.85, 0.30, sz=8, color=WF_T1)
        T(slide, abril, cx3+2.3,  ry, 0.85, 0.30, sz=8, color=WF_T2)
        T(slide, diff,  cx3+cw3-0.85, ry, 0.72, 0.30,
          sz=8, bold=True, color=dc, align=PP_ALIGN.RIGHT)
        ry += 0.38

    # RIGHT: Top transactions
    rx3 = cx3 + cw3 + 0.1
    wcard(slide, rx3, col_t, rw3+pad-0.15, col_h)
    T(slide, "Maiores gastos", rx3+0.12, col_t+0.1, rw3, 0.26,
      sz=10, bold=True, color=WF_T1)
    trans = [
        ("Aluguel",     "R$1.058", "Moradia",  C_BRAND),
        ("Supermercado","R$428",   "Aliment.", C_ORANGE),
        ("iFood",       "R$387",   "Lazer",    C_PURPLE),
        ("Farmácia",    "R$253",   "Saúde",    C_PINK),
        ("Uber",        "R$212",   "Transp.",  C_TEAL),
        ("Netflix",     "R$59",    "Assin.",   WF_PURP),
    ]
    ty = col_t + 0.44
    for name, val, cat, col in trans:
        R(slide, rx3+0.12, ty, rw3-0.24, 0.40, WF_BORD)
        R(slide, rx3+0.12, ty, 0.06, 0.40, col)
        T(slide, name, rx3+0.26, ty+0.08, 1.4, 0.24, sz=9, color=WF_T1)
        T(slide, cat,  rx3+rw3-1.35, ty+0.08, 0.8, 0.24, sz=8, color=WF_T2)
        T(slide, val,  rx3+rw3-0.52, ty+0.08, 0.45, 0.24,
          sz=9, bold=True, color=WF_RED, align=PP_ALIGN.RIGHT)
        ty += 0.48

    ann(slide, [
        ("🎯", "Dores atacadas", "#1 Não saber onde foi · #7 Impulsos (drill-down mostra os culpados)", C_RED),
        ("💡", "Insight chave",  "Comparativo mensal mostra evolução — usuário vê progresso ou regressão", C_WARN),
        ("👤", "Personas",       "P1 Jovem · P2 Endividado · P3 Iniciante · P5 Profissional Ocupado", C_PURPLE),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — T5: ALERTA INTELIGENTE (CARTÃO)
# ══════════════════════════════════════════════════════════════════════════════

def slide_alerta(prs):
    slide = blank(prs)
    wf_header(slide, "T5", "Alerta Inteligente — Cartão de Crédito",
              ["P2","P1","P5"], "Dor #2 + #3 (timing)")
    chrome(slide, "app.financeapp.com.br/alertas/cartao")
    app_nav(slide, active="Gastos")
    app_sidebar(slide, active="Gastos")

    pad = 0.12
    mx  = MC_X + pad
    mw  = MC_W - pad*2

    # Big alert card
    R(slide, MC_X, MC_T, MC_W, 1.05, WF_REDB)
    R(slide, MC_X, MC_T, 0.08, 1.05, WF_RED)
    T(slide, "⚠  Ação necessária — Fatura Nubank",
      mx, MC_T+0.08, mw*0.65, 0.28, sz=13, bold=True, color=WF_RED)
    T(slide, "Vence em 3 dias (14/05/2026)",
      mx, MC_T+0.38, mw*0.65, 0.24, sz=10, color=WF_AMB)
    T(slide, "R$ 1.847,50",
      mx, MC_T+0.62, 3.5, 0.40, sz=26, bold=True, color=WF_RED)
    wbtn(slide, MC_X+mw-1.8, MC_T+0.38, 1.7, "Quitar agora →", WF_RED, C_WHITE)
    wbtn(slide, MC_X+mw-1.8, MC_T+0.72, 1.7, "Ver parcelamento", WF_BORD, WF_AMB)

    # 3-col: bill breakdown | simulation | tip
    col_t = MC_T + 1.12
    col_h = MC_H - 1.12
    lw5 = mw * 0.36
    cw5 = mw * 0.35
    rw5 = mw - lw5 - cw5 - 0.2

    # LEFT: Bill items
    wcard(slide, MC_X, col_t, lw5+pad, col_h)
    T(slide, "O que está na fatura", mx, col_t+0.1, lw5, 0.26,
      sz=10, bold=True, color=WF_T1)
    items_f = [
        ("iFood — 12 pedidos",    "R$387,00"),
        ("Supermercado Extra",    "R$312,50"),
        ("Posto Ipiranga",        "R$198,00"),
        ("Netflix",               "R$59,90"),
        ("Amazon Prime",          "R$19,90"),
        ("Farmácia São João",     "R$127,40"),
        ("Outros (47 itens)",     "R$742,80"),
    ]
    iy = col_t + 0.44
    for name, val in items_f:
        R(slide, mx, iy, lw5-0.05, 0.34, WF_BORD if iy % 0.68 > 0.34 else WF_CARD)
        T(slide, name, mx+0.1, iy+0.07, lw5-0.75, 0.24, sz=8, color=WF_T1)
        T(slide, val,  mx+lw5-0.65, iy+0.07, 0.58, 0.24,
          sz=8, bold=True, color=WF_RED, align=PP_ALIGN.RIGHT)
        iy += 0.36

    # CENTER: Simulation (the CRITICAL EDUCATION MOMENT)
    cx5 = MC_X + lw5 + pad + 0.1
    R(slide, cx5, col_t, cw5, col_h, RGBColor(0x1A,0x07,0x07))
    R(slide, cx5, col_t, cw5, 0.05, WF_RED)
    T(slide, "⚡ Se você pagar só o mínimo...",
      cx5+0.15, col_t+0.12, cw5-0.25, 0.28,
      sz=11, bold=True, color=WF_RED)
    T(slide, "Mínimo da fatura: R$ 92,38",
      cx5+0.15, col_t+0.50, cw5-0.25, 0.24, sz=10, color=C_WHITE)
    # Horror numbers
    horror = [
        ("Juros que pagará:", "R$ 14.290", WF_RED),
        ("Tempo para quitar:", "31 meses", WF_AMB),
        ("Total final:", "R$ 16.137", WF_RED),
        ("Taxa mensal:", "~7,5% a.m.", WF_AMB),
    ]
    hy = col_t + 0.88
    for lbl, val, col in horror:
        T(slide, lbl, cx5+0.15, hy, cw5*0.6, 0.30, sz=9, color=WF_T2)
        T(slide, val, cx5+cw5*0.6, hy, cw5*0.35, 0.30,
          sz=11, bold=True, color=col, align=PP_ALIGN.RIGHT)
        hy += 0.38

    R(slide, cx5+0.15, hy+0.05, cw5-0.3, 0.03, C_RED)
    T(slide, "Se você quitar agora:",
      cx5+0.15, hy+0.14, cw5-0.25, 0.24, sz=9, color=C_ACCENT)
    T(slide, "Economiza R$ 14.290 em juros",
      cx5+0.15, hy+0.40, cw5-0.25, 0.30,
      sz=11, bold=True, color=C_ACCENT)
    wbtn(slide, cx5+0.15, hy+0.82, cw5-0.3,
         "Quitar R$1.847,50 agora", C_ACCENT, C_DARK, sz=9)

    # RIGHT: Education tip
    rx5 = cx5 + cw5 + 0.1
    wcard(slide, rx5, col_t, rw5+pad-0.1, col_h, RGBColor(0x05,0x2A,0x1A))
    R(slide, rx5, col_t, rw5+pad-0.1, 0.05, C_ACCENT)
    T(slide, "📚 Por que o rotativo é tão perigoso?",
      rx5+0.12, col_t+0.12, rw5-0.15, 0.28,
      sz=10, bold=True, color=C_ACCENT)
    edu_lines = [
        "O cartão de crédito tem duas",
        "faces: parceiro ou inimigo.",
        "",
        "Quando você paga o mínimo,",
        "a dívida restante entra no",
        "rotativo — a taxa mais cara",
        "do Brasil: ~90% ao ano.",
        "",
        "Cada R$100 não pago vira",
        "R$190 em 12 meses.",
        "",
        "O FinanceApp avisa você",
        "ANTES do vencimento para",
        "você nunca cair nessa armadilha.",
    ]
    ML(slide, edu_lines, rx5+0.12, col_t+0.50,
       rw5-0.15, col_h-0.55, sz=9, color=C_LIGHT, sp=1.25)

    ann(slide, [
        ("🎯", "Dor #2 atacada",   "Rotativo (90% a.a.) — mostra o custo real ANTES de o usuário errar", C_RED),
        ("💡", "Pilar 3 — Timing", "Educação financeira chega NO MOMENTO DA AÇÃO, não em tutorial separado", C_ACCENT),
        ("👤", "Personas",         "P2 Endividado Crônico (core) · P1 Jovem · P5 Profissional Ocupado", C_PURPLE),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — T6: METAS E RESERVA DE EMERGÊNCIA
# ══════════════════════════════════════════════════════════════════════════════

def slide_metas(prs):
    slide = blank(prs)
    wf_header(slide, "T6", "Metas e Reserva de Emergência",
              ["P1","P3","P4","P5"], "Dor #3 + #5")
    chrome(slide, "app.financeapp.com.br/metas")
    app_nav(slide, active="Metas")
    app_sidebar(slide, active="Metas")

    pad = 0.12
    mx  = MC_X + pad
    mw  = MC_W - pad*2

    # Header
    R(slide, MC_X, MC_T, MC_W, 0.36, WF_CARD)
    T(slide, "Minhas Metas", mx, MC_T+0.07, 4.0, 0.26,
      sz=13, bold=True, color=WF_T1)
    wbtn(slide, MC_X+mw-1.35, MC_T+0.06, 1.3, "+ Nova meta", C_BRAND, C_WHITE, sz=9)

    goals_t = MC_T + 0.44

    # RESERVA DE EMERGÊNCIA (always first and prominent)
    R(slide, MC_X, goals_t, MC_W, 1.1, RGBColor(0x05,0x2A,0x1A))
    R(slide, MC_X, goals_t, 0.06, 1.1, C_ACCENT)
    T(slide, "🛡  Reserva de Emergência",
      mx, goals_t+0.08, 5.0, 0.28, sz=12, bold=True, color=C_ACCENT)
    T(slide, "Prioridade máxima · 6× seus gastos mensais · R$ 36.000 alvo",
      mx, goals_t+0.38, 7.0, 0.22, sz=9, color=WF_T2)
    T(slide, "R$ 23.400  /  R$ 36.000",
      MC_X+mw-3.8, goals_t+0.15, 3.7, 0.35,
      sz=16, bold=True, color=C_ACCENT, align=PP_ALIGN.RIGHT)
    R(slide, mx, goals_t+0.65, mw-0.1, 0.14, WF_BORD)
    R(slide, mx, goals_t+0.65, (mw-0.1)*0.65, 0.14, C_ACCENT)
    T(slide, "65%", mx+(mw-0.1)*0.67, goals_t+0.66, 0.5, 0.14,
      sz=8, bold=True, color=C_ACCENT)
    T(slide, "Você poupa R$800/mês. Faltam 15 meses para completar a reserva.",
      mx, goals_t+0.84, mw-0.1, 0.22, sz=9, color=WF_T2)

    # Other goals (2 per row)
    goals = [
        ("🏠", "Entrada do imóvel", "R$50.000", "R$18.200", 36, C_BRAND,
         "Poupar R$1.200/mês. Faltam 26 meses."),
        ("✈️", "Viagem Europa", "R$12.000", "R$7.400", 62, C_PURPLE,
         "Faltam R$4.600. Com R$800/mês: 6 meses."),
        ("🚗", "Troca do carro", "R$35.000", "R$4.200", 12, C_TEAL,
         "Meta nova. Definir aporte mensal."),
        ("📈", "Previdência priv.", "R$200.000", "R$31.000", 15, WF_PURP,
         "Aportando R$500/mês. Projeção: 25 anos."),
    ]
    gy = goals_t + 1.20
    for i, (icon, name, alvo, atual, pct, col, tip) in enumerate(goals):
        gx = MC_X + (i % 2) * (mw/2 + 0.1)
        gw = mw/2 - 0.12
        if i % 2 == 0 and i > 0:
            gy += 1.20
        wcard(slide, gx, gy, gw, 1.10)
        T(slide, f"{icon}  {name}", gx+0.12, gy+0.08, gw-0.2, 0.28,
          sz=10, bold=True, color=WF_T1)
        T(slide, f"{atual} / {alvo}",
          gx+0.12, gy+0.38, gw-0.2, 0.24,
          sz=9, color=WF_T2)
        R(slide, gx+0.12, gy+0.65, gw-0.25, 0.10, WF_BORD)
        R(slide, gx+0.12, gy+0.65, (gw-0.25)*pct/100, 0.10, col)
        T(slide, f"{pct}%", gx+0.12+(gw-0.25)*pct/100+0.05, gy+0.64, 0.4, 0.14,
          sz=7, color=col)
        T(slide, tip, gx+0.12, gy+0.80, gw-0.2, 0.25, sz=8, color=WF_T2)

    ann(slide, [
        ("🎯", "Dor #3 atacada",  "Ausência de reserva — 31% dos brasileiros sem nenhuma. Reserva é o primeiro objetivo.", C_RED),
        ("💡", "Dor #5 atacada",  "\"Não sei se estou investindo o suficiente\" → simulador mostra prazo real por meta", C_WARN),
        ("👤", "Personas",        "P1 Jovem · P3 Iniciante · P4 Mulher · P5 Profissional Ocupado", C_PURPLE),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — T7: EDUCAÇÃO CONTEXTUAL
# ══════════════════════════════════════════════════════════════════════════════

def slide_educacao(prs):
    slide = blank(prs)
    wf_header(slide, "T7", "Educação Contextual — No Momento da Ação",
              ["P1","P2","P3","P4"], "Pilar 3 — Timing")
    chrome(slide, "app.financeapp.com.br/gastos/categoria/alimentacao")
    app_nav(slide, active="Educação")
    app_sidebar(slide, active="Gastos")

    pad = 0.12
    mx  = MC_X + pad
    mw  = MC_W - pad*2

    # Context: user is viewing Alimentação category
    R(slide, MC_X, MC_T, MC_W, 0.55, WF_CARD)
    T(slide, "Gastos › Alimentação › Maio 2026",
      mx, MC_T+0.05, 5.0, 0.22, sz=9, color=WF_T2)
    T(slide, "Você gastou R$1.189 em alimentação — 13,8% acima do mês passado",
      mx, MC_T+0.28, 8.0, 0.24, sz=10, bold=True, color=WF_T1)

    # INLINE INSIGHT CARD (education trigger)
    ins_t = MC_T + 0.62
    R(slide, MC_X, ins_t, MC_W, 1.25, RGBColor(0xFF,0xF7,0xE0))
    R(slide, MC_X, ins_t, 0.06, 1.25, C_WARN)
    T(slide, "💡 Você sabia? Brasileiros gastam em média 28% da renda em alimentação.",
      mx, ins_t+0.08, mw*0.75, 0.28, sz=11, bold=True, color=WF_AMB)
    T(slide, "Seu gasto ficou nessa faixa, mas o delivery representa 41% do total. "
             "Preparar refeições em casa pode economizar até R$200/mês.",
      mx, ins_t+0.42, mw*0.75, 0.35, sz=10, color=WF_T1)
    wbtn(slide, MC_X+mw-2.1, ins_t+0.47, 2.0,
         "Ver dica completa →", C_WARN, C_WHITE, sz=9)
    T(slide, "✗ Fechar insight", MC_X+mw-2.2, ins_t+0.98, 2.1, 0.22,
      sz=8, color=WF_T2, align=PP_ALIGN.RIGHT)

    # 2-col: left = transaction list, right = micro-lesson
    col_t = ins_t + 1.32
    col_h = MC_H - (col_t - MC_T) - 0.05
    lw7 = mw * 0.50
    rw7 = mw - lw7 - 0.15

    # LEFT: transactions with categories
    wcard(slide, MC_X, col_t, lw7+pad, col_h)
    T(slide, "Transações de alimentação", mx, col_t+0.08, lw7, 0.26,
      sz=10, bold=True, color=WF_T1)
    trans7 = [
        ("iFood — Burguer King",  "12/05", "R$48,90",  "🍕 Delivery"),
        ("Supermercado Pão de Açúcar", "10/05", "R$213,40", "🛒 Mercado"),
        ("iFood — Sushi Naka",    "09/05", "R$89,00",  "🍕 Delivery"),
        ("iFood — Pizza Hut",     "07/05", "R$67,80",  "🍕 Delivery"),
        ("Hortifruti São Paulo",  "06/05", "R$45,20",  "🥦 Feira"),
        ("McDonald's",            "05/05", "R$32,70",  "🍔 Lanche"),
    ]
    ty7 = col_t + 0.42
    for name, date, val, cat in trans7:
        T(slide, name, mx+0.1, ty7+0.04, lw7*0.52, 0.24, sz=9, color=WF_T1)
        T(slide, date, mx+lw7*0.54, ty7+0.04, 0.6, 0.24, sz=8, color=WF_T2)
        T(slide, cat,  mx+lw7*0.64, ty7+0.04, 0.85, 0.24, sz=8, color=WF_T2)
        T(slide, val,  mx+lw7-0.05, ty7+0.04, 0.6, 0.24,
          sz=9, bold=True, color=WF_RED, align=PP_ALIGN.RIGHT)
        R(slide, mx, ty7+0.30, lw7-0.05, 0.02, WF_BORD)
        ty7 += 0.34

    # RIGHT: Micro-lesson card
    rx7 = MC_X + lw7 + pad + 0.15
    R(slide, rx7, col_t, rw7+pad-0.2, col_h, RGBColor(0x05,0x2A,0x1A))
    R(slide, rx7, col_t, rw7+pad-0.2, 0.05, C_ACCENT)
    T(slide, "📚 Lição: Delivery vs. Cozinha",
      rx7+0.15, col_t+0.12, rw7-0.15, 0.28,
      sz=11, bold=True, color=C_ACCENT)

    lesson = [
        "Seu delivery custou R$487",
        "este mês (31 pedidos).",
        "",
        "Média por pedido: R$15,70.",
        "",
        "Se você cozinhar 3× por",
        "semana ao invés de pedir:",
        "",
        "  ✅ Economiza ~R$180/mês",
        "  ✅ R$2.160 por ano",
        "  ✅ O suficiente para a",
        "      passagem da viagem",
        "      que você quer fazer.",
        "",
        "Não precisa ser perfeito.",
        "Comece com 1× por semana.",
    ]
    ML(slide, lesson, rx7+0.15, col_t+0.48,
       rw7-0.18, col_h-0.55, sz=9, color=C_LIGHT, sp=1.22)

    ann(slide, [
        ("🎯", "Pilar 3 — Timing",   "Educação ocorre DENTRO da ação, não em seção separada de 'Aprender'", C_ACCENT),
        ("💡", "Personalizado",       "A lição conecta o gasto com a META do usuário (viagem). Contexto = retenção.", C_WARN),
        ("👤", "Personas",            "P1 Jovem · P2 Endividado · P3 Iniciante · P4 Mulher Invisível", C_PURPLE),
    ])
    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — PERSONAS × TELAS
# ══════════════════════════════════════════════════════════════════════════════

def slide_matrix(prs):
    slide = blank(prs)
    pres_header(slide, "Personas × Telas — Quem é Servido por Quê",
                "Cada tela atende múltiplos perfis. Nenhuma persona fica sem cobertura.")

    personas = ["P1\nJovem", "P2\nEndivid.", "P3\nInic.", "P4\nMulher",
                "P5\nOcupado", "P6\nAvanç.", "P7\nAposent.", "P8\nMEI"]
    screens  = ["T1\nLanding", "T2\nOnboard.", "T3\nDashboard",
                "T4\nGastos", "T5\nAlerta", "T6\nMetas", "T7\nEducação"]

    # Matrix: which personas are served by which screen
    # 1 = primary, 0.5 = secondary, 0 = not addressed in MVP
    matrix = [
        #P1  P2    P3    P4    P5    P6    P7    P8
        [1,   1,    1,    1,    1,    0.5,  0.5,  0.5],  # T1 Landing
        [1,   1,    1,    1,    1,    0,    0,    0.5],  # T2 Onboarding
        [1,   1,    1,    1,    1,    0.5,  0,    1  ],  # T3 Dashboard
        [1,   1,    1,    0.5,  1,    0,    0,    0.5],  # T4 Gastos
        [1,   1,    0.5,  0,    1,    0,    0,    0  ],  # T5 Alerta
        [1,   0.5,  1,    1,    1,    0,    0.5,  0  ],  # T6 Metas
        [1,   1,    1,    1,    0,    0,    0,    0  ],  # T7 Educação
    ]

    # Draw grid
    cell_w = 1.38
    cell_h = 0.54
    hdr_w  = 1.55
    col_start = 0.25 + hdr_w

    persona_colors = [C_WARN, C_RED, C_ACCENT, C_PINK, C_BRAND, C_TEAL, C_PURPLE, C_ORANGE]
    screen_colors  = [C_ORANGE, C_WARN, C_RED, C_BRAND, C_RED, C_ACCENT, C_PURPLE]

    # Column headers (personas)
    for j, (p, pc) in enumerate(zip(personas, persona_colors)):
        x = col_start + j * cell_w
        R(slide, x, 1.15, cell_w-0.05, 0.52, pc)
        T(slide, p, x, 1.18, cell_w-0.05, 0.48,
          sz=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Row headers (screens) + cells
    for i, (scr, sc, row) in enumerate(zip(screens, screen_colors, matrix)):
        y = 1.73 + i * cell_h
        # Row header
        R(slide, 0.25, y, hdr_w-0.05, cell_h-0.04, sc)
        T(slide, scr, 0.25, y+0.03, hdr_w-0.05, cell_h-0.1,
          sz=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        # Cells
        for j, val in enumerate(row):
            x = col_start + j * cell_w
            if val == 1:
                bg_c = persona_colors[j]
                sym, tc = "●  Core", C_DARK
            elif val == 0.5:
                bg_c = RGBColor(0x2D, 0x3A, 0x5E)
                sym, tc = "◐  Sec.", C_GRAY
            else:
                bg_c = RGBColor(0x16, 0x23, 0x3A)
                sym, tc = "—", RGBColor(0x2D, 0x3A, 0x4E)
            R(slide, x, y, cell_w-0.05, cell_h-0.04, bg_c)
            T(slide, sym, x, y+0.14, cell_w-0.05, 0.24,
              sz=8, bold=(val==1), color=tc, align=PP_ALIGN.CENTER)

    # Legend
    R(slide, 0.25, 5.55, 12.83, 0.40, RGBColor(0x16,0x23,0x3A))
    T(slide, "●  Core = tela principal para essa persona",
      0.4, 5.62, 4.0, 0.26, sz=10, color=C_ACCENT)
    T(slide, "◐  Sec. = tela também serve (secundariamente)",
      4.8, 5.62, 4.5, 0.26, sz=10, color=C_GRAY)
    T(slide, "—  Fora do escopo MVP para esse perfil",
      9.5, 5.62, 3.5, 0.26, sz=10, color=RGBColor(0x2D,0x3A,0x4E))

    # Note
    T(slide, "P6 (Investidor Avançado) e P7 (Aposentado) têm cobertura parcial no MVP "
             "— são alvo de versões futuras com funcionalidades específicas.",
      0.25, 6.05, 12.83, 0.30, sz=9, color=C_GRAY)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — ESCOPO MVP
# ══════════════════════════════════════════════════════════════════════════════

def slide_scope(prs):
    slide = blank(prs)
    pres_header(slide, "Escopo MVP — O Que Entra em v1.0",
                "Foco máximo na Dor #1. Tudo o mais é v2+")

    in_items = [
        ("T1", "Landing page com problem statement e CTA",         C_ORANGE),
        ("T2", "Onboarding + diagnóstico via Open Finance",         C_WARN),
        ("T3", "Dashboard «Onde foi meu dinheiro?»",               C_RED),
        ("T4", "Análise de gastos por categoria + comparativo",    C_BRAND),
        ("T5", "Alertas de fatura de cartão + simulador rotativo", C_RED),
        ("T6", "Metas e reserva de emergência",                    C_ACCENT),
        ("T7", "Educação contextual inline (insights automáticos)", C_PURPLE),
        ("",   "Open Finance: Pluggy (Nubank, Itaú, Bradesco...)", C_TEAL),
        ("",   "Autenticação: email + Google OAuth",               C_BRAND),
        ("",   "Freemium: tier Free + Essentials (R$19,90/mês)",  C_ACCENT),
    ]

    out_items = [
        "Investimentos e carteira (v2)",
        "Análise por IA / categorização automática avançada (v2)",
        "Simulador de IR sobre investimentos (v2)",
        "App mobile nativo iOS/Android (v2)",
        "Canal B2B / benefício corporativo (v2)",
        "Módulo de crédito e antecipação de recebíveis (v3)",
        "Previdência privada e simulador de aposentadoria (v3)",
        "Integração multi-corretora para P6 (v3)",
        "Funcionalidades específicas para MEI/PJ (v3)",
        "Suporte a P7 (Aposentado) com interface acessível (v3)",
    ]

    # IN column
    R(slide, 0.25, 1.15, 7.3, 0.38, C_ACCENT)
    T(slide, "✅  Entra no MVP v1.0", 0.4, 1.22, 7.0, 0.28,
      sz=13, bold=True, color=C_DARK)
    for i, (tag, desc, col) in enumerate(in_items):
        y = 1.60 + i * 0.52
        R(slide, 0.25, y, 7.3, 0.48, RGBColor(0x16,0x23,0x3A) if i%2==0 else C_DARK)
        if tag:
            R(slide, 0.25, y, 0.52, 0.48, col)
            T(slide, tag, 0.25, y+0.10, 0.52, 0.28,
              sz=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
            T(slide, desc, 0.83, y+0.11, 6.6, 0.28, sz=10, color=C_WHITE)
        else:
            R(slide, 0.25, y, 0.52, 0.48, RGBColor(0x16,0x23,0x3A))
            T(slide, "—", 0.25, y+0.10, 0.52, 0.28,
              sz=9, color=C_GRAY, align=PP_ALIGN.CENTER)
            T(slide, desc, 0.83, y+0.11, 6.6, 0.28, sz=10, color=C_GRAY)

    # OUT column
    R(slide, 7.85, 1.15, 5.23, 0.38, RGBColor(0x2D,0x3A,0x4E))
    T(slide, "🔲  Fora do MVP — versões futuras", 8.0, 1.22, 5.0, 0.28,
      sz=13, bold=True, color=C_GRAY)
    for i, desc in enumerate(out_items):
        y = 1.60 + i * 0.52
        R(slide, 7.85, y, 5.23, 0.48, RGBColor(0x14,0x1D,0x2E) if i%2==0 else C_DARK)
        T(slide, "→ " + desc, 8.0, y+0.11, 5.0, 0.28, sz=10, color=C_GRAY)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — JORNADA DE ATIVAÇÃO D1→D7
# ══════════════════════════════════════════════════════════════════════════════

def slide_journey(prs):
    slide = blank(prs)
    pres_header(slide, "Jornada de Ativação — Do Sign-up ao Hábito",
                "Os primeiros 7 dias determinam se o usuário fica. Meta: 40% completando o ciclo.")

    days = [
        ("D1", "Sign-up\n+ Diagnóstico",
         "Conecta banco → vê onde\nfoi o dinheiro do mês",
         "Aha moment\nem <5 min", C_ORANGE),
        ("D2", "1ª Meta\ndefinida",
         "App sugere: \"Que tal criar\numa reserva de emergência?\"",
         "Orientação\nativa", C_WARN),
        ("D3", "1º Alerta\nrecebido",
         "\"Você gastou R$80 em delivery\nnos últimos 3 dias.\"",
         "Insight\ncontextual", C_RED),
        ("D5", "1ª Ação\ntomada",
         "Usuário cancela assinatura\nnão utilizada (R$59/mês)",
         "Primeira\nvitória", C_ACCENT),
        ("D7", "Check-in\nsemanal",
         "Dashboard semanal:\n\"Você economizou R$280!\"",
         "Loop\nformado", C_BRAND),
    ]

    arrow_x = []
    for i, (day, title, desc, badge, col) in enumerate(days):
        x = 0.4 + i * 2.55
        arrow_x.append(x + 2.0)
        # Card
        R(slide, x, 1.4, 2.2, 4.8, col)
        T(slide, day, x, 1.5, 2.2, 0.55,
          sz=26, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        R(slide, x, 2.05, 2.2, 0.03, C_DARK)
        T(slide, title, x+0.1, 2.12, 2.0, 0.55,
          sz=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        T(slide, desc, x+0.1, 2.72, 2.0, 0.80,
          sz=10, color=C_DARK, align=PP_ALIGN.CENTER)
        R(slide, x+0.25, 3.7, 1.7, 0.50, RGBColor(0,0,0))
        T(slide, badge, x+0.25, 3.77, 1.7, 0.40,
          sz=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Arrows between cards
    for ax in arrow_x[:-1]:
        T(slide, "→", ax, 3.55, 0.45, 0.40,
          sz=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # KPIs strip
    R(slide, 0.25, 6.35, 12.83, 0.75, RGBColor(0x16,0x23,0x3A))
    kpis = [
        ("40%", "completam ciclo D1→D7"),
        ("60%", "retornam na semana 2"),
        ("NPS >50", "antes de escalar marketing"),
        ("<5 min", "para o primeiro aha moment"),
    ]
    kw = 12.83 / len(kpis)
    for i, (val, lbl) in enumerate(kpis):
        kx = 0.25 + i * kw
        T(slide, val, kx, 6.40, kw, 0.38,
          sz=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        T(slide, lbl, kx, 6.80, kw, 0.25,
          sz=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — MÉTRICAS DE SUCESSO
# ══════════════════════════════════════════════════════════════════════════════

def slide_metrics(prs):
    slide = blank(prs)
    pres_header(slide, "Métricas de Sucesso — Como Saberemos que Funcionou",
                "Cada métrica valida um pilar do Problem Statement")

    groups = [
        ("Visibilidade", C_RED, [
            ("Taxa de ativação D1",   "≥ 70%",   "conectam banco e veem diagnóstico no 1º dia"),
            ("Tempo até aha moment",  "< 5 min",  "da criação de conta até ver \"onde foi meu dinheiro\""),
            ("Sessões/semana",        "≥ 3×",     "usuário retorna espontaneamente para ver o dashboard"),
        ]),
        ("Orientação", C_WARN, [
            ("Ações tomadas D1-D7",   "≥ 1 ação",  "cancelar assinatura, criar meta, ou quitar fatura"),
            ("Meta criada D1-D7",     "≥ 60%",     "dos usuários ativos criam ao menos 1 meta"),
            ("Alertas clicados",      "≥ 40%",     "taxa de clique em alertas de cartão/vencimento"),
        ]),
        ("Timing (Educação)", C_ACCENT, [
            ("Insights lidos",        "≥ 35%",     "dos insights inline gerados são lidos até o fim"),
            ("NPS 30 dias",           "> 50",       "net promoter score no final do primeiro mês"),
            ("Retenção D30",          "≥ 45%",     "usuários ainda ativos 30 dias após o cadastro"),
        ]),
        ("Negócio", C_BRAND, [
            ("Conversão free→pago",   "≥ 8%",      "após 30 dias de uso no plano gratuito"),
            ("CAC orgânico",          "< R$50",     "custo de aquisição via conteúdo e boca a boca"),
            ("LTV/CAC",               "> 8×",       "indicador de saúde do modelo SaaS"),
        ]),
    ]

    for i, (group, col, metrics) in enumerate(groups):
        col_i = i % 2
        row_i = i // 2
        x = 0.25 + col_i * 6.55
        y = 1.15 + row_i * 3.05
        gw = 6.30
        gh = 2.90

        R(slide, x, y, gw, gh, RGBColor(0x16,0x23,0x3A))
        R(slide, x, y, gw, 0.40, col)
        T(slide, f"Pilar: {group}", x+0.15, y+0.08, gw-0.3, 0.28,
          sz=13, bold=True, color=C_DARK)

        for j, (metric, target, desc) in enumerate(metrics):
            my = y + 0.48 + j * 0.77
            R(slide, x+0.15, my, gw-0.3, 0.70, C_DARK)
            R(slide, x+0.15, my, 0.05, 0.70, col)
            T(slide, metric, x+0.28, my+0.06, 3.5, 0.26, sz=10, bold=True, color=C_WHITE)
            R(slide, x+gw-1.55, my+0.06, 1.35, 0.30, col)
            T(slide, target, x+gw-1.55, my+0.08, 1.35, 0.26,
              sz=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
            T(slide, desc, x+0.28, my+0.38, gw-0.6, 0.26, sz=9, color=C_GRAY)

    footer_tag(slide)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build(prs):
    slide_title(prs)
    slide_anchor(prs)
    slide_arch(prs)
    slide_landing(prs)
    slide_onboarding(prs)
    slide_dashboard(prs)
    slide_gastos(prs)
    slide_alerta(prs)
    slide_metas(prs)
    slide_educacao(prs)
    slide_matrix(prs)
    slide_scope(prs)
    slide_journey(prs)
    slide_metrics(prs)


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = new_prs()
    build(prs)
    out = os.path.join(OUT_DIR, f"prototype-{VERSION}.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
