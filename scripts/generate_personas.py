"""
Generate FinanceApp Personas PowerPoint presentation.
Usage: python scripts/generate_personas.py
Output: docs/presentations/personas-v1.0.0.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
C_BRAND  = RGBColor(0x00, 0x84, 0xFF)
C_ACCENT = RGBColor(0x00, 0xD4, 0x8A)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x6B, 0x72, 0x80)
C_WARN   = RGBColor(0xF5, 0x9E, 0x0B)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_PINK   = RGBColor(0xEC, 0x48, 0x99)
C_TEAL   = RGBColor(0x14, 0xB8, 0xA6)
C_ORANGE = RGBColor(0xF9, 0x73, 0x16)

VERSION = "v1.0.0"
DATE    = "2026-05-11"
OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations")


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def add_lines(slide, lines, left, top, width, height,
              font_size=12, color=C_WHITE, bold=False, spacing=1.25):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(font_size * (spacing - 1))
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb

def slide_header(slide, title, subtitle="", accent=C_BRAND):
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, accent)
    add_rect(slide, 0, 7.42, 13.33, 0.08, accent)
    add_text(slide, title, 0.5, 0.18, 12.3, 0.72,
             font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.88, 12.3, 0.38,
                 font_size=13, color=C_GRAY)

def footer(slide, text):
    add_text(slide, text, 0, 7.12, 13.33, 0.32,
             font_size=10, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 1: Title ─────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0,    13.33, 0.12, C_PURPLE)
    add_rect(slide, 0, 7.38, 13.33, 0.12, C_PURPLE)
    add_rect(slide, 0, 2.85, 13.33, 0.04, C_PURPLE)
    add_text(slide, "FinanceApp", 1, 0.45, 11, 0.9,
             font_size=18, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, "Estudo de Personas", 1, 1.4, 11, 1.3,
             font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "8 perfis · Dores reais · Pain point central · Problem Statement",
             1, 3.05, 11, 0.75,
             font_size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Baseado em: ANBIMA · B3 · Serasa · SPC Brasil · CNC · Febraban · BCB · IBGE · Creditas · Nubank/Ipsos",
             1, 4.0, 11, 0.5,
             font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, f"{VERSION}  ·  {DATE}", 1, 6.85, 11, 0.4,
             font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 2: Números do mercado ────────────────────────────────────────────────
def slide_market_numbers(prs):
    slide = blank_slide(prs)
    slide_header(slide, "O Contexto: Números que Definem o Problema",
                 "Dados públicos — ANBIMA, Serasa, CNC, Febraban, BCB (2025/2026)",
                 accent=C_PURPLE)

    stats = [
        ("82,8M",  "brasileiros com\nCPF negativado",      C_RED,    "Serasa mar/2026\n(recorde histórico)"),
        ("80,9%",  "das famílias\nendividadas",             C_ORANGE, "CNC abr/2026\n(recorde histórico)"),
        ("107,7M", "brasileiros que\nnão investem",         C_WARN,   "ANBIMA 2025\n(64% da população)"),
        ("48%",    "não controlam\nseus gastos",            C_PURPLE, "SPC Brasil 2025"),
        ("31%",    "sem nenhuma\nreserva de emergência",    C_PINK,   "ANBIMA 2025"),
        ("90%",    "juros anuais do\nrotativo do cartão",   C_RED,    "CNC 2026"),
    ]

    for i, (val, label, color, source) in enumerate(stats):
        col, row = i % 3, i // 3
        x = 0.35 + col * 4.3
        y = 1.4  + row * 2.85
        add_rect(slide, x, y, 3.9, 2.6, color)
        add_text(slide, val,    x, y + 0.1,  3.9, 1.0,
                 font_size=38, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, label,  x, y + 1.1,  3.9, 0.75,
                 font_size=13, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, source, x, y + 1.9,  3.9, 0.55,
                 font_size=10, align=PP_ALIGN.CENTER, color=C_DARK)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ── Slide 3: Mapa das 8 personas ───────────────────────────────────────────────
def slide_persona_map(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Os 8 Perfis Identificados",
                 "Segmentados por faixa etária, renda e comportamento financeiro",
                 accent=C_PURPLE)

    personas = [
        ("P1", "Jovem Desorganizado",      "Gen Z · 18-24a · até R$2.500",   C_WARN),
        ("P2", "Endividado Crônico",        "25-45a · 1-5 SM · cartão no limit", C_RED),
        ("P3", "Investidor Iniciante",      "25-40a · classe C · começando",   C_ACCENT),
        ("P4", "Mulher Invisível",          "25-60a · subatendida · 31% invest.", C_PINK),
        ("P5", "Profissional Ocupado",      "30-50a · 5+ SM · sem tempo",      C_BRAND),
        ("P6", "Investidor Avançado",       "28-50a · renda alta · multi-ativo", C_TEAL),
        ("P7", "Aposentado/Pré-Aposent.",   "50-70a · renda fixa · segurança",  C_PURPLE),
        ("P8", "MEI/Empreendedor",          "25-55a · renda variável · PF=PJ",  C_ORANGE),
    ]

    for i, (code, name, desc, color) in enumerate(personas):
        col, row = i % 4, i // 4
        x = 0.3  + col * 3.25
        y = 1.35 + row * 2.85
        add_rect(slide, x, y, 3.0, 2.6, color)
        add_text(slide, code, x, y + 0.08, 3.0, 0.55,
                 font_size=24, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_rect(slide, x, y + 0.62, 3.0, 0.04, C_DARK)
        add_text(slide, name, x + 0.1, y + 0.72, 2.8, 0.65,
                 font_size=14, bold=True, color=C_DARK)
        add_text(slide, desc, x + 0.1, y + 1.4,  2.8, 1.0,
                 font_size=11, color=C_DARK)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ── Helper: persona detail slide ───────────────────────────────────────────────
def slide_persona(prs, code, name, age, income, profile_desc,
                  color, pains, wants, quote, data_point, source):
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0,    13.33, 0.08, color)
    add_rect(slide, 0, 7.42, 13.33, 0.08, color)

    # Header badge
    add_rect(slide, 0.35, 0.18, 0.85, 0.85, color)
    add_text(slide, code, 0.35, 0.22, 0.85, 0.75,
             font_size=22, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)

    add_text(slide, name, 1.3, 0.18, 9.0, 0.55,
             font_size=26, bold=True, color=C_WHITE)
    add_text(slide, f"{age}  ·  {income}", 1.3, 0.72, 9.0, 0.38,
             font_size=13, color=C_GRAY)

    # Profile strip
    add_rect(slide, 0.35, 1.2, 12.63, 0.55, RGBColor(0x1A, 0x22, 0x38))
    add_text(slide, profile_desc, 0.5, 1.25, 12.3, 0.45,
             font_size=12, color=C_LIGHT)

    # Left: pains
    add_rect(slide, 0.35, 1.9, 5.9, 4.7, RGBColor(0x1A, 0x07, 0x07))
    add_text(slide, "Principais Dores", 0.5, 1.98, 5.6, 0.45,
             font_size=15, bold=True, color=C_RED)
    add_lines(slide, ["• " + p for p in pains],
              0.5, 2.5, 5.6, 3.9, font_size=12, color=C_LIGHT, spacing=1.45)

    # Right: wants
    add_rect(slide, 6.5, 1.9, 6.48, 4.7, RGBColor(0x05, 0x2A, 0x1A))
    add_text(slide, "O Que Quer do App", 6.65, 1.98, 6.1, 0.45,
             font_size=15, bold=True, color=C_ACCENT)
    add_lines(slide, ["✔ " + w for w in wants],
              6.65, 2.5, 6.1, 3.9, font_size=12, color=C_LIGHT, spacing=1.45)

    # Quote strip
    add_rect(slide, 0.35, 6.68, 8.5, 0.65, color)
    add_text(slide, f'"{quote}"', 0.5, 6.71, 8.2, 0.58,
             font_size=11, bold=False, color=C_DARK)

    # Data badge
    add_rect(slide, 9.1, 6.68, 4.23, 0.65, RGBColor(0x16, 0x23, 0x3A))
    add_text(slide, f"{data_point}  —  {source}", 9.2, 6.72, 4.0, 0.55,
             font_size=10, color=C_GRAY)


# ── Slide 12: Ranking de dores ─────────────────────────────────────────────────
def slide_pains_ranking(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Ranking das Dores — Transversal a Todos os Perfis",
                 "Evidência quantitativa de cada dor  ·  Quanto mais vermelho, mais crítico",
                 accent=C_PURPLE)

    rows = [
        ("1", "Não sabe para onde o dinheiro vai",        "48% não controlam gastos",              C_RED),
        ("2", "Dívida no cartão de crédito",              "85,1% dos endividados têm cartão",      C_RED),
        ("3", "Sem reserva de emergência",                 "31% sem nenhuma reserva",               C_RED),
        ("4", "Renda que nunca sobra",                    "51% dos jovens; 95% querem renda extra", C_ORANGE),
        ("5", "Não sabe como nem onde investir",          "52% sem planejamento; 69% mulheres",    C_ORANGE),
        ("6", "Imprevistos que destroem o planejamento",  "32% citam como causa principal",        C_ORANGE),
        ("7", "Compras por impulso",                      "44-56% cedem frequentemente",           C_WARN),
        ("8", "Estresse emocional com dinheiro",          "77% endividados com impacto na saúde",  C_WARN),
        ("9", "Não se prepara para aposentadoria",        "75% da Gen Z ignora o tema",            C_PURPLE),
        ("10","Fraudes e golpes digitais",                 "39% já sofreram tentativas",            C_PURPLE),
    ]

    rh = 0.49
    hy = 1.35
    add_rect(slide, 0.3, hy, 12.7, rh, C_PURPLE)
    for h, x, w in zip(["#", "Dor", "Evidência (fonte pública)"],
                        [0.4, 1.05, 8.4], [0.55, 7.2, 4.5]):
        add_text(slide, h, x, hy + 0.07, w, rh - 0.14,
                 font_size=13, bold=True, color=C_DARK)

    for i, (num, pain, evidence, color) in enumerate(rows):
        y = hy + rh * (i + 1)
        bg = RGBColor(0x16, 0x23, 0x3A) if i % 2 == 0 else C_DARK
        add_rect(slide, 0.3, y, 12.7, rh, bg)
        add_rect(slide, 0.3, y, 0.65, rh, color)
        add_text(slide, num, 0.3, y + 0.08, 0.65, rh - 0.16,
                 font_size=14, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_text(slide, pain,     1.05, y + 0.09, 7.2,  rh - 0.18, font_size=12, color=C_WHITE)
        add_text(slide, evidence, 8.4,  y + 0.09, 4.5,  rh - 0.18, font_size=11, color=C_GRAY)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ── Slide 13: O que todos querem ───────────────────────────────────────────────
def slide_wants(prs):
    slide = blank_slide(prs)
    slide_header(slide, "O Que os Brasileiros Querem de um App Financeiro",
                 "Síntese de: Ipsos/Nubank · Creditas · Febraban · ANBIMA · SPC Brasil",
                 accent=C_PURPLE)

    items = [
        (C_ACCENT,  "Controle automático e categorizado de gastos",
                    "80% usam apps para controle; querem menos trabalho manual (Febraban)"),
        (C_BRAND,   "Recomendações personalizadas de investimento",
                    "66% esperam isso do app/banco (Ipsos/Nubank)"),
        (C_WARN,    "Reserva de emergência como objetivo padrão",
                    "50% citam como meta principal (Ipsos/Nubank)"),
        (C_PURPLE,  "Educação financeira integrada à experiência",
                    "84% buscam aprender; querem aprender no momento da ação (Nubank)"),
        (C_TEAL,    "Visão consolidada de todo o patrimônio",
                    "Open Finance: 62M consentimentos ativos; crescimento de 44% (BCB)"),
        (C_ORANGE,  "Simulador de metas financeiras",
                    "'Quanto poupar para comprar um imóvel em 3 anos?' (Creditas)"),
        (C_PINK,    "Alerta de vencimentos e fatura do cartão",
                    "Evitar juros do rotativo (90% a.a.) e negativação (SPC Brasil)"),
        (C_RED,     "Plano claro de saída das dívidas",
                    "29,7% com contas em atraso; querem ver a luz no fim do túnel (CNC)"),
    ]

    for i, (color, want, evidence) in enumerate(items):
        col, row = i % 2, i // 2
        x = 0.3  + col * 6.5
        y = 1.4  + row * 1.45
        add_rect(slide, x, y, 6.2, 1.3, RGBColor(0x16, 0x23, 0x3A))
        add_rect(slide, x, y, 0.18, 1.3, color)
        add_text(slide, want,     x + 0.3, y + 0.08, 5.8, 0.5,
                 font_size=13, bold=True, color=C_WHITE)
        add_text(slide, evidence, x + 0.3, y + 0.62, 5.8, 0.55,
                 font_size=11, color=C_GRAY)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ── Slide 14: Pain point central ──────────────────────────────────────────────
def slide_central_pain(prs):
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0,    13.33, 0.08, C_RED)
    add_rect(slide, 0, 7.42, 13.33, 0.08, C_RED)

    add_text(slide, "O Pain Point Central", 0.5, 0.2, 12.3, 0.55,
             font_size=20, bold=True, color=C_RED)
    add_text(slide, "O problema que aparece em todos os 8 perfis, independente de renda ou geração",
             0.5, 0.75, 12.3, 0.4, font_size=13, color=C_GRAY)

    # Big statement
    add_rect(slide, 0.5, 1.3, 12.33, 2.5, RGBColor(0x1A, 0x07, 0x07))
    add_text(slide,
             "O brasileiro não sabe para onde\nseu dinheiro vai.",
             0.7, 1.45, 11.9, 1.6,
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "48% não controlam seus gastos · 85% dos endividados têm cartão como principal dívida · "
             "31% sem reserva de emergência · 107M não investem",
             0.7, 3.0, 11.9, 0.55,
             font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Why it's the root
    add_rect(slide, 0.5, 3.75, 12.33, 0.45, C_RED)
    add_text(slide, "Por que essa é a raiz de tudo?", 0.65, 3.82, 12.0, 0.35,
             font_size=14, bold=True, color=C_DARK)

    causes = [
        ("Sem visibilidade dos gastos →", "não consegue poupar, nem montar reserva de emergência"),
        ("Sem controle do cartão →",       "cai no rotativo (90% a.a.), se endivida mais a cada mês"),
        ("Sem sobra de dinheiro →",        "não investe; não se prepara para aposentadoria"),
        ("Sem planejamento →",             "imprevistos viram catástrofes financeiras"),
        ("Sem educação no momento da ação →", "repete os mesmos erros mês após mês"),
    ]
    for i, (cause, effect) in enumerate(causes):
        y = 4.3 + i * 0.56
        add_rect(slide, 0.5, y, 4.5, 0.5, RGBColor(0x2A, 0x0D, 0x0D))
        add_rect(slide, 5.2, y, 7.63, 0.5, RGBColor(0x16, 0x23, 0x3A))
        add_text(slide, cause,  0.65, y + 0.07, 4.2, 0.38, font_size=12, bold=True, color=C_RED)
        add_text(slide, effect, 5.35, y + 0.07, 7.3, 0.38, font_size=12, color=C_LIGHT)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ── Slide 15: Problem Statement ────────────────────────────────────────────────
def slide_problem_statement(prs):
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK)
    add_rect(slide, 0, 0,    13.33, 0.08, C_ACCENT)
    add_rect(slide, 0, 7.42, 13.33, 0.08, C_ACCENT)

    add_text(slide, "Problem Statement", 0.5, 0.2, 12.3, 0.55,
             font_size=20, bold=True, color=C_ACCENT)
    add_text(slide, "O problema que o FinanceApp precisa resolver",
             0.5, 0.75, 12.3, 0.38, font_size=13, color=C_GRAY)

    # Main statement box
    add_rect(slide, 0.5, 1.25, 12.33, 3.0, RGBColor(0x05, 0x2A, 0x1A))
    add_text(slide,
             "O brasileiro ganha, gasta, e no fim do mês\n"
             "não sabe o que aconteceu com o dinheiro.",
             0.8, 1.4, 11.7, 1.5,
             font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Sem visibilidade real dos seus gastos, ele não consegue sair do vermelho,\n"
             "montar uma reserva, nem investir para o futuro — mesmo quando quer.",
             0.8, 2.9, 11.7, 0.9,
             font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Three pillars
    pillars = [
        (C_RED,    "Falta de\nVisibilidade",
                   "Não enxerga para onde\no dinheiro vai em tempo real",
                   "48% sem controle\n(SPC Brasil)"),
        (C_WARN,   "Falta de\nOrientação",
                   "Sabe que tem problema mas\nnão sabe o próximo passo",
                   "52% sem planejamento\n(Creditas/Leve)"),
        (C_ACCENT, "Falta de\nMomento Certo",
                   "A educação financeira chega\ntarde, fora do contexto da ação",
                   "84% querem aprender\nmas apenas 36% investem"),
    ]
    for i, (color, title, desc, data) in enumerate(pillars):
        x = 0.5 + i * 4.15
        add_rect(slide, x, 4.45, 3.85, 2.7, color)
        add_text(slide, title, x, 4.55, 3.85, 0.75,
                 font_size=20, bold=True, align=PP_ALIGN.CENTER, color=C_DARK)
        add_rect(slide, x, 5.32, 3.85, 0.04, C_DARK)
        add_text(slide, desc, x + 0.1, 5.42, 3.65, 0.85,
                 font_size=12, color=C_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, data, x + 0.1, 6.3,  3.65, 0.65,
                 font_size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    footer(slide, f"FinanceApp · Estudo de Personas {VERSION} · {DATE}")


# ══════════════════════════════════════════════════════════════════════════════

def build(prs):
    slide_title(prs)
    slide_market_numbers(prs)
    slide_persona_map(prs)

    # ── P1: Jovem Desorganizado ───────────────────────────────────────────────
    slide_persona(prs,
        code="P1", name="O Jovem Desorganizado", color=C_WARN,
        age="Gen Z · 18–24 anos",
        income="Até R$2.500/mês · CLT inicial ou informal",
        profile_desc="Entrou no mercado de trabalho recentemente. Tem smartphone, usa Nubank, mas nunca aprendeu a gerir dinheiro. Gasta por impulso, sem noção do que sobra.",
        pains=[
            "Não sabe como começar a controlar os gastos",
            "Nunca sobra dinheiro no fim do mês (51%)",
            "Falta de disciplina e compras por impulso (56%)",
            "37% já tiveram o nome negativado",
            "75% não pensam em aposentadoria",
            "Vergonha do estado atual das finanças",
        ],
        wants=[
            "Interface visual e intuitiva (sem planilhas)",
            "Feedback imediato: quanto gastei hoje?",
            "Metas pequenas e alcançáveis (gamificação)",
            "Conteúdo de educação em formato de reels/stories",
            "Sem julgamento — linguagem acolhedora",
        ],
        quote="No fim do mês fico me perguntando pra onde foi o dinheiro. Sai e eu nem vi.",
        data_point="47% da Gen Z não controlam finanças",
        source="CNDL/SPC Brasil 2025"
    )

    # ── P2: Endividado Crônico ────────────────────────────────────────────────
    slide_persona(prs,
        code="P2", name="O Endividado Crônico", color=C_RED,
        age="25–45 anos",
        income="1–5 salários mínimos · R$1.412 a R$7.060",
        profile_desc="Vive no limite do cartão de crédito. Paga apenas o mínimo e vê os juros (90% a.a.) consumirem a renda. Sente que nunca vai sair do vermelho.",
        pains=[
            "Espiral de dívidas: paga mínimo, juros crescem",
            "Rotativo do cartão: 90% a.a. (CNC 2026)",
            "Comprometimento de renda com dívidas: ~30%",
            "Estresse emocional severo: 77% com impacto na saúde",
            "Medo constante de negativação",
            "Vergonha de falar sobre a situação financeira",
        ],
        wants=[
            "Simulador: quando vou sair do vermelho?",
            "Alerta de fatura ANTES do vencimento",
            "Plano de quitação passo a passo",
            "Visão clara dos juros que está pagando",
            "Renegociação de dívidas dentro do app",
        ],
        quote="Pago o mínimo do cartão todo mês mas a dívida só aumenta. Parece um buraco sem fundo.",
        data_point="82,8M com CPF negativado · dívida média R$6.177",
        source="Serasa mar/2026 (recorde histórico)"
    )

    # ── P3: Investidor Iniciante ──────────────────────────────────────────────
    slide_persona(prs,
        code="P3", name="O Investidor Iniciante", color=C_ACCENT,
        age="25–40 anos",
        income="R$4.000–10.000/mês · Classe C ascendente",
        profile_desc="Começou a investir nos últimos 2 anos. Tem Tesouro Selic ou CDB, mas não entende a diferença entre os produtos. Usa Excel para controle de gastos.",
        pains=[
            "Não sabe distinguir CDB, LCI, Tesouro, fundos",
            "Medo de perder dinheiro em renda variável",
            "52% sem capacidade de montar planejamento",
            "46% não se sentem confiantes para metas longas",
            "Sente que chegou tarde ao mercado",
            "Não sabe se está investindo o suficiente",
        ],
        wants=[
            "Comparativo visual claro entre produtos",
            "Recomendação personalizada por perfil de risco",
            "Dashboard de patrimônio consolidado",
            "Calculadora: quanto guardar para atingir X?",
            "Educação contextualizada (aprende enquanto faz)",
        ],
        quote="Sei que preciso investir mas não faço ideia se o que escolhi é bom pra mim.",
        data_point="50% começaram a investir nos últimos 2 anos",
        source="Ipsos/Nubank 2024"
    )

    # ── P4: Mulher Invisível ──────────────────────────────────────────────────
    slide_persona(prs,
        code="P4", name="A Mulher Invisível ao Sistema", color=C_PINK,
        age="25–60 anos",
        income="Variada · frequentemente abaixo da média masculina",
        profile_desc="Cuida das finanças da casa mas raramente das próprias. O mercado financeiro nunca falou com ela. É a maioria dos inadimplentes mas a minoria dos investidores.",
        pains=[
            "Apenas 31% das mulheres investem (vs. 41% homens)",
            "69% desconhecem produtos financeiros (ANBIMA)",
            "50,4% dos inadimplentes são mulheres",
            "Exclui-se do mercado: 'isso não é para mim'",
            "Impacto emocional do endividamento mais profundo",
            "Menor acesso a crédito formal",
        ],
        wants=[
            "Linguagem acessível, sem jargão técnico",
            "Educação financeira básica integrada e não intimidadora",
            "Comparativos visuais simples (sem gráficos complexos)",
            "Funcionalidades de organização doméstica",
            "Comunidade e pertencimento ('não estou sozinha')",
        ],
        quote="Sempre achei que investimento era coisa de homem rico. Ninguém nunca me explicou que era pra mim também.",
        data_point="31% das mulheres investem · 69% desconhecem produtos",
        source="ANBIMA Raio-X do Investidor 2025"
    )

    # ── P5: Profissional Ocupado ──────────────────────────────────────────────
    slide_persona(prs,
        code="P5", name="O Profissional Ocupado", color=C_BRAND,
        age="30–50 anos · Millennials / início Gen X",
        income="Acima de R$7.000/mês · Emprego formal · Classe B",
        profile_desc="Tem vários produtos financeiros em bancos e corretoras diferentes mas não tem visão consolidada. Sabe que deveria organizar mas não tem tempo.",
        pains=[
            "Dispersão: 2-3 bancos, 2 corretoras, sem visão unificada",
            "Não sabe se está investindo o suficiente para aposentadoria",
            "Cartão de crédito como maior dívida (mesmo com boa renda)",
            "Sem tempo para acompanhar portfólio regularmente",
            "62% preocupados com inflação do consumo rotineiro",
            "Apenas 4% dos aposentados têm previdência privada",
        ],
        wants=[
            "Open Finance: tudo em um lugar automaticamente",
            "Alertas inteligentes sem precisar olhar todo dia",
            "Relatório mensal de evolução patrimonial",
            "Projeção de aposentadoria realista",
            "Respostas rápidas — pouco tempo de interação",
        ],
        quote="Tenho dinheiro em três bancos e duas corretoras. Não sei qual é meu patrimônio total.",
        data_point="Open Finance: 62M consentimentos ativos (+44% em 12 meses)",
        source="Banco Central do Brasil 2025"
    )

    # ── P6: Investidor Avançado ───────────────────────────────────────────────
    slide_persona(prs,
        code="P6", name="O Investidor Avançado", color=C_TEAL,
        age="28–50 anos",
        income="Acima de R$14.000/mês · Classe A/B",
        profile_desc="Acompanha a B3 diariamente. Conhece CDI, IPCA, SELIC. Tem carteira diversificada com ações, FIIs, renda fixa e cripto. Usa múltiplos apps e quer consolidação.",
        pains=[
            "Falta consolidador de carteira multi-corretora",
            "IR sobre investimentos é complexo e manual",
            "Sem comparativo de performance vs. benchmarks",
            "Apps de banco não têm profundidade analítica",
            "Não sabe quanto pagar de DARF mensalmente",
            "Rebalanceamento manual e demorado",
        ],
        wants=[
            "Consolidação de ativos de múltiplas corretoras",
            "Performance vs. CDI, IPCA, Ibovespa",
            "Simulador de IR e cálculo automático de DARF",
            "Alertas de rebalanceamento por alocação-alvo",
            "Análise de correlação entre ativos",
        ],
        quote="Uso cinco apps diferentes para ver minha carteira completa. Quero um que faça tudo.",
        data_point="PF movimentou R$517,3 bilhões em ações em 2025",
        source="B3 2025"
    )

    # ── P7: Aposentado / Pré-Aposentado ──────────────────────────────────────
    slide_persona(prs,
        code="P7", name="O Aposentado / Pré-Aposentado", color=C_PURPLE,
        age="50–70 anos · Geração X e Boomers",
        income="Renda fixa · INSS + aplicações conservadoras",
        profile_desc="Acumulou patrimônio e quer preservar. Preocupa-se com longevidade e com não acabar o dinheiro antes de acabar a vida. Cresce na B3: +7,96% em 2025.",
        pains=[
            "Medo de não ter renda suficiente até o fim da vida",
            "Dificuldade com interfaces digitais complexas",
            "Fraudes e golpes: 39% já sofreram tentativas",
            "Desconfia de fintechs e marcas desconhecidas",
            "Apenas 4% dos aposentados têm previdência privada",
            "Custo crescente de saúde não previsto",
        ],
        wants=[
            "Interface simples: fonte grande, poucos cliques",
            "Projeção: 'meu dinheiro dura até quando?'",
            "Alertas de fraude e proteção ativa",
            "Conteúdo sobre previdência e longevidade",
            "Suporte humano disponível quando precisar",
        ],
        quote="Não quero ficar rico, só quero ter certeza de que não vou passar necessidade.",
        data_point="Investidores 60+ cresceram 7,96% na B3 em 2025",
        source="B3 2025"
    )

    # ── P8: MEI / Empreendedor Informal ──────────────────────────────────────
    slide_persona(prs,
        code="P8", name="O MEI / Empreendedor Informal", color=C_ORANGE,
        age="25–55 anos",
        income="R$2.000–12.000/mês · variável e incerto",
        profile_desc="Mistura finanças pessoais com as do negócio. Renda irregular dificulta qualquer planejamento. 7,58 milhões de MEIs inadimplentes em 2025.",
        pains=[
            "Não sabe quanto 'é dele' vs. do negócio",
            "Fluxo de caixa irregular: não sabe quando recebe",
            "7,58M de MEIs inadimplentes (recorde, Serasa 2025)",
            "Sem acesso a crédito formal (alto índice de rejeição)",
            "Sem reserva de emergência para o negócio",
            "Nome negativado limita crescimento da empresa",
        ],
        wants=[
            "Separação clara finanças PF × PJ na mesma tela",
            "Previsão de fluxo de caixa para os próximos 30 dias",
            "Categorização automática de despesas do negócio",
            "Relatório simples de faturamento (para dar ao contador)",
            "Crédito rápido e acessível integrado ao app",
        ],
        quote="Recebi R$8.000 esse mês, mas não consigo separar o que é meu do que é da empresa.",
        data_point="7,58M de MEIs inadimplentes · recorde histórico",
        source="Serasa 2025"
    )

    slide_pains_ranking(prs)
    slide_wants(prs)
    slide_central_pain(prs)
    slide_problem_statement(prs)


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = new_prs()
    build(prs)
    out = os.path.join(OUT_DIR, f"personas-{VERSION}.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
