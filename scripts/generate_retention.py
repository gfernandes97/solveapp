"""
Solve — Retention Loop Presentation
Gera docs/presentations/retention-v1.0.0.pptx (10 slides)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

VERSION  = "v1.0.0"
DATE     = "2026-05-13"
ROOT     = os.path.join(os.path.dirname(__file__), "..")
OUT_FILE = os.path.join(ROOT, "docs", "presentations", f"retention-{VERSION}.pptx")

C_OBS = RGBColor(0x0F, 0x17, 0x2A)
C_GRN = RGBColor(0x00, 0xC9, 0x7A)
C_SIG = RGBColor(0x25, 0x63, 0xEB)
C_SOL = RGBColor(0xF5, 0x9E, 0x0B)
C_CLO = RGBColor(0xF8, 0xFA, 0xFC)
C_SMO = RGBColor(0x64, 0x74, 0x8B)
C_WHI = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xEF, 0x44, 0x44)
C_GRY = RGBColor(0xE2, 0xE8, 0xF0)
C_D20 = RGBColor(0x1E, 0x29, 0x3B)

LOOPS = [
    {
        "num": "01", "name": "Alerta Contextual",
        "freq": "Toda vez que há upload + lembretes preventivos",
        "color": C_RED,
        "trigger":    "Lembrete baseado em calendário:\nfatura próxima do fechamento",
        "action":     "Importa o extrato\nou fatura no app",
        "reward":     "Análise imediata com anomalias\ne impacto no orçamento",
        "investment": "Histórico categorizado\nenriquece padrões futuros",
        "notif_pre":  "⚡  Sua fatura costuma fechar no dia 15. Faltam 3 dias — importe agora para ver sua margem.",
        "notif_post": "✅  Fatura analisada: R$180 em delivery — 60% acima da sua média. Ainda dá tempo de ajustar.",
    },
    {
        "num": "02", "name": "Check-in Semanal",
        "freq": "1× por semana — domingo ou segunda",
        "color": C_SIG,
        "trigger":    "Lembrete semanal se não houve\nupload nos últimos 7 dias",
        "action":     "Importa extrato\nda semana no app",
        "reward":     "Resumo da semana com\ncomparação vs. semana anterior",
        "investment": "Preferências de categorias\naprimoram sugestões futuras",
        "notif_pre":  "📋  Faz 7 dias sem análise. Importe seu extrato de domingo para ver como foi a semana.",
        "notif_post": "📊  Semana analisada: R$1.840 gastos · 3 categorias acima do normal · R$420 economizados.",
    },
    {
        "num": "03", "name": "Fechamento Mensal",
        "freq": "1× por mês — início do mês seguinte",
        "color": C_GRN,
        "trigger":    "Lembrete no início do mês:\nextrato do mês anterior disponível",
        "action":     "Importa o extrato mensal\ndo banco ou fatura",
        "reward":     "Relatório completo com evolução\ne comparação mês a mês",
        "investment": "Define meta\npara o próximo mês",
        "notif_pre":  "📅  É dia 3 de maio. Seu extrato de abril já está disponível — importe para fechar o mês.",
        "notif_post": "🎯  Abril fechado: você gastou 12% menos que março. Sobrou R$680 disponíveis para investir.",
    },
    {
        "num": "04", "name": "Progresso de Metas",
        "freq": "Sempre que há upload com meta ativa",
        "color": C_SOL,
        "trigger":    "Meta ativa + X dias sem\natualização de progresso",
        "action":     "Importa extrato para\natualizar o progresso da meta",
        "reward":     "Barra de progresso atualizada\ne celebração de milestones",
        "investment": "Meta mais detalhada\ne prazo calibrado",
        "notif_pre":  "🎯  Sua meta de reserva está ativa. Importe seu extrato para ver se você está no caminho.",
        "notif_post": "🏆  Progresso atualizado: 70% da reserva de emergência. Faltam R$1.200 — mais 2 meses no ritmo atual.",
    },
]

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=14, bold=False, color=C_OBS,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def ML(slide, lines, l, t, w, h, sz=12, color=C_OBS, sp=2.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(sp * 2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.color.rgb = color
    return tb

def bg_fill(slide, color):
    R(slide, 0, 0, 13.33, 7.5, color)

def footer_line(slide, dark_bg=False):
    c = RGBColor(0x4A, 0x55, 0x68) if dark_bg else C_SMO
    T(slide, f"Solve · Retention Loop · {VERSION} · {DATE}",
      0.4, 7.22, 12.5, 0.22, sz=7, color=c)

def chapter_label(slide, text, color=C_GRN):
    T(slide, text.upper(), 0.5, 0.22, 12.0, 0.22, sz=8, bold=True, color=color)

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    R(slide, 0, 0, 0.06, 7.5, C_GRN)
    R(slide, 0, 7.2, 13.33, 0.06, C_GRN)

    T(slide, "solve", 0.6, 0.5, 5, 0.72, sz=42, bold=True, color=C_GRN)
    T(slide, "Agora você sabe.", 0.6, 1.22, 9, 0.4, sz=16, italic=True, color=C_SMO)

    T(slide, "Retention Loop", 0.6, 2.4, 12, 0.92, sz=52, bold=True, color=C_WHI)
    T(slide, "Como o Solve forma um hábito financeiro mensal e retém o usuário",
      0.6, 3.5, 11.5, 0.38, sz=14, color=C_SMO)

    R(slide, 0.6, 4.05, 4.0, 0.035, C_GRN)
    T(slide, f"Versão {VERSION}  ·  {DATE}  ·  Confidencial",
      0.6, 4.22, 9, 0.3, sz=11, color=C_SMO)


def slide_problem(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    footer_line(slide, dark_bg=True)
    chapter_label(slide, "O Problema de Retenção em Fintech", color=C_RED)

    T(slide, "A maioria dos apps de finanças não forma hábito.",
      0.5, 0.52, 12.5, 0.65, sz=28, bold=True, color=C_WHI)
    R(slide, 0.5, 1.18, 12.0, 0.035, C_RED)

    # Churn stats
    stats = [
        ("~70%", "abandonam\nem 30 dias"),
        ("~15%", "chegam ao\nterceiro mês"),
        ("~8%",  "convertem\npara pago"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 4.0
        R(slide, x, 1.32, 3.7, 0.04, C_RED)
        T(slide, num,   x, 1.42, 3.7, 0.65, sz=40, bold=True, color=C_RED)
        T(slide, label, x, 2.08, 3.7, 0.5,  sz=13, color=C_SMO)

    R(slide, 0.5, 2.72, 12.0, 0.025, C_D20)

    # Two failure modes
    T(slide, "DOIS MODOS DE FALHA", 0.5, 2.85, 12, 0.22, sz=8, bold=True, color=C_RED)

    modes = [
        ("App útil sem razão para voltar",
         "Nenhum trigger externo, nenhum hábito formado. O usuário não "
         "volta porque o app não chama — só responde quando procurado."),
        ("App que pede esforço sem retornar valor",
         "Alto custo cognitivo (lançamento manual, categorização), recompensa "
         "fraca e demorada. O usuário desiste antes do primeiro insight real."),
    ]
    for i, (title, desc) in enumerate(modes):
        x = 0.5 + i * 6.4
        R(slide, x, 3.15, 6.1, 0.05, C_RED)
        R(slide, x, 3.2, 6.1, 1.88, C_D20)
        T(slide, title, x + 0.2, 3.32, 5.7, 0.36, sz=13, bold=True, color=C_WHI)
        T(slide, desc,  x + 0.2, 3.74, 5.7, 1.1,  sz=11, color=C_SMO)

    R(slide, 0, 5.28, 13.33, 0.98, C_GRN)
    T(slide, "Hábito não é uma feature. É um loop projetado.",
      0, 5.42, 13.33, 0.48, sz=22, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)

    T(slide, "O Solve usa 4 loops complementares para transformar o app em parte da rotina financeira mensal do usuário.",
      0.5, 6.38, 12.33, 0.55, sz=12, italic=True, color=C_SMO)


def slide_hook(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "O Modelo · Hook Framework — Nir Eyal")

    T(slide, "Os 4 componentes de qualquer hábito digital",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    steps = [
        ("01", "TRIGGER",           C_RED,
         "O que dispara a ação.\nExterno (push, email) → interno (ansiedade, curiosidade).",
         "Push de alerta · Resumo semanal · Email de relatório"),
        ("02", "AÇÃO",              C_SIG,
         "O menor esforço possível.\nA ação deve ser simples o suficiente para acontecer sem motivação extra.",
         "Abrir o app · Ver o resumo · Confirmar uma categoria"),
        ("03", "RECOMPENSA\nVARIÁVEL", C_GRN,
         "Surpresa + progresso visível.\nVariabilidade mantém o loop interessante a cada ciclo.",
         "Insight novo · Comparação com semana anterior · Milestone"),
        ("04", "INVESTIMENTO",      C_SOL,
         "O que o usuário deixa que aumenta o valor futuro.\nMais histórico = melhor experiência no próximo ciclo.",
         "Categorias salvas · Meta criada · Histórico acumulado"),
    ]

    bw = 2.72
    gap = 0.40
    for i, (num, label, color, desc, example) in enumerate(steps):
        x = 0.5 + i * (bw + gap)
        R(slide, x, 1.12, bw, 0.06, color)
        R(slide, x, 1.18, bw, 4.65, C_WHI)
        T(slide, num,   x + 0.15, 1.28, bw - 0.25, 0.38, sz=20, bold=True, color=color)
        T(slide, label, x + 0.15, 1.68, bw - 0.25, 0.52, sz=12, bold=True, color=C_OBS)
        T(slide, desc,  x + 0.15, 2.28, bw - 0.25, 1.55, sz=11, color=C_SMO)
        R(slide, x + 0.15, 3.88, bw - 0.3, 0.025, C_GRY)
        T(slide, "Ex: " + example, x + 0.15, 3.98, bw - 0.25, 1.62, sz=10, italic=True, color=C_SMO)
        if i < 3:
            ax = x + bw + 0.09
            T(slide, "→", ax, 2.45, 0.22, 0.45, sz=18, bold=True, color=C_GRY, align=PP_ALIGN.CENTER)

    R(slide, 0.5, 5.98, 12.33, 0.55, C_OBS)
    T(slide,
      "O investimento acumula contexto que torna o próximo trigger mais relevante — o loop se fortalece a cada ciclo.",
      0.72, 6.08, 11.8, 0.36, sz=11, bold=True, color=C_GRN)


def slide_upload_model(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    chapter_label(slide, "O Upload como Evento Central")

    T(slide, "Sem agregação automática, o upload é o eixo de toda a experiência.",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    # Central upload box
    R(slide, 5.42, 2.62, 2.5, 1.12, C_OBS)
    T(slide, "UPLOAD", 5.42, 2.78, 2.5, 0.38, sz=18, bold=True,
      color=C_GRN, align=PP_ALIGN.CENTER)
    T(slide, "extrato ou fatura", 5.42, 3.15, 2.5, 0.38, sz=10,
      color=C_SMO, align=PP_ALIGN.CENTER)

    # Left: Before upload (triggers that bring user in)
    R(slide, 0.5, 1.18, 4.62, 4.6, C_CLO)
    R(slide, 0.5, 1.18, 4.62, 0.06, C_SIG)
    T(slide, "ANTES DO UPLOAD  ·  trazer o usuário ao app", 0.7, 1.3, 4.25, 0.24,
      sz=9, bold=True, color=C_SIG)
    T(slide, "Gatilhos baseados em:", 0.7, 1.62, 4.25, 0.24, sz=10, bold=True, color=C_OBS)
    ML(slide, [
        "📅  Calendário — fatura fecha no dia X",
        "⏱  Tempo — faz N dias sem análise",
        "🎯  Meta ativa — progresso desatualizado",
        "📆  Hábito semanal — segunda-feira",
        "📣  Extrato mensal disponível no banco",
    ], 0.7, 1.95, 4.2, 2.6, sz=12, color=C_OBS, sp=3.0)

    T(slide, "Objetivo: criar o hábito de importar regularmente,\nnão esperar o usuário lembrar sozinho.",
      0.7, 4.55, 4.2, 0.8, sz=10, italic=True, color=C_SMO)

    # Arrow left → center
    T(slide, "→", 5.12, 3.0, 0.35, 0.45, sz=20, bold=True, color=C_SIG, align=PP_ALIGN.CENTER)

    # Right: After upload (immediate reward)
    R(slide, 8.22, 1.18, 4.6, 4.6, C_CLO)
    R(slide, 8.22, 1.18, 4.6, 0.06, C_GRN)
    T(slide, "DEPOIS DO UPLOAD  ·  recompensar o ato", 8.42, 1.3, 4.2, 0.24,
      sz=9, bold=True, color=C_GRN)
    T(slide, "Triggers disparados pela análise:", 8.42, 1.62, 4.2, 0.24, sz=10, bold=True, color=C_OBS)
    ML(slide, [
        "⚡  Anomalia detectada no extrato",
        "💰  Economia vs. período anterior",
        "📊  Resumo semanal ou mensal pronto",
        "🏆  Milestone de meta atingido",
        "💡  Insight de categoria incomum",
    ], 8.42, 1.95, 4.2, 2.6, sz=12, color=C_OBS, sp=3.0)

    T(slide, "Regra crítica: o primeiro insight deve aparecer\nem menos de 3 segundos após o upload.",
      8.42, 4.55, 4.2, 0.8, sz=10, italic=True, color=C_SMO)

    # Arrow center → right
    T(slide, "→", 7.92, 3.0, 0.35, 0.45, sz=20, bold=True, color=C_GRN, align=PP_ALIGN.CENTER)

    # Bottom principle
    R(slide, 0, 5.98, 13.33, 1.0, C_OBS)
    T(slide, "PRINCÍPIO DE DESIGN", 0.5, 6.08, 4, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide,
      "Se o upload não for recompensado com um insight relevante imediatamente, "
      "o usuário não cria o hábito de importar — e todos os 4 loops deixam de funcionar.",
      0.5, 6.32, 12.3, 0.5, sz=12, color=C_WHI)


def slide_loops_overview(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    chapter_label(slide, "Os 4 Loops do Solve")

    T(slide, "Loops complementares que cobrem diferentes frequências e motivações",
      0.5, 0.5, 12.5, 0.42, sz=20, bold=True, color=C_OBS)
    R(slide, 0.5, 0.94, 12.0, 0.035, C_GRN)

    for i, loop in enumerate(LOOPS):
        y = 1.1 + i * 1.52
        color = loop["color"]
        R(slide, 0.5, y, 12.33, 1.44, C_CLO)
        R(slide, 0.5, y, 0.08, 1.44, color)

        # Loop number badge
        R(slide, 0.75, y + 0.47, 0.58, 0.5, color)
        T(slide, loop["num"], 0.75, y + 0.49, 0.6, 0.46,
          sz=14, bold=True, color=C_OBS if color == C_SOL else C_WHI, align=PP_ALIGN.CENTER)

        T(slide, loop["name"], 1.5, y + 0.18, 3.8, 0.38, sz=15, bold=True, color=C_OBS)
        T(slide, loop["freq"], 1.5, y + 0.6,  3.8, 0.28, sz=11, italic=True, color=C_SMO)

        # Freq pill
        R(slide, 5.6, y + 0.3, 0.06, 0.84, color)

        T(slide, "GATILHO", 5.82, y + 0.18, 2.2, 0.22, sz=8, bold=True, color=color)
        T(slide, loop["trigger"].replace("\n", " "), 5.82, y + 0.45, 3.0, 0.65, sz=11, color=C_OBS)

        R(slide, 9.1, y + 0.3, 0.06, 0.84, color)

        T(slide, "RECOMPENSA PRINCIPAL", 9.32, y + 0.18, 3.5, 0.22, sz=8, bold=True, color=color)
        T(slide, loop["reward"].replace("\n", " "), 9.32, y + 0.45, 3.5, 0.65, sz=11, color=C_OBS)


def slide_loop(prs, loop):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    color = loop["color"]

    R(slide, 0, 0, 13.33, 0.12, color)
    T(slide, f"Loop {loop['num']}  ·  Frequência: {loop['freq']}", 0.5, 0.18, 12, 0.22,
      sz=8, bold=True, color=color)

    T(slide, loop["name"], 0.5, 0.46, 10, 0.62, sz=32, bold=True, color=C_OBS)
    R(slide, 0.5, 1.1, 12.0, 0.03, C_GRY)

    # 4 component columns
    labels    = ["TRIGGER", "AÇÃO", "RECOMPENSA VARIÁVEL", "INVESTIMENTO"]
    descs     = [loop["trigger"], loop["action"], loop["reward"], loop["investment"]]
    col_clrs  = [color, C_SIG, C_GRN, C_SMO]
    bw, gap   = 2.72, 0.41

    for i, (label, desc, cc) in enumerate(zip(labels, descs, col_clrs)):
        x = 0.5 + i * (bw + gap)
        R(slide, x, 1.18, bw, 0.06, cc)
        R(slide, x, 1.24, bw, 3.38, C_CLO)
        T(slide, label, x + 0.15, 1.35, bw - 0.25, 0.24, sz=8, bold=True, color=cc)
        T(slide, desc,  x + 0.15, 1.66, bw - 0.25, 2.72, sz=13, color=C_OBS)
        if i < 3:
            ax = x + bw + 0.09
            T(slide, "→", ax, 2.55, 0.23, 0.45,
              sz=16, bold=True, color=C_GRY, align=PP_ALIGN.CENTER)

    R(slide, 0.5, 4.74, 12.0, 0.03, C_GRY)
    T(slide,
      "↺  O investimento acumula contexto que torna o próximo trigger mais relevante — o loop se fortalece a cada ciclo.",
      0.5, 4.82, 12.3, 0.32, sz=10, italic=True, color=C_SMO)

    # Two-phase notification band
    R(slide, 0, 5.28, 13.33, 1.94, C_OBS)
    # Pre-upload
    R(slide, 0, 5.28, 6.55, 0.04, C_SIG)
    T(slide, "1  LEMBRETE — antes do upload", 0.5, 5.38, 5.8, 0.22, sz=8, bold=True, color=C_SIG)
    T(slide, loop["notif_pre"],  0.5, 5.62, 6.0, 0.98, sz=13, italic=True, color=C_WHI)
    # Divider
    R(slide, 6.62, 5.35, 0.04, 1.72, RGBColor(0x2D, 0x3A, 0x55))
    # Post-upload
    R(slide, 6.72, 5.28, 6.55, 0.04, C_GRN)
    T(slide, "2  PÓS-UPLOAD — após análise", 6.72, 5.38, 5.8, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, loop["notif_post"], 6.72, 5.62, 6.1, 0.98, sz=13, italic=True, color=C_WHI)


def slide_triggers(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Triggers por Persona · Copy de Notificação")

    T(slide, "A mesma dor, o gatilho certo para cada perfil",
      0.5, 0.5, 12.5, 0.42, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.94, 12.0, 0.035, C_GRN)

    personas = [
        {
            "code": "P1", "name": "Jovem Desorganizado",
            "color": C_SIG,
            "triggers": [
                ("Lembrete semanal",
                 '"Faz 7 dias sem ver seus gastos. Importe o extrato e veja para onde foi o dinheiro."'),
                ("Pós-upload: impulso detectado",
                 '"Extrato analisado: R$340 em lazer este mês — 80% foi nos últimos 3 finais de semana."'),
                ("Fatura próxima do fechamento",
                 '"Sua fatura fecha em 4 dias. Importe agora para saber quanto você ainda pode gastar."'),
            ],
        },
        {
            "code": "P2", "name": "Endividado Crônico",
            "color": C_RED,
            "triggers": [
                ("Alerta preventivo de fatura",
                 '"Fatura fecha em 3 dias. Importe para ver sua margem antes de gastar mais."'),
                ("Pós-upload: mínimo detectado",
                 '"Fatura importada: pagamento mínimo identificado. Isso custa R$142 em juros este mês."'),
                ("Progresso na dívida",
                 '"Extrato do mês importado. Você reduziu R$800 da dívida — 4 meses para zerar no ritmo atual."'),
            ],
        },
        {
            "code": "P3", "name": "Investidor Iniciante",
            "color": C_GRN,
            "triggers": [
                ("Fechamento do mês",
                 '"É dia 3. Importe o extrato de abril para calcular quanto sobrou para investir."'),
                ("Pós-upload: sobra calculada",
                 '"Abril fechado: sobrou R$530. Aplicados em CDB renderiam R$62 até dezembro."'),
                ("Meta de investimento",
                 '"Sua meta de R$10.000 em renda fixa está ativa. Importe o extrato para atualizar."'),
            ],
        },
        {
            "code": "P5", "name": "Profissional Ocupado",
            "color": C_SMO,
            "triggers": [
                ("Relatório mensal disponível",
                 '"É início do mês. Importe os extratos de abril para gerar seu relatório em 1 minuto."'),
                ("Pós-upload: anomalia detectada",
                 '"Extrato analisado: 2 gastos incomuns detectados e R$1.100 de sobra identificada."'),
                ("Meta de patrimônio",
                 '"Importe seu extrato para atualizar a projeção de patrimônio. Última atualização: 18 dias atrás."'),
            ],
        },
    ]

    cw = 2.9
    for i, persona in enumerate(personas):
        x = 0.5 + i * (cw + 0.31)
        color = persona["color"]
        R(slide, x, 1.1, cw, 0.06, color)
        R(slide, x, 1.16, cw, 0.32, C_OBS)
        T(slide, persona["code"],  x + 0.12, 1.19, 0.5,  0.26, sz=11, bold=True, color=color)
        T(slide, persona["name"],  x + 0.62, 1.19, cw - 0.7, 0.26, sz=10, bold=True, color=C_WHI)

        for j, (label, notif) in enumerate(persona["triggers"]):
            y = 1.56 + j * 1.86
            bg = C_WHI if j % 2 == 0 else C_CLO
            R(slide, x, y, cw, 1.78, bg)
            T(slide, label, x + 0.12, y + 0.1,  cw - 0.2, 0.24, sz=8, bold=True, color=color)
            T(slide, notif, x + 0.12, y + 0.38, cw - 0.2, 1.25, sz=10, italic=True, color=C_OBS)


def slide_journey_metrics(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    chapter_label(slide, "Jornada de Hábito + Métricas de Retenção")

    T(slide, "Do primeiro acesso ao hábito consolidado",
      0.5, 0.5, 12.5, 0.42, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.94, 12.0, 0.035, C_GRN)

    # Journey timeline (left 7.8")
    phases = [
        ("D1–D3",   "Ativação",       C_RED,  "Primeiro insight — aha moment",     "Alerta Contextual"),
        ("D4–D14",  "Formação",       C_SIG,  "Primeira semana revisada",          "Check-in Semanal"),
        ("D15–D30", "Consolidação",   C_GRN,  "Primeiro relatório mensal",         "Fechamento Mensal"),
        ("D31–D60", "Comprometimento",C_SOL,  "Primeira meta criada e acompanhada","Progresso de Metas"),
        ("D60+",    "Hábito formado", C_GRN,  "Ciclo mensal autossustentado",      "Todos os 4 loops"),
    ]

    for i, (period, phase, color, milestone, loop) in enumerate(phases):
        y = 1.12 + i * 1.16
        R(slide, 0.5, y, 0.06, 1.06, color)
        T(slide, period, 0.72, y + 0.08, 1.1,  0.26, sz=10, bold=True, color=color)
        T(slide, phase,  0.72, y + 0.38, 1.5,  0.28, sz=10, color=C_SMO, italic=True)
        T(slide, milestone, 2.4, y + 0.08, 3.5, 0.36, sz=12, bold=True, color=C_OBS)
        T(slide, "Loop: " + loop, 2.4, y + 0.5, 3.5, 0.38, sz=10, color=C_SMO)

    R(slide, 6.1, 6.02, 6.72, 0.38, C_GRN)
    T(slide, "M3  ·  Momento de conversão para plano pago", 6.1, 6.08, 6.72, 0.26,
      sz=11, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)

    # Metrics table (right 5.8")
    R(slide, 6.25, 1.1, 6.58, 4.82, C_CLO)
    T(slide, "MÉTRICA", 6.42, 1.22, 1.4, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, "META",    8.02, 1.22, 1.1, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, "MERCADO", 9.32, 1.22, 1.2, 0.22, sz=8, bold=True, color=C_SMO)
    T(slide, "SE ABAIXO → OTIMIZAR", 10.72, 1.22, 2.0, 0.22, sz=8, bold=True, color=C_SMO)
    R(slide, 6.25, 1.46, 6.58, 0.025, C_GRY)

    metrics = [
        ("D1  Volta no dia seguinte",   "50%", "30–40%", "Onboarding → primeiro insight mais rápido"),
        ("D7  Completa primeira semana", "30%", "15–25%", "Trigger do check-in semanal"),
        ("D30  Vê primeiro relatório",   "20%", "8–15%",  "Valor percebido nas semanas 2–3"),
        ("M3  Hábito formado",           "12%", "5–8%",   "Loop de metas + personalização"),
        ("M3  Converte para pago",        "8%", "3–6%",   "Paywall no momento certo do ciclo"),
    ]
    for i, (metric, meta, mkt, fix) in enumerate(metrics):
        y = 1.55 + i * 0.9
        bg = C_WHI if i % 2 == 0 else C_CLO
        R(slide, 6.25, y, 6.58, 0.84, bg)
        T(slide, metric, 6.42, y + 0.18, 1.55, 0.45, sz=10, bold=True, color=C_OBS)
        T(slide, meta,   8.02, y + 0.18, 1.1,  0.38, sz=13, bold=True, color=C_GRN)
        T(slide, mkt,    9.32, y + 0.22, 1.2,  0.32, sz=11, color=C_SMO)
        T(slide, fix,   10.72, y + 0.12, 2.0,  0.6,  sz=9,  color=C_SMO, italic=True)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_hook(prs)
    slide_upload_model(prs)
    slide_loops_overview(prs)
    for loop in LOOPS:
        slide_loop(prs, loop)
    slide_triggers(prs)
    slide_journey_metrics(prs)

    os.makedirs(os.path.dirname(OUT_FILE), exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Saved: {os.path.abspath(OUT_FILE)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
