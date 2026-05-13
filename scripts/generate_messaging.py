"""
Solve — Messaging Framework Presentation
Gera docs/presentations/messaging-v1.0.0.pptx (11 slides)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Config ────────────────────────────────────────────────────────────────────
VERSION  = "v1.0.0"
DATE     = "2026-05-12"
ROOT     = os.path.join(os.path.dirname(__file__), "..")
OUT_FILE = os.path.join(ROOT, "docs", "presentations", f"messaging-{VERSION}.pptx")

# ── Paleta Solve ──────────────────────────────────────────────────────────────
C_OBS = RGBColor(0x0F, 0x17, 0x2A)  # Obsidian
C_GRN = RGBColor(0x00, 0xC9, 0x7A)  # Solve Green
C_SIG = RGBColor(0x25, 0x63, 0xEB)  # Signal Blue
C_SOL = RGBColor(0xF5, 0x9E, 0x0B)  # Solar Yellow
C_CLO = RGBColor(0xF8, 0xFA, 0xFC)  # Cloud
C_SMO = RGBColor(0x64, 0x74, 0x8B)  # Smoke
C_WHI = RGBColor(0xFF, 0xFF, 0xFF)  # White
C_RED = RGBColor(0xEF, 0x44, 0x44)  # Red
C_GRY = RGBColor(0xE2, 0xE8, 0xF0)  # Light gray
C_D20 = RGBColor(0x1E, 0x29, 0x3B)  # Dark separator

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=14, bold=False, color=C_OBS,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def ML(slide, lines, l, t, w, h, sz=12, color=C_OBS, sp=2.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(sp * 2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.color.rgb = color
    return tb

def bg_fill(slide, color):
    R(slide, 0, 0, 13.33, 7.5, color)

def footer_line(slide, dark_bg=False):
    c = RGBColor(0x4A, 0x55, 0x68) if dark_bg else C_SMO
    T(slide, f"Solve · Messaging Framework · {VERSION} · {DATE}",
      0.4, 7.22, 12.5, 0.22, sz=7, color=c)

def chapter_label(slide, text, color=C_GRN):
    T(slide, text.upper(), 0.5, 0.22, 12.0, 0.22, sz=8, bold=True, color=color)

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    R(slide, 0, 0, 0.06, 7.5, C_GRN)
    R(slide, 0, 7.2, 13.33, 0.06, C_GRN)

    T(slide, "solve", 0.6, 0.5, 5, 0.72, sz=42, bold=True, color=C_GRN)
    T(slide, "Agora você sabe.", 0.6, 1.22, 9, 0.42, sz=16, italic=True, color=C_SMO)

    T(slide, "Messaging Framework", 0.6, 2.4, 12, 0.95,
      sz=52, bold=True, color=C_WHI)
    T(slide,
      "Copy de landing page · Ads segmentados · Onboarding · Regras de aplicação",
      0.6, 3.55, 11.5, 0.38, sz=14, color=C_SMO)

    R(slide, 0.6, 4.1, 4.0, 0.035, C_GRN)
    T(slide, f"Versão {VERSION}  ·  {DATE}  ·  Confidencial",
      0.6, 4.25, 9, 0.3, sz=11, color=C_SMO)


def slide_pain_point(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    footer_line(slide, dark_bg=True)
    chapter_label(slide, "O Problema que Atacamos", color=C_GRN)

    T(slide,
      '"O brasileiro ganha, gasta, e no fim\ndo mês não sabe o que aconteceu\ncom o dinheiro."',
      1.0, 0.65, 10.5, 3.4, sz=32, bold=True, color=C_WHI)
    T(slide, "— Problem Statement · Solve", 1.0, 4.1, 7, 0.3,
      sz=12, italic=True, color=C_SMO)

    pillars = [
        ("48%",  "não controlam",       "seus gastos"),
        ("85%",  "dos endividados",      "devem no cartão"),
        ("31%",  "sem reserva",          "de emergência"),
    ]
    for i, (num, l1, l2) in enumerate(pillars):
        x = 1.0 + i * 4.0
        R(slide, x, 4.95, 3.4, 0.04, C_GRN)
        T(slide, num, x, 5.05, 3.4, 0.62, sz=36, bold=True, color=C_GRN)
        T(slide, f"{l1}\n{l2}", x, 5.68, 3.4, 0.52, sz=12, color=C_SMO)

    R(slide, 0, 6.75, 13.33, 0.5, C_GRN)
    T(slide, "Agora você sabe.", 0, 6.8, 13.33, 0.38, sz=18, bold=True,
      color=C_OBS, align=PP_ALIGN.CENTER)


def slide_hero(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Hero Universal · Landing Page")

    T(slide, "A mensagem que funciona para qualquer perfil", 0.5, 0.5, 12.5, 0.5,
      sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 1.02, 5.0, 0.035, C_GRN)

    # Dark card
    R(slide, 0.5, 1.2, 12.33, 4.3, C_OBS)

    T(slide, "HEADLINE", 0.8, 1.33, 4, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, "Agora você sabe para onde vai o seu dinheiro.",
      0.8, 1.56, 11.7, 0.72, sz=26, bold=True, color=C_WHI)

    R(slide, 0.8, 2.32, 11.5, 0.02, C_D20)

    T(slide, "SUBHEADLINE", 0.8, 2.42, 4, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide,
      "O Solve conecta todas as suas contas e mostra, em tempo real, o que você ganha, "
      "o que você gasta — e o que você pode fazer diferente.",
      0.8, 2.66, 11.7, 0.65, sz=14, color=C_SMO)

    R(slide, 0.8, 3.35, 11.5, 0.02, C_D20)

    T(slide, "CTA PRIMÁRIO", 0.8, 3.45, 4, 0.22, sz=8, bold=True, color=C_GRN)
    R(slide, 0.8, 3.7, 3.1, 0.42, C_GRN)
    T(slide, "Entrar na lista de espera", 0.82, 3.75, 3.06, 0.32,
      sz=12, bold=True, color=C_OBS, align=PP_ALIGN.CENTER)

    T(slide, "SUB-COPY", 4.2, 3.45, 4, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, "Grátis para começar · Sem cartão de crédito · Configure em minutos",
      4.2, 3.72, 8.3, 0.35, sz=11, color=C_SMO)

    # Why it works
    R(slide, 0.5, 5.6, 12.33, 1.55, RGBColor(0xEF, 0xF6, 0xFF))
    T(slide, "POR QUE FUNCIONA", 0.8, 5.7, 5, 0.22, sz=8, bold=True, color=C_SIG)
    T(slide,
      "Inverte literalmente o pain point central. Não promete riqueza — promete visibilidade. "
      "Qualquer um dos 8 perfis se reconhece nessa frase antes de saber mais sobre o produto.",
      0.8, 5.96, 11.7, 0.9, sz=12, color=C_OBS)


def slide_persona(prs, code, name, age_income, channels, headline, subheadline,
                  bullets, cta, tone, objection, counter, accent):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)

    # Top accent bar
    R(slide, 0, 0, 13.33, 0.12, accent)
    T(slide, f"Persona · {code} — {name}", 0.5, 0.18, 12, 0.25,
      sz=8, bold=True, color=accent)

    # Headline
    T(slide, headline, 0.5, 0.46, 12.5, 0.98, sz=27, bold=True, color=C_OBS)
    # Subheadline
    T(slide, subheadline, 0.5, 1.5, 12.5, 0.62, sz=13, color=C_SMO)

    R(slide, 0.5, 2.16, 12.0, 0.025, C_GRY)

    # LEFT — persona info
    R(slide, 0.35, 2.26, 3.05, 3.2, C_CLO)
    T(slide, "QUEM É",  0.55, 2.36, 2.85, 0.22, sz=8, bold=True, color=accent)
    T(slide, age_income, 0.55, 2.6,  2.85, 0.38, sz=12, color=C_OBS)
    T(slide, "CANAIS",  0.55, 3.18, 2.85, 0.22, sz=8, bold=True, color=accent)
    ML(slide, channels,  0.55, 3.42, 2.85, 1.55, sz=12, color=C_OBS, sp=1.5)
    T(slide, "TOM",     0.55, 4.84, 2.85, 0.22, sz=8, bold=True, color=accent)
    T(slide, tone,      0.55, 5.08, 2.85, 0.38, sz=11, color=C_SMO, italic=True)

    # RIGHT — messaging
    T(slide, "MENSAGEM", 3.65, 2.36, 9.3, 0.22, sz=8, bold=True, color=accent)
    ML(slide, bullets,   3.65, 2.6,  9.1, 1.95, sz=13, color=C_OBS, sp=4.0)

    cta_fg = C_OBS if accent in (C_GRN, C_SOL) else C_WHI
    R(slide, 3.65, 4.68, 2.75, 0.44, accent)
    T(slide, cta, 3.67, 4.73, 2.71, 0.34, sz=12, bold=True,
      color=cta_fg, align=PP_ALIGN.CENTER)

    # Objection band
    R(slide, 0, 5.58, 13.33, 1.62, C_CLO)
    R(slide, 0, 5.58, 13.33, 0.025, C_GRY)

    T(slide, "OBJEÇÃO PRINCIPAL",  0.5, 5.67, 5.5, 0.22, sz=8, bold=True, color=C_RED)
    T(slide, f'"{objection}"',     0.5, 5.92, 6.0, 0.52, sz=12, italic=True, color=C_OBS)

    T(slide, "COMO RESPONDER", 7.0, 5.67, 5.8, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, counter,          7.0, 5.92, 6.0, 0.92, sz=12, color=C_OBS)


def slide_secondary(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Personas Secundárias · Pós-MVP")

    T(slide, "Mensagens para expansão após consolidar P1, P2, P3 e P5",
      0.5, 0.5, 12.5, 0.35, sz=13, color=C_SMO)
    R(slide, 0.5, 0.88, 12.33, 0.035, C_GRN)

    personas = [
        ("P4", "Mulher Invisível ao Sistema",
         '"Dinheiro não é coisa de homem. É coisa sua."',
         "Instagram · Comunidades femininas"),
        ("P6", "Investidor Avançado",
         '"Seu portfólio bate o CDI? Agora você pode saber."',
         "LinkedIn · Fintechs B2B"),
        ("P7", "Aposentado / Pré-Aposentado",
         '"Quanto tempo o seu dinheiro dura? Agora você calcula."',
         "Facebook · WhatsApp"),
        ("P8", "MEI / Empreendedor",
         '"PF aqui, PJ ali. O Solve separa tudo automaticamente."',
         "Instagram · LinkedIn · Contabilidade online"),
    ]
    accents = [
        RGBColor(0xEC, 0x48, 0x99),
        RGBColor(0x06, 0xB6, 0xD4),
        RGBColor(0x84, 0xCC, 0x16),
        C_SOL,
    ]

    for i, ((code, name, headline, channels), accent) in enumerate(zip(personas, accents)):
        y = 1.08 + i * 1.46
        R(slide, 0.5, y, 0.08, 1.22, accent)
        T(slide, code,     0.75, y + 0.05, 0.9,  0.26, sz=10, bold=True, color=accent)
        T(slide, name,     0.75, y + 0.32, 3.6,  0.3,  sz=12, bold=True, color=C_OBS)
        T(slide, headline, 4.5,  y + 0.05, 8.3,  0.52, sz=14, bold=True, color=C_OBS)
        T(slide, channels, 4.5,  y + 0.6,  8.3,  0.3,  sz=11, italic=True, color=C_SMO)
        if i < 3:
            R(slide, 0.5, y + 1.3, 12.33, 0.02, C_GRY)


def slide_rules(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    chapter_label(slide, "Regras de Aplicação")

    T(slide, "Como e onde usar cada mensagem", 0.5, 0.5, 12.5, 0.4,
      sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.93, 12.0, 0.035, C_GRN)

    cols = [
        {
            "title": "Landing Page",
            "color": C_GRN,
            "items": [
                "Topo: hero universal",
                "O visitante ainda não",
                "sabe que é uma persona",
                "",
                "Meio: 1 seção por perfil",
                "com headline e screenshot",
                "específico da persona",
                "",
                "Rodapé: prova social",
                "+ CTA secundário",
            ],
        },
        {
            "title": "Ads Pagos",
            "color": C_SIG,
            "items": [
                "Nunca misturar personas",
                "no mesmo anúncio",
                "",
                "Headline = dor ou objeção",
                "Subheadline = solução",
                "CTA = próximo passo",
                "",
                "P1/P2: pessoas reais,",
                "ambiente cotidiano",
                "P3/P5: dashboard, números",
            ],
        },
        {
            "title": "Onboarding",
            "color": C_SOL,
            "items": [
                '1ª tela: "Como você se',
                'descreveria?" — selecionar',
                "perfil antes de qualquer coisa",
                "",
                "Redirecionar para fluxo",
                "da persona escolhida",
                "",
                "Email D+1: headline da persona",
                "Email D+7: case de resultado",
                "do mesmo perfil",
            ],
        },
        {
            "title": "A/B Testing",
            "color": C_SMO,
            "items": [
                "1 elemento por vez:",
                "headline OU sub OU CTA",
                "(nunca os três juntos)",
                "",
                "P1/P2: métrica = CTR",
                "P3/P5: métrica = ativação",
                "(conexão de conta)",
                "",
                "Ciclo mínimo: 7 dias",
                "ou 500 impressões/variante",
            ],
        },
    ]

    col_w = 2.9
    for i, col in enumerate(cols):
        x = 0.5 + i * 3.19
        R(slide, x, 1.1, col_w, 0.06, col["color"])
        T(slide, col["title"], x, 1.22, col_w, 0.34, sz=14, bold=True, color=col["color"])
        R(slide, x, 1.6, col_w, 5.35, C_CLO)
        ML(slide, col["items"], x + 0.14, 1.72, col_w - 0.2, 5.1,
           sz=11, color=C_OBS, sp=1.2)


def slide_dont(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Lista Negra · O que Nunca Dizer")

    T(slide, "Palavras e promessas que destroem credibilidade ou geram risco regulatório",
      0.5, 0.5, 12.5, 0.36, sz=13, color=C_SMO)
    R(slide, 0.5, 0.88, 12.33, 0.035, C_RED)

    rows = [
        ('Inteligência Artificial',
         'Virou buzzword sem credibilidade — diga o que o produto faz, não a tecnologia'),
        ('Fique rico / Multiplique seu dinheiro',
         'Gera desconfiança + risco regulatório real (BACEN / CVM)'),
        ('"Finanças pessoais"',
         'Soa chato e genérico para P1 e P2 — o maior volume no topo do funil'),
        ('"Plataforma" / "Solução" / "Ecossistema"',
         'Jargão corporativo. Use "app" ou simplesmente "Solve"'),
        ('Comparação direta com concorrentes',
         'Posicione pelo problema que resolvemos, não pela competição'),
        ('"Simples e fácil" como benefício',
         'Todo app afirma isso. Mostre via copy o quanto é rápido — não afirme'),
    ]

    for i, (avoid, reason) in enumerate(rows):
        y = 1.05 + i * 1.02
        bg_row = C_WHI if i % 2 == 0 else C_CLO
        R(slide, 0.5, y, 12.33, 0.96, bg_row)
        R(slide, 0.55, y + 0.22, 0.3, 0.3, C_RED)
        T(slide, "x", 0.55, y + 0.2, 0.34, 0.34, sz=13, bold=True,
          color=C_WHI, align=PP_ALIGN.CENTER)
        T(slide, avoid,  1.05, y + 0.1,  11.5, 0.3,  sz=12, bold=True, color=C_OBS)
        T(slide, reason, 1.05, y + 0.46, 11.5, 0.38, sz=11, color=C_SMO)


def slide_tagline(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    footer_line(slide, dark_bg=True)
    R(slide, 0, 0, 0.06, 7.5, C_GRN)
    chapter_label(slide, "Tagline e Derivações", color=C_GRN)

    T(slide, '"Agora você sabe."', 0.6, 0.52, 12.5, 1.0,
      sz=44, bold=True, color=C_WHI)
    T(slide,
      "A tagline inverte o pain point. Cada derivação aplica essa lógica ao contexto específico.",
      0.6, 1.58, 11.5, 0.36, sz=13, color=C_SMO)
    R(slide, 0.6, 2.02, 11.5, 0.035, C_D20)

    derivations = [
        ("Hero — Landing page",
         "Agora você sabe para onde vai o seu dinheiro.",
         C_GRN),
        ("Email de boas-vindas",
         "A partir de hoje, nada passa despercebido.",
         C_SIG),
        ("Push — insight mensal",
         "Você gastou 34% mais em alimentação este mês. Agora você sabe.",
         C_SOL),
        ("Push — alerta de fatura",
         "Sua fatura fecha em 5 dias. Você ainda tem R$ 320 de margem.",
         RGBColor(0xEC, 0x48, 0x99)),
    ]

    for i, (label, copy, accent) in enumerate(derivations):
        y = 2.22 + i * 1.17
        R(slide, 0.6, y, 0.06, 0.92, accent)
        T(slide, label, 0.84, y,        4.0,  0.28, sz=9, bold=True, color=accent)
        T(slide, copy,  0.84, y + 0.3, 11.5,  0.55, sz=14, color=C_WHI)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_cover(prs)
    slide_pain_point(prs)
    slide_hero(prs)

    slide_persona(
        prs,
        code="P1", name="O Jovem Desorganizado",
        age_income="18–24 anos · até R$2.500/mês",
        channels=["Instagram", "TikTok", "Influenciador financeiro", "YouTube Shorts"],
        headline="Não sabe para onde vai o seu dinheiro? Você não é o único.",
        subheadline="O Solve te mostra, em segundos, quanto você gastou hoje, essa semana e esse mês — sem planilha, sem lançamento manual.",
        bullets=[
            "✓  Conecta automaticamente com seu banco",
            "✓  Avisa antes de você estourar o limite",
            "✓  Resume seus gastos em categorias simples",
        ],
        cta="Começar grátis",
        tone="Sem julgamento · Informal · Visual · Direto",
        objection="Mais um app financeiro que não vou usar.",
        counter="O Solve não depende da sua disciplina. Ele lê suas transações automaticamente — você só abre e vê.",
        accent=C_SIG,
    )

    slide_persona(
        prs,
        code="P2", name="O Endividado Crônico",
        age_income="25–45 anos · 1–5 salários mínimos",
        channels=["Google Search", "Facebook Ads", "YouTube pre-roll"],
        headline="Sabendo para onde vai cada real, você para de cair no rotativo.",
        subheadline="O Solve mostra exatamente quanto você pode gastar antes do fechamento da fatura — e avisa quando está chegando no limite.",
        bullets=[
            "✓  Alerta de fatura antes do vencimento",
            "✓  Veja quanto entra e quanto sai por semana",
            "✓  Plano de quitação em passos simples",
        ],
        cta="Ver como funciona",
        tone="Empático · Sem julgamento · Esperançoso — nunca alarmista",
        objection="Já tentei vários apps e não funcionou.",
        counter="A diferença é que o Solve não te pede nada. Ele lê suas transações sozinho e te avisa antes — não depois que o estrago já foi feito.",
        accent=C_SOL,
    )

    slide_persona(
        prs,
        code="P3", name="O Investidor Iniciante",
        age_income="25–40 anos · R$4.000–10.000/mês",
        channels=["Google Search", "Newsletters financeiras", "YouTube"],
        headline="Organize o que gasta. Invista o que sobra.",
        subheadline="O Solve consolida gastos, conta corrente e investimentos em um único painel. Você vê onde cortar e quanto pode aplicar este mês.",
        bullets=[
            "✓  Visão consolidada de gastos e investimentos num só lugar",
            "✓  Sobra calculada automaticamente todo mês",
            "✓  Metas de investimento com acompanhamento real",
        ],
        cta="Testar grátis 14 dias",
        tone="Racional · Aspiracional · Orientado a resultado concreto",
        objection="Já uso o app da minha corretora.",
        counter="A corretora te mostra onde o dinheiro está. O Solve te mostra de onde ele vem — e como fazer mais chegar até lá.",
        accent=C_GRN,
    )

    slide_persona(
        prs,
        code="P5", name="O Profissional Ocupado",
        age_income="30–50 anos · acima de R$7.000/mês",
        channels=["LinkedIn Ads", "Google Search", "Newsletters executivas"],
        headline="Controle total das suas finanças. Sem planilha, sem esforço.",
        subheadline="O Solve categoriza seus gastos automaticamente e entrega um relatório financeiro mensal completo. Você enxerga o que importa, sem perder tempo com lançamentos.",
        bullets=[
            "✓  Categorização automática de todos os seus gastos",
            "✓  Relatório mensal completo gerado automaticamente",
            "✓  Projeção de patrimônio e planejamento de longo prazo",
        ],
        cta="Começar agora",
        tone="Preciso · Executivo · Orientado à eficiência · Sem firulas",
        objection="Não tenho tempo para configurar mais um app.",
        counter="Configure em menos de 5 minutos. O Solve já começa a organizar e categorizar tudo. Sem planilha, sem lançamento manual tedioso.",
        accent=C_OBS,
    )

    slide_secondary(prs)
    slide_rules(prs)
    slide_dont(prs)
    slide_tagline(prs)

    os.makedirs(os.path.dirname(OUT_FILE), exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Saved: {os.path.abspath(OUT_FILE)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
