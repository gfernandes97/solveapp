"""
Solve — Positioning Statement Presentation
Gera docs/presentations/positioning-v1.0.0.pptx (9 slides)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Config ────────────────────────────────────────────────────────────────────
VERSION  = "v1.0.0"
DATE     = "2026-05-12"
ROOT     = os.path.join(os.path.dirname(__file__), "..")
OUT_FILE = os.path.join(ROOT, "docs", "presentations", f"positioning-{VERSION}.pptx")

# ── Paleta ────────────────────────────────────────────────────────────────────
C_OBS = RGBColor(0x0F, 0x17, 0x2A)
C_GRN = RGBColor(0x00, 0xC9, 0x7A)
C_SIG = RGBColor(0x25, 0x63, 0xEB)
C_SOL = RGBColor(0xF5, 0x9E, 0x0B)
C_CLO = RGBColor(0xF8, 0xFA, 0xFC)
C_SMO = RGBColor(0x64, 0x74, 0x8B)
C_WHI = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xEF, 0x44, 0x44)
C_GRY = RGBColor(0xE2, 0xE8, 0xF0)
C_D20 = RGBColor(0x1E, 0x29, 0x3B)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=14, bold=False, color=C_OBS,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def ML(slide, lines, l, t, w, h, sz=12, color=C_OBS, sp=2.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(sp * 2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.color.rgb = color
    return tb

def bg_fill(slide, color):
    R(slide, 0, 0, 13.33, 7.5, color)

def footer_line(slide, dark_bg=False):
    c = RGBColor(0x4A, 0x55, 0x68) if dark_bg else C_SMO
    T(slide, f"Solve · Positioning Statement · {VERSION} · {DATE}",
      0.4, 7.22, 12.5, 0.22, sz=7, color=c)

def chapter_label(slide, text, color=C_GRN):
    T(slide, text.upper(), 0.5, 0.22, 12.0, 0.22, sz=8, bold=True, color=color)

def badge(slide, num, color=C_GRN):
    R(slide, 0.5, 0.17, 0.38, 0.38, color)
    T(slide, num, 0.5, 0.19, 0.4, 0.34, sz=14, bold=True,
      color=C_OBS, align=PP_ALIGN.CENTER)

def statement_band(slide, label, text, y=6.18):
    R(slide, 0, y, 13.33, 1.04, C_OBS)
    T(slide, label, 0.5, y + 0.08, 3.5, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide, text,  0.5, y + 0.35, 12.3, 0.58, sz=12, italic=True, color=C_SMO)

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    R(slide, 0, 0, 0.06, 7.5, C_GRN)
    R(slide, 0, 7.2, 13.33, 0.06, C_GRN)

    T(slide, "solve", 0.6, 0.5, 5, 0.72, sz=42, bold=True, color=C_GRN)
    T(slide, "Agora você sabe.", 0.6, 1.22, 9, 0.4, sz=16, italic=True, color=C_SMO)

    T(slide, "Positioning Statement", 0.6, 2.4, 12, 0.92,
      sz=52, bold=True, color=C_WHI)
    T(slide, "Âncora estratégica de produto, marketing e time",
      0.6, 3.5, 11.5, 0.38, sz=14, color=C_SMO)

    R(slide, 0.6, 4.05, 4.0, 0.035, C_GRN)
    T(slide, f"Versão {VERSION}  ·  {DATE}  ·  Confidencial",
      0.6, 4.22, 9, 0.3, sz=11, color=C_SMO)


def slide_framework(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Framework · Os 5 Componentes")

    T(slide, "O que compõe um positioning statement eficaz",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    components = [
        ("01", "Para quem",       "O segmento-alvo com a dor mais urgente",    C_GRN),
        ("02", "O Problema",      "A necessidade específica não atendida",      C_SOL),
        ("03", "Nossa Categoria", "Onde competimos e como nos definimos",       C_SIG),
        ("04", "Diferencial",     "O benefício que só nós entregamos",          C_GRN),
        ("05", "Proof Points",    "A razão que torna a promessa crível",        C_SMO),
    ]

    bw = 2.2
    gap = 0.35
    for i, (num, label, desc, color) in enumerate(components):
        x = 0.5 + i * (bw + gap)
        R(slide, x, 1.18, bw, 0.06, color)
        R(slide, x, 1.24, bw, 3.1, C_WHI)
        T(slide, num,   x + 0.15, 1.35, bw - 0.2, 0.55, sz=28, bold=True, color=color)
        T(slide, label, x + 0.15, 1.94, bw - 0.2, 0.35, sz=13, bold=True, color=C_OBS)
        T(slide, desc,  x + 0.15, 2.36, bw - 0.2, 1.72, sz=11, color=C_SMO)
        if i < 4:
            R(slide, x + bw + 0.08, 2.62, gap - 0.16, 0.025, C_GRY)

    R(slide, 0.5, 4.5, 12.0, 0.035, C_GRY)
    T(slide,
      "Positioning não é slogan. É a âncora estratégica que alinha produto, marketing e time "
      "em torno do mesmo problema e da mesma promessa.",
      0.5, 4.6, 12.3, 0.55, sz=13, italic=True, color=C_SMO)

    R(slide, 0.5, 5.35, 12.33, 1.6, C_OBS)
    T(slide, "Sem positioning claro: cada channel fala de um produto diferente.",
      0.75, 5.48, 11.5, 0.36, sz=13, bold=True, color=C_WHI)
    T(slide,
      "Com positioning claro: qualquer pessoa no time explica o Solve em uma frase — "
      "e todas as frases convergem para o mesmo problema e diferencial.",
      0.75, 5.88, 11.5, 0.55, sz=12, color=C_SMO)


def slide_target(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    badge(slide, "01", C_GRN)
    T(slide, "Para quem · Target Customer", 1.05, 0.22, 11, 0.22,
      sz=8, bold=True, color=C_GRN)

    T(slide, "Quem é o nosso usuário prioritário",
      0.5, 0.58, 12.5, 0.46, sz=24, bold=True, color=C_OBS)
    R(slide, 0.5, 1.06, 12.0, 0.035, C_GRN)

    # Left panel
    R(slide, 0.5, 1.18, 6.85, 4.82, C_CLO)
    T(slide, "PERFIL DEMOGRÁFICO", 0.7, 1.3, 6.5, 0.22, sz=8, bold=True, color=C_GRN)
    ML(slide, [
        "Brasileiros economicamente ativos",
        "25 a 50 anos · Classes B e C",
        "Renda mensal: R$2.500 a R$14.000",
        "Smartphone como canal principal",
        "Pelo menos 1 conta bancária ativa",
    ], 0.7, 1.56, 6.4, 1.78, sz=13, color=C_OBS, sp=2.0)

    T(slide, "MOTIVAÇÃO LATENTE", 0.7, 3.5, 6.5, 0.22, sz=8, bold=True, color=C_GRN)
    T(slide,
      "A motivação existe — falta visibilidade e orientação. "
      "23 milhões planejam começar a investir (ANBIMA 2025), "
      "mas 48% não controlam nem os próprios gastos mensais.",
      0.7, 3.74, 6.4, 1.12, sz=12, color=C_SMO)

    # Right: data cards
    cards = [
        ("75%",   "das transações bancárias via mobile",      "Febraban 2025"),
        ("23M+",  "planejam começar a investir",              "ANBIMA 2025"),
        ("80,9%", "das famílias brasileiras endividadas",     "CNC abr/2026"),
        ("82,8M", "CPFs negativados — recorde histórico",     "Serasa mar/2026"),
    ]
    for i, (num, label, src) in enumerate(cards):
        y = 1.18 + i * 1.21
        R(slide, 7.6, y, 5.2, 1.12, C_CLO)
        R(slide, 7.6, y, 0.06, 1.12, C_GRN)
        T(slide, num,   7.82, y + 0.1,  4.8, 0.46, sz=26, bold=True, color=C_OBS)
        T(slide, label, 7.82, y + 0.56, 4.8, 0.3,  sz=11, color=C_SMO)
        T(slide, src,   7.82, y + 0.87, 4.8, 0.2,  sz=9,  italic=True, color=C_GRY)

    statement_band(slide,
        "FRAGMENTO DO STATEMENT:",
        '"Para brasileiros de 25 a 50 anos, das classes B e C, que ganham '
        'mas não sabem para onde vai o dinheiro..."')


def slide_problem(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    footer_line(slide, dark_bg=True)
    badge(slide, "02", C_SOL)
    T(slide, "O Problema · Need Statement", 1.05, 0.22, 11, 0.22,
      sz=8, bold=True, color=C_SOL)

    T(slide, "A necessidade específica não atendida",
      0.5, 0.58, 12.5, 0.46, sz=24, bold=True, color=C_WHI)
    R(slide, 0.5, 1.06, 12.0, 0.035, C_SOL)

    T(slide, '"Sem visibilidade, não há comportamento."',
      0.5, 1.22, 12.2, 0.68, sz=28, bold=True, italic=True, color=C_WHI)
    T(slide,
      "O brasileiro não sabe para onde vai o seu dinheiro. Não é falta de vontade — "
      "é falta de visibilidade no momento da ação. Sem isso, não há como mudar o comportamento.",
      0.5, 1.96, 12.2, 0.72, sz=13, color=C_SMO)

    R(slide, 0.5, 2.76, 12.0, 0.035, C_D20)

    dims = [
        ("Falta de Visibilidade",
         "48% não controlam gastos\n85% dos endividados devem no cartão",
         C_SOL),
        ("Falta de Orientação",
         "52% sem planejamento financeiro\nEducação chega tarde, fora do contexto",
         C_SIG),
        ("Falta do Momento Certo",
         "O insight chega depois do erro\nNunca antes da decisão",
         C_GRN),
    ]
    for i, (title, sub, color) in enumerate(dims):
        x = 0.5 + i * 4.24
        R(slide, x, 2.88, 3.94, 0.06, color)
        R(slide, x, 2.94, 3.94, 2.1, C_D20)
        T(slide, title, x + 0.2, 3.06, 3.55, 0.38, sz=13, bold=True, color=C_WHI)
        T(slide, sub,   x + 0.2, 3.5,  3.55, 0.78, sz=11, color=C_SMO)

    statement_band(slide,
        "FRAGMENTO DO STATEMENT:",
        '"...que por isso não conseguem poupar, investir nem sair do ciclo de endividamento..."',
        y=6.18)


def slide_category_diff(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)

    R(slide, 0.5, 0.17, 0.65, 0.38, C_GRN)
    T(slide, "03+04", 0.5, 0.19, 0.68, 0.34, sz=11, bold=True,
      color=C_OBS, align=PP_ALIGN.CENTER)
    T(slide, "Nossa Categoria e Diferencial Único", 1.32, 0.22, 11, 0.22,
      sz=8, bold=True, color=C_GRN)

    T(slide, "Onde competimos — e o que nos torna únicos",
      0.5, 0.58, 12.5, 0.46, sz=24, bold=True, color=C_OBS)
    R(slide, 0.5, 1.06, 12.0, 0.035, C_GRN)

    # Left: Category
    R(slide, 0.5, 1.18, 5.92, 4.82, C_CLO)
    R(slide, 0.5, 1.18, 5.92, 0.06, C_SIG)
    T(slide, "03  NOSSA CATEGORIA", 0.7, 1.32, 5.6, 0.22, sz=9, bold=True, color=C_SIG)
    T(slide, "Gestão financeira pessoal\ncom análise inteligente e orientação contextual",
      0.7, 1.6, 5.6, 0.72, sz=14, bold=True, color=C_OBS)
    ML(slide, [
        "Não somos um app de lançamento manual",
        "Não somos uma corretora de investimentos",
        "Somos a camada de inteligência entre",
        "o banco e o comportamento do usuário",
    ], 0.7, 2.42, 5.6, 1.75, sz=12, color=C_SMO, sp=2.0)
    T(slide, "Por que essa definição importa:", 0.7, 4.28, 5.6, 0.26,
      sz=9, bold=True, color=C_SIG)
    T(slide,
      "Define com quem competimos, como o produto é descrito em PR "
      "e como vendemos ao usuário.",
      0.7, 4.58, 5.6, 0.58, sz=11, color=C_SMO)

    # Right: Differentiator
    R(slide, 6.68, 1.18, 6.15, 4.82, C_CLO)
    R(slide, 6.68, 1.18, 6.15, 0.06, C_GRN)
    T(slide, "04  DIFERENCIAL ÚNICO", 6.88, 1.32, 5.8, 0.22, sz=9, bold=True, color=C_GRN)
    T(slide, "O único que combina os três:",
      6.88, 1.6, 5.8, 0.36, sz=14, bold=True, color=C_OBS)

    pillars = [
        ("Análise Inteligente",
         "Seus extratos viram categorias e padrões automaticamente"),
        ("Orientação Contextual",
         "A dica certa antes da decisão — não depois do erro"),
        ("Sem Julgamento",
         "Linguagem acessível · sem jargão · sem culpa"),
    ]
    for i, (title, sub) in enumerate(pillars):
        y = 2.1 + i * 1.28
        R(slide, 6.88, y, 0.06, 1.02, C_GRN)
        T(slide, title, 7.1, y,        5.55, 0.32, sz=12, bold=True, color=C_OBS)
        T(slide, sub,   7.1, y + 0.38, 5.55, 0.52, sz=11, color=C_SMO)

    statement_band(slide,
        "FRAGMENTO DO STATEMENT:",
        '"...é o app de gestão financeira pessoal que analisa seus extratos, '
        'entrega visibilidade real dos gastos e orientação no momento certo..."')


def slide_competition(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Vs. Concorrência · O Espaço que Ninguém Ocupa")

    T(slide, "Por que somos únicos no mercado",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    # Table header
    R(slide, 0.5, 1.08, 12.33, 0.42, C_OBS)
    for x, w, h in [(0.65, 3.0, "Concorrente"), (3.85, 3.1, "O que fazem"),
                    (7.15, 5.5, "Nossa diferença")]:
        T(slide, h, x, 1.16, w, 0.26, sz=10, bold=True, color=C_GRN)

    rows = [
        ("Organizze · Mobills",
         "Lançamento manual de transações",
         "Categorização inteligente + orientação contextual antes da decisão"),
        ("Kinvo · Gorila",
         "Consolidação de investimentos",
         "Gastos + investimentos integrados, foco no dia a dia"),
        ("Guiabolso (descontinuado)",
         "Agregação básica, sem orientação",
         "Orientação contextual + UX moderno + suporte ativo"),
        ("Bancos · Nubank",
         "Visão de uma única instituição",
         "Cross-bank, sem conflito de interesse, visão total"),
        ("YNAB",
         "Metodologia de orçamento manual",
         "Categorização inteligente, orientação contextual, mercado BR"),
    ]
    for i, (comp, they, us) in enumerate(rows):
        y = 1.58 + i * 1.0
        bg = C_WHI if i % 2 == 0 else C_CLO
        R(slide, 0.5, y, 12.33, 0.94, bg)
        T(slide, comp, 0.65, y + 0.16, 3.0, 0.55, sz=11, bold=True, color=C_OBS)
        T(slide, they, 3.85, y + 0.16, 3.1, 0.55, sz=11, color=C_SMO)
        R(slide, 7.08, y + 0.1, 0.06, 0.74, C_GRN)
        T(slide, us,   7.22, y + 0.08, 5.4, 0.78, sz=11, color=C_OBS)

    R(slide, 0.5, 6.65, 12.33, 0.52, C_OBS)
    T(slide,
      "Nenhum concorrente combina: agregação automática + orientação contextual "
      "+ sem jargão + modelo freemium brasileiro",
      0.72, 6.73, 11.9, 0.36, sz=11, bold=True, color=C_GRN)


def slide_statement(prs):
    slide = blank(prs)
    bg_fill(slide, C_OBS)
    footer_line(slide, dark_bg=True)
    R(slide, 0, 0, 0.06, 7.5, C_GRN)

    T(slide, "O Positioning Statement · Versão Completa", 0.6, 0.22, 12, 0.22,
      sz=8, bold=True, color=C_GRN)

    # PARA
    R(slide, 0.6, 0.5, 0.06, 1.32, C_GRN)
    T(slide, "PARA", 0.84, 0.5, 1.5, 0.24, sz=9, bold=True, color=C_GRN)
    T(slide,
      "brasileiros de 25 a 50 anos, das classes B e C, que ganham "
      "mas não sabem para onde vai o dinheiro — e por isso não conseguem "
      "poupar, investir nem sair do ciclo de endividamento,",
      0.84, 0.76, 12.0, 0.96, sz=13, color=C_SMO)

    R(slide, 0.6, 1.95, 12.5, 0.025, C_D20)

    # O SOLVE É
    T(slide, "o Solve", 0.84, 2.08, 5.5, 0.62, sz=34, bold=True, color=C_GRN)
    T(slide, "é o app de gestão financeira pessoal",
      0.84, 2.74, 9.5, 0.36, sz=15, color=C_WHI)

    # QUE
    R(slide, 0.6, 3.22, 0.06, 1.18, C_WHI)
    T(slide, "QUE", 0.84, 3.22, 1.2, 0.24, sz=9, bold=True, color=C_WHI)
    T(slide,
      "analisa seus extratos e transações, entrega visibilidade real dos gastos "
      "e orientação no momento certo da decisão financeira —",
      0.84, 3.48, 12.0, 0.66, sz=13, color=C_WHI)

    R(slide, 0.6, 4.28, 12.5, 0.025, C_D20)

    # AO CONTRÁRIO
    R(slide, 0.6, 4.42, 0.06, 1.32, C_SMO)
    T(slide, "AO CONTRÁRIO", 0.84, 4.42, 2.5, 0.24, sz=9, bold=True, color=C_SMO)
    T(slide,
      "de apps que dependem de lançamento manual ou focam exclusivamente "
      "em investimentos, o Solve age na raiz do problema: "
      "transformar visibilidade em comportamento.",
      0.84, 4.68, 12.0, 0.82, sz=13, color=C_SMO)

    # Short version band
    R(slide, 0, 6.18, 13.33, 1.0, C_GRN)
    T(slide, "VERSÃO CURTA", 0.5, 6.26, 3.0, 0.22, sz=8, bold=True, color=C_OBS)
    T(slide,
      '"O Solve é o app que organiza seus gastos, mostra para onde vai o dinheiro '
      '— e orienta no momento certo para mudar isso."',
      0.5, 6.5, 12.33, 0.52, sz=13, italic=True, color=C_OBS)


def slide_proof(prs):
    slide = blank(prs)
    bg_fill(slide, C_WHI)
    footer_line(slide)
    chapter_label(slide, "Proof Points · Razões para Acreditar")

    T(slide, "Dados que sustentam cada parte do statement",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    proofs = [
        ("O problema é real",    "48% dos brasileiros não controlam seus gastos",    "SPC Brasil 2025",    C_RED),
        ("O canal é mobile",     "75% das transações bancárias via smartphone",      "Febraban 2025",      C_SIG),
        ("A motivação existe",   "23M+ planejam começar a investir",                 "ANBIMA 2025",        C_GRN),
        ("107M não investem",    "64% dos brasileiros ainda não investem o próprio dinheiro", "ANBIMA 2025",        C_SMO),
        ("O timing é urgente",   "82,8M CPFs negativados — recorde histórico",       "Serasa mar/2026",    C_SOL),
        ("O mercado paga",       "Freemium: R$0 → R$19,90 → R$34,90/mês",           "Análise interna",    C_GRN),
    ]

    cw, ch, gap = 3.9, 2.28, 0.22
    for i, (claim, data, src, color) in enumerate(proofs):
        col = i % 3
        row = i // 3
        x = 0.5 + col * (cw + gap)
        y = 1.18 + row * (ch + gap)
        R(slide, x, y, cw, 0.06, color)
        R(slide, x, y + 0.06, cw, ch - 0.06, C_CLO)
        T(slide, claim, x + 0.18, y + 0.18, cw - 0.3, 0.28, sz=9, bold=True, color=color)
        T(slide, data,  x + 0.18, y + 0.52, cw - 0.3, 0.88, sz=13, bold=True, color=C_OBS)
        T(slide, src,   x + 0.18, y + 1.56, cw - 0.3, 0.28, sz=9, italic=True, color=C_SMO)


def slide_apply(prs):
    slide = blank(prs)
    bg_fill(slide, C_CLO)
    footer_line(slide)
    chapter_label(slide, "Como Aplicar · Três Contextos")

    T(slide, "O positioning guia decisões de produto, marketing e time",
      0.5, 0.5, 12.5, 0.45, sz=22, bold=True, color=C_OBS)
    R(slide, 0.5, 0.96, 12.0, 0.035, C_GRN)

    contexts = [
        {
            "title": "No Produto",
            "color": C_GRN,
            "items": [
                "Onboarding em < 3 min",
                "é a prova viva do diferencial",
                "",
                "Cada notificação: antes",
                "da decisão — nunca depois",
                "",
                "Feature test: isso aumenta",
                "visibilidade ou orientação?",
                "Se não → reconsidere",
                "",
                "Copy: visibilidade → ação",
                "específica, nunca genérica",
            ],
        },
        {
            "title": "No Marketing",
            "color": C_SIG,
            "items": [
                "Hero da landing page ancora",
                "no positioning — não na feature",
                "",
                "Ads por persona, mas todos",
                "convergem ao mesmo diferencial",
                "",
                'PR: "camada entre o banco',
                'e o comportamento financeiro"',
                "",
                "Nunca comparação direta",
                "com concorrentes no copy",
            ],
        },
        {
            "title": "No Time",
            "color": C_SOL,
            "items": [
                "Hiring: candidato explica",
                "o Solve em 1 frase?",
                "= alinhamento cultural",
                "",
                "Decisões: isso reforça",
                "o positioning?",
                "",
                "Investor pitch: diferencial",
                "defensável vs. mais um",
                "app de finanças",
                "",
                "OKRs derivados do",
                "diferencial, não de features",
            ],
        },
    ]

    cw = 3.85
    gap = 0.32
    for i, ctx in enumerate(contexts):
        x = 0.5 + i * (cw + gap)
        R(slide, x, 1.15, cw, 0.06, ctx["color"])
        T(slide, ctx["title"], x, 1.27, cw, 0.36,
          sz=16, bold=True, color=ctx["color"])
        R(slide, x, 1.68, cw, 5.38, C_WHI)
        ML(slide, ctx["items"], x + 0.18, 1.82, cw - 0.28, 5.1,
           sz=11, color=C_OBS, sp=1.2)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_cover(prs)
    slide_framework(prs)
    slide_target(prs)
    slide_problem(prs)
    slide_category_diff(prs)
    slide_competition(prs)
    slide_statement(prs)
    slide_proof(prs)
    slide_apply(prs)

    os.makedirs(os.path.dirname(OUT_FILE), exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Saved: {os.path.abspath(OUT_FILE)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
